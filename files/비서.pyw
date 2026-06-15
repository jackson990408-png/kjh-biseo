# -*- coding: utf-8 -*-
"""로컬 AI 비서 — Ollama + Edge 웹뷰. 토큰 비용 없음. 클로드 스타일 채팅 UI.
채팅기록 · 파일 지능검색 · 프로그램/파일/웹 실행 · 프로필 학습 · 암호화 금고
· 코드 자동 실행(승인 프리패스) · 저사양 자동 호환
"""
import base64
import ctypes
import glob as globmod
import hashlib
import json
import os
import re
import shutil
import sqlite3
import string
import subprocess
import sys
import threading
import time
import urllib.parse
import urllib.request
import webbrowser

import webview

OLLAMA = "http://localhost:11434"
# 사양에 따라 자동 선택 (위가 우선) — 텍스트용과 이미지용(비전) 분리
# qwen2.5:7b = 도구블록(```search``` 등)을 가장 안정적으로 따르는 모델 → 최우선
TEXT_PREFS = ["qwen2.5:7b", "exaone3.5:7.8b", "gemma3:12b", "exaone3.5:2.4b",
              "gemma3:4b", "gemma3:1b"]
VISION_PREFS = ["gemma3:12b", "gemma3:4b", "gemma3:1b"]
EMBED_MODEL = "bge-m3"          # 시맨틱 검색·장기기억용 임베딩 모델 (다국어, 한국어 강함)
MODEL_SAFE_PULL = "gemma3:4b"   # 아무 모델도 없을 때 받을 저사양 호환 모델
NUM_CTX = 8192                  # RTX2060 6GB + q8_0 KV캐시 (GPU 84% 확인, 드라이버 610.47)
CTX_MSGS = 14                   # 모델에 보내는 최근 메시지 수
HOME = os.path.join(os.path.expanduser("~"), "AI비서")
OUT_DIR = os.path.join(HOME, "결과물")
HIST_DIR = os.path.join(HOME, "채팅기록")
PROFILE = os.path.join(HOME, "프로필.md")
VAULT = os.path.join(HOME, "금고.dat")
SETTINGS = os.path.join(HOME, "설정.json")
SCHED = os.path.join(HOME, "일정.json")           # 일정·메모 (달력 + 알림)
INDEX_DB = os.path.join(HOME, "색인.db")          # 컴퓨터 전체 파일 색인
ANALYSIS = os.path.join(HOME, "컴퓨터분석.md")     # 자동 생성되는 컴퓨터 이해 요약
MEMORY_DB = os.path.join(HOME, "기억.json")        # 장기기억 (임베딩 기반 자동 회상)
TITLES_DB = os.path.join(HOME, "대화제목.json")      # 사용자가 직접 수정한 대화 제목
ERR_LOG = os.path.join(HOME, "오류.log")           # 전역 예외·디버그 로그
EXEC_LOG = os.path.join(HOME, "실행로그.txt")       # 자동 실행한 코드 감사 로그
INBOX = os.path.join(HOME, "받은파일")              # 사용자가 처리할 파일을 떨구는 폴더 (자동 감시)
HIST_VEC = os.path.join(HOME, "대화색인.json")       # 과거 대화 의미검색용 임베딩 캐시
WORKFLOW = os.path.join(HOME, "워크플로우.json")      # 프로젝트·업무 진척도 보드

VERSION = "KJH비서_1.3.7"   # 업데이트 시 이 값만 올리면 됨
# GitHub raw URL — version.json 위치. 빈 문자열이면 업데이트 체크 안 함.
# 예) "https://raw.githubusercontent.com/내아이디/kjh-biseo/main/version.json"
UPDATE_JSON_URL = "https://raw.githubusercontent.com/jackson990408-png/kjh-biseo/main/version.json"

ANSWER_STYLE = "표준"      # 답변 말투·길이: 간결 / 표준 / 자세히 (설정.json의 answer_style)
ASSISTANT_NAME = "KJH비서"  # 설정.json의 name 키로 변경 가능
STYLE_HINT = {
    "간결": "**답변은 2~4줄로 아주 짧게.** 결론만 말하고 군더더기·예시·부연은 생략한다.",
    "표준": "",
    "자세히": "필요하면 충분히 길게 설명해도 된다. 이유·단계·예시를 곁들여 친절하고 자세히.",
}


def log(msg):
    """오류·이벤트를 파일에 남긴다 (.pyw는 콘솔이 없어 에러가 조용히 사라지는 것 방지)."""
    try:
        os.makedirs(HOME, exist_ok=True)
        with open(ERR_LOG, "a", encoding="utf-8") as f:
            f.write(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] {msg}\n")
    except Exception:
        pass
IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}

USER_ROOTS = [
    os.path.expanduser("~/Desktop"), os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"), os.path.expanduser("~/Pictures"), HOME,
]
INDEX_SKIP = {"windows", "windows.old", "program files", "program files (x86)",
              "programdata", "appdata", "$recycle.bin", "system volume information",
              "recovery", "perflogs", "intel", "nvidia", "amd", "drivers", "msocache",
              "node_modules", "__pycache__", "venv", ".venv", "site-packages",
              "onedrivetemp"}

# ---- 설치 프로그램 자동 감지 (새 컴퓨터에서도 알아서 찾음. * = 아무 버전) ----
PROGRAM_PATTERNS = {
    "스케치업": [r"C:\Program Files\SketchUp\SketchUp *\SketchUp\SketchUp.exe",
              r"C:\Program Files\SketchUp\SketchUp *\SketchUp.exe"],
    "캐드": [r"C:\Program Files\Autodesk\AutoCAD *\acad.exe"],
    "포토샵": [r"C:\Program Files\Adobe\Adobe Photoshop *\Photoshop.exe"],
    "일러스트": [r"C:\Program Files\Adobe\Adobe Illustrator *\Support Files\Contents\Windows\Illustrator.exe"],
    "인디자인": [r"C:\Program Files\Adobe\Adobe InDesign *\InDesign.exe"],
    "프리미어": [r"C:\Program Files\Adobe\Adobe Premiere Pro *\Adobe Premiere Pro.exe"],
    "아크로뱃": [r"C:\Program Files\Adobe\Acrobat DC\Acrobat\Acrobat.exe"],
    "워드": [r"C:\Program Files\Microsoft Office\root\Office16\WINWORD.EXE",
           r"C:\Program Files (x86)\Microsoft Office\root\Office16\WINWORD.EXE"],
    "엑셀": [r"C:\Program Files\Microsoft Office\root\Office16\EXCEL.EXE",
           r"C:\Program Files (x86)\Microsoft Office\root\Office16\EXCEL.EXE"],
    "파워포인트": [r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
              r"C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE"],
}


def detect_programs():
    lad  = os.path.expandvars("%LOCALAPPDATA%")
    pf   = os.environ.get("ProgramFiles",         r"C:\Program Files")
    pf86 = os.environ.get("ProgramFiles(x86)",    r"C:\Program Files (x86)")
    appd = os.environ.get("APPDATA",
           os.path.join(os.path.expanduser("~"), "AppData", "Roaming"))

    # ── 기본 Windows 내장 도구 + 한국어 별칭 ──────────────────────────
    progs = {
        "메모장": "notepad.exe",
        "계산기": "calc.exe",
        "그림판": "mspaint.exe",
        "탐색기": "explorer.exe",
        "파일탐색기": "explorer.exe",
        # PowerShell
        "파워셀": "powershell.exe",
        "파워쉘": "powershell.exe",
        "powershell": "powershell.exe",
        "윈도우파워셀": "powershell.exe",
        "윈도우파워쉘": "powershell.exe",
        "파워셀열어줘": "powershell.exe",   # 조사 포함 문장도 처리
        # CMD
        "cmd": "cmd.exe",
        "명령프롬프트": "cmd.exe",
        "명령창": "cmd.exe",
        "도스창": "cmd.exe",
        "명령줄": "cmd.exe",
        # 시스템 관리
        "작업관리자": "taskmgr.exe",
        "제어판": "control.exe",
        "레지스트리": "regedit.exe",
        "레지스트리편집기": "regedit.exe",
        "장치관리자": "devmgmt.msc",
        "서비스": "services.msc",
        "디스크관리": "diskmgmt.msc",
        "이벤트뷰어": "eventvwr.msc",
        "시스템정보": "msinfo32.exe",
        "방화벽": "firewall.cpl",
        "네트워크연결": "ncpa.cpl",
        "사운드설정": "mmsys.cpl",
        "날짜및시간": "timedate.cpl",
        # Windows 설정 (ms- URI)
        "설정": "ms-settings:",
        "윈도우설정": "ms-settings:",
        "윈도우보안": "windowsdefender:",
        "디펜더": "windowsdefender:",
        "윈도우디펜더": "windowsdefender:",
        "업데이트": "ms-settings:windowsupdate",
        "윈도우업데이트": "ms-settings:windowsupdate",
        "블루투스설정": "ms-settings:bluetooth",
        "와이파이설정": "ms-settings:network-wifi",
        # 캡처
        "스니핑도구": "SnippingTool.exe",
        "캡처도구": "SnippingTool.exe",
        "화면캡처": "SnippingTool.exe",
        "스티커메모": "stikynot.exe",
        # 미디어
        "윈도우미디어플레이어": "wmplayer.exe",
        "미디어플레이어": "wmplayer.exe",
        # 기타
        "원드라이브": "OneDrive.exe",
        "onedrive": "OneDrive.exe",
        "작업표시줄설정": "ms-settings:taskbar",
    }

    # ── PROGRAM_PATTERNS (기존 설계 프로그램 등) ──────────────────────
    for name, pats in PROGRAM_PATTERNS.items():
        for pat in pats:
            hits = sorted(globmod.glob(pat), reverse=True)  # 최신 버전 우선
            if hits:
                progs[name] = hits[0]
                break

    # ── 별칭 확장 ─────────────────────────────────────────────────────
    if "캐드" in progs:
        progs["오토캐드"] = progs["캐드"]
        progs["autocad"]  = progs["캐드"]
    if "스케치업" in progs:
        progs["sketchup"] = progs["스케치업"]
    if "포토샵" in progs:
        progs["photoshop"] = progs["포토샵"]
    if "일러스트" in progs:
        progs["illustrator"] = progs["일러스트"]
        progs["일러스트레이터"] = progs["일러스트"]
    if "프리미어" in progs:
        progs["premiere"] = progs["프리미어"]
        progs["프리미어프로"] = progs["프리미어"]
    if "인디자인" in progs:
        progs["indesign"] = progs["인디자인"]

    # ── 브라우저·일반 프로그램 패턴 탐색 ─────────────────────────────
    extra = {
        # 브라우저
        "크롬": [
            os.path.join(lad,  r"Google\Chrome\Application\chrome.exe"),
            os.path.join(pf,   r"Google\Chrome\Application\chrome.exe"),
            os.path.join(pf86, r"Google\Chrome\Application\chrome.exe"),
        ],
        "구글크롬": "크롬",
        "웨일": [
            os.path.join(lad, r"Naver\Naver Whale\Application\whale.exe"),
            os.path.join(pf,  r"Naver\Naver Whale\Application\whale.exe"),
        ],
        "네이버웨일": "웨일",
        "웨일브라우저": "웨일",
        "엣지": [
            os.path.join(pf,   r"Microsoft\Edge\Application\msedge.exe"),
            os.path.join(pf86, r"Microsoft\Edge\Application\msedge.exe"),
        ],
        "마이크로소프트엣지": "엣지",
        "파이어폭스": [
            os.path.join(pf,   r"Mozilla Firefox\firefox.exe"),
            os.path.join(pf86, r"Mozilla Firefox\firefox.exe"),
        ],
        "브레이브": [
            os.path.join(lad, r"BraveSoftware\Brave-Browser\Application\brave.exe"),
        ],
        "오페라": [
            os.path.join(lad, r"Programs\Opera\opera.exe"),
            os.path.join(lad, r"Programs\Opera GX\opera.exe"),
        ],
        # 개발 도구
        "vscode": [
            os.path.join(lad, r"Programs\Microsoft VS Code\Code.exe"),
            os.path.join(pf,  r"Microsoft VS Code\Code.exe"),
        ],
        "비주얼스튜디오코드": "vscode",
        "코드": "vscode",
        "비주얼스튜디오": [
            os.path.join(pf,   r"Microsoft Visual Studio\*\*\Common7\IDE\devenv.exe"),
            os.path.join(pf86, r"Microsoft Visual Studio\*\*\Common7\IDE\devenv.exe"),
        ],
        "깃": [
            os.path.join(pf,  r"Git\git-bash.exe"),
            os.path.join(pf86,r"Git\git-bash.exe"),
        ],
        "깃배시": "깃",
        "파이참": [
            os.path.join(lad, r"JetBrains\Toolbox\apps\PyCharm-P\*\*\bin\pycharm64.exe"),
            os.path.join(pf,  r"JetBrains\PyCharm *\bin\pycharm64.exe"),
        ],
        # 커뮤니케이션
        "카카오톡": [
            os.path.join(lad, r"Kakao\KakaoTalk\KakaoTalk.exe"),
            os.path.join(lad, r"Programs\Kakao\KakaoTalk\KakaoTalk.exe"),
        ],
        "카톡": "카카오톡",
        "카카오워크": [
            os.path.join(lad, r"Kakao\KakaoWork\KakaoWork.exe"),
        ],
        "줌": [
            os.path.join(lad, r"Zoom\bin\Zoom.exe"),
        ],
        "zoom": "줌",
        "팀즈": [
            os.path.join(lad, r"Microsoft\Teams\Update.exe"),
            os.path.join(pf,  r"Microsoft\Teams\current\Teams.exe"),
            os.path.join(appd,r"Microsoft Teams\Update.exe"),
        ],
        "마이크로소프트팀즈": "팀즈",
        "teams": "팀즈",
        "디스코드": [
            os.path.join(lad, r"Discord\Update.exe"),
            os.path.join(lad, r"Discord\app-*\Discord.exe"),
        ],
        "슬랙": [
            os.path.join(lad, r"slack\slack.exe"),
        ],
        # 미디어
        "팟플레이어": [
            os.path.join(pf,   r"DAUM\PotPlayer\PotPlayerMini64.exe"),
            os.path.join(pf86, r"DAUM\PotPlayer\PotPlayerMini.exe"),
            os.path.join(pf,   r"DAUM\PotPlayer\PotPlayerMini.exe"),
        ],
        "팟플": "팟플레이어",
        "곰플레이어": [
            os.path.join(pf,   r"GOM\GOM Player\GOM.exe"),
            os.path.join(pf86, r"GOM\GOM Player\GOM.exe"),
        ],
        "곰플": "곰플레이어",
        "vlc": [
            os.path.join(pf,   r"VideoLAN\VLC\vlc.exe"),
            os.path.join(pf86, r"VideoLAN\VLC\vlc.exe"),
        ],
        # 유틸리티
        "7zip": [
            os.path.join(pf,   r"7-Zip\7zFM.exe"),
            os.path.join(pf86, r"7-Zip\7zFM.exe"),
        ],
        "반디집": [
            os.path.join(lad, r"Bandizip\Bandizip.exe"),
            os.path.join(pf,  r"Bandizip\Bandizip.exe"),
        ],
        "알집": [
            os.path.join(pf,   r"ESTsoft\ALZip\ALZip.exe"),
            os.path.join(pf86, r"ESTsoft\ALZip\ALZip.exe"),
        ],
        "노트패드플플": [
            os.path.join(pf,   r"Notepad++\notepad++.exe"),
            os.path.join(pf86, r"Notepad++\notepad++.exe"),
        ],
        "notepad++": "노트패드플플",
        # 한글오피스
        "한글": [
            os.path.join(pf,   r"HNC\Hwp*\Hwp.exe"),
            os.path.join(pf86, r"HNC\Hwp*\Hwp.exe"),
            os.path.join(pf,   r"Hnc\Office*\HOffice*\Bin\Hwp.exe"),
            os.path.join(pf86, r"Hnc\Office*\HOffice*\Bin\Hwp.exe"),
        ],
        "한컴오피스": "한글",
        "hwp": "한글",
        # 3D / 렌더링
        "d5render": [
            os.path.join(pf, r"D5 Render\D5Render.exe"),
            os.path.join(lad, r"D5 Render\D5Render.exe"),
        ],
        "d5": "d5render",
        "디파이브": "d5render",
        "enscape": [
            os.path.join(pf,  r"Enscape\Bin64\Enscape.exe"),
            os.path.join(lad, r"Enscape\Bin64\Enscape.exe"),
        ],
        "엔스케이프": "enscape",
        "rhino": [
            os.path.join(pf, r"Rhino *\System\Rhino.exe"),
        ],
        "라이노": "rhino",
        "3dmax": [
            os.path.join(pf, r"Autodesk\3ds Max *\3dsmax.exe"),
        ],
        "쓰리디맥스": "3dmax",
        "마야": [
            os.path.join(pf, r"Autodesk\Maya *\bin\maya.exe"),
        ],
        "maya": "마야",
        # PDF
        "pdf리더": [
            os.path.join(pf86, r"Adobe\Acrobat Reader DC\Reader\AcroRd32.exe"),
            os.path.join(pf,   r"Adobe\Acrobat Reader DC\Reader\AcroRd32.exe"),
        ],
        "어도비리더": "pdf리더",
        "foxit": [
            os.path.join(pf86, r"Foxit Software\Foxit PDF Reader\FoxitPDFReader.exe"),
            os.path.join(pf,   r"Foxit Software\Foxit PDF Reader\FoxitPDFReader.exe"),
        ],
        "폭싯": "foxit",
        # 스팀·게임
        "스팀": [
            os.path.join(pf86, r"Steam\Steam.exe"),
            os.path.join(pf,   r"Steam\Steam.exe"),
        ],
        "steam": "스팀",
    }

    # extra 처리: 별칭(str)이면 이미 등록된 키에서 복사, 리스트면 경로 탐색
    for name, val in extra.items():
        if isinstance(val, str):
            if val in progs:
                progs[name] = progs[val]
        else:
            for p in val:
                # glob 패턴 지원 (* 포함 시)
                hits = sorted(globmod.glob(p), reverse=True) if "*" in p else ([p] if os.path.isfile(p) else [])
                if hits:
                    progs[name] = hits[0]
                    break

    # ── Windows 레지스트리에서 설치 프로그램 자동 탐색 ─────────────────
    try:
        import winreg
        reg_paths = [
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\Uninstall"),
            (winreg.HKEY_CURRENT_USER,  r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
        ]
        for hive, reg_path in reg_paths:
            try:
                with winreg.OpenKey(hive, reg_path) as key:
                    i = 0
                    while True:
                        try:
                            sub = winreg.EnumKey(key, i); i += 1
                        except OSError:
                            break
                        try:
                            with winreg.OpenKey(key, sub) as sk:
                                def rv(n, default=""):
                                    try: return winreg.QueryValueEx(sk, n)[0]
                                    except OSError: return default
                                display = rv("DisplayName")
                                loc     = rv("InstallLocation")
                                exe_raw = rv("DisplayIcon")
                                if not display:
                                    continue
                                # exe 경로 추출
                                exe_path = ""
                                if exe_raw:
                                    exe_path = exe_raw.split(",")[0].strip('"').strip()
                                elif loc:
                                    # InstallLocation 에서 실행파일 추정
                                    for fn in os.listdir(loc) if os.path.isdir(loc) else []:
                                        if fn.lower().endswith(".exe"):
                                            exe_path = os.path.join(loc, fn); break
                                if not exe_path or not os.path.isfile(exe_path):
                                    continue
                                # 한국어 이름 매핑 (영문 DisplayName → 한국어 키 추가)
                                dn_low = display.lower().replace(" ", "").replace("-", "")
                                # 이미 등록된 프로그램이면 스킵 (기존 키 우선)
                                if dn_low in progs:
                                    continue
                                # 영어 이름 그대로도 등록 (소문자·공백제거)
                                progs[dn_low] = exe_path
                                # 원래 DisplayName 도 등록
                                progs[display] = exe_path
                        except Exception:
                            pass
            except Exception:
                pass
    except Exception:
        pass

    return progs


PROGRAMS = detect_programs()
SITES = {
    "구글드라이브": "https://drive.google.com", "드라이브": "https://drive.google.com",
    "제미나이": "https://gemini.google.com", "지메일": "https://mail.google.com",
    "유튜브": "https://www.youtube.com", "네이버": "https://www.naver.com",
    "클로드": "https://claude.ai",
    "구글": "https://www.google.com", "쿠팡": "https://www.coupang.com",
    "넷플릭스": "https://www.netflix.com", "인스타": "https://www.instagram.com",
    "인스타그램": "https://www.instagram.com", "페이스북": "https://www.facebook.com",
    "엑스": "https://x.com", "트위터": "https://x.com",
    "지마켓": "https://www.gmarket.co.kr", "11번가": "https://www.11st.co.kr",
    "옥션": "https://www.auction.co.kr", "알리": "https://ko.aliexpress.com",
    "아마존": "https://www.amazon.com", "당근": "https://www.daangn.com",
    "무신사": "https://www.musinsa.com", "멜론": "https://www.melon.com",
    "챗지피티": "https://chatgpt.com", "챗gpt": "https://chatgpt.com",
    "깃허브": "https://github.com", "노션": "https://www.notion.so",
    "피그마": "https://www.figma.com", "핀터레스트": "https://www.pinterest.co.kr",
    "네이버지도": "https://map.naver.com", "카카오맵": "https://map.kakao.com",
    "구글맵": "https://maps.google.com", "번역기": "https://papago.naver.com",
    "파파고": "https://papago.naver.com", "위키": "https://ko.wikipedia.org",
    "야놀자": "https://www.yanolja.com", "여기어때": "https://www.goodchoice.kr",
    "에어비앤비": "https://www.airbnb.co.kr", "아고다": "https://www.agoda.com/ko-kr/",
    "스카이스캐너": "https://www.skyscanner.co.kr",
}

# ---- 브라우저: 네이버 웨일이 있으면 웨일로 연다 ----
def _find_browser():
    """기본 브라우저를 자동 감지한다.
    ① Windows 레지스트리에서 시스템 기본 브라우저 실행 경로를 읽는다 (가장 정확).
    ② 레지스트리 실패 시, 설치 경로를 순서대로 탐색한다.
    ③ 아무것도 없으면 None (→ webbrowser.open 으로 폴백)."""
    # ① 레지스트리: 현재 사용자의 기본 브라우저 ProgId
    try:
        import winreg
        for scheme in ("https", "http"):
            try:
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                        rf"Software\Microsoft\Windows\Shell\Associations"
                        rf"\UrlAssociations\{scheme}\UserChoice") as k:
                    prog_id = winreg.QueryValueEx(k, "ProgId")[0]
            except OSError:
                continue
            # ProgId → 실행 명령 조회 (HKCR\<ProgId>\shell\open\command)
            for hive in (winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE,
                         winreg.HKEY_CLASSES_ROOT):
                try:
                    with winreg.OpenKey(hive,
                            rf"{prog_id}\shell\open\command") as k2:
                        cmd = winreg.QueryValueEx(k2, "")[0].strip()
                except OSError:
                    continue
                # 명령에서 실행 파일 경로만 추출 ("C:\...\chrome.exe" -- "%1" 형태)
                import shlex
                try:
                    parts = shlex.split(cmd.replace("\\", "/"))
                    exe = parts[0].replace("/", "\\")
                except Exception:
                    exe = re.sub(r'^"([^"]+)".*', r'\1', cmd)
                if os.path.isfile(exe):
                    return exe
    except Exception:
        pass

    # ② 알려진 브라우저 경로를 우선순위 순으로 탐색
    lad = os.path.expandvars("%LOCALAPPDATA%")
    pf  = os.environ.get("ProgramFiles",  r"C:\Program Files")
    pf86= os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")
    candidates = [
        # 네이버 웨일
        os.path.join(lad, r"Naver\Naver Whale\Application\whale.exe"),
        os.path.join(pf,  r"Naver\Naver Whale\Application\whale.exe"),
        # 구글 크롬
        os.path.join(lad, r"Google\Chrome\Application\chrome.exe"),
        os.path.join(pf,  r"Google\Chrome\Application\chrome.exe"),
        os.path.join(pf86,r"Google\Chrome\Application\chrome.exe"),
        # Microsoft Edge
        os.path.join(pf,  r"Microsoft\Edge\Application\msedge.exe"),
        os.path.join(pf86,r"Microsoft\Edge\Application\msedge.exe"),
        # Mozilla Firefox
        os.path.join(pf,  r"Mozilla Firefox\firefox.exe"),
        os.path.join(pf86,r"Mozilla Firefox\firefox.exe"),
        # Brave
        os.path.join(lad, r"BraveSoftware\Brave-Browser\Application\brave.exe"),
        # Opera
        os.path.join(lad, r"Programs\Opera\opera.exe"),
        os.path.join(lad, r"Programs\Opera GX\opera.exe"),
        # Vivaldi
        os.path.join(lad, r"Vivaldi\Application\vivaldi.exe"),
    ]
    for c in candidates:
        if os.path.isfile(c):
            return c
    return None


BROWSER = _find_browser()
log(f"브라우저: {BROWSER or '시스템 기본(webbrowser)'}")


def open_url(url):
    if BROWSER:
        subprocess.Popen([BROWSER, url])
    else:
        webbrowser.open(url)


def load_profile():
    if not os.path.exists(PROFILE):
        with open(PROFILE, "w", encoding="utf-8") as f:
            f.write("# 사용자 프로필 (AI가 학습한 내용)\n"
                    "- 건축/공간 디자인 작업. 스케치업, 캐드, 포토샵, 엑셀 사용.\n")
    with open(PROFILE, encoding="utf-8") as f:
        return f.read()


def load_analysis():
    try:
        with open(ANALYSIS, encoding="utf-8") as f:
            return f.read()[:1200]
    except Exception:
        return "(아직 분석 전 — 백그라운드에서 이 컴퓨터 분석이 진행 중)"


def build_system():
    # 시스템 프롬프트에 포함할 프로그램 목록 (영문 소문자·중복·내장 단순도구 제외)
    _skip = {"메모장", "계산기", "그림판", "탐색기", "파일탐색기", "오토캐드",
             "autocad", "sketchup", "photoshop", "illustrator", "illustrator",
             "premiere", "indesign", "hwp", "zoom", "steam", "teams"}
    progs = ", ".join(k for k in PROGRAMS
                      if k not in _skip
                      and not k.startswith("http")
                      and len(k) <= 20   # 너무 긴 레지스트리 이름 제외
                      and not any(c in k for c in (r"\\", "/", ".", " ")))
    return (
        f"너는 사용자의 개인 AI 비서 '{ASSISTANT_NAME}'다. 정중하고 유능하며, "
        "군더더기 없이 핵심을 말하고, 요청을 끝까지 해결한다. 항상 한국어. "
        "호칭은 '대표님'. 답변 첫 문장은 결론부터. "
        "답변은 마크다운(목록 -, **굵게**)으로 읽기 좋게 쓴다.\n"
        "\n=== 대화 규칙 ===\n"
        "- 사용자의 **마지막 메시지에만** 답하라. 이전 질문·이전 답변 내용을 다시 꺼내거나 "
        "반복하지 마라.\n"
        "- 짧고 간결하게. 대부분 8줄 이내. 같은 말을 두 번 하지 마라.\n"
        "- 묻지 않은 것을 설명하지 마라.\n"
        + (("- " + STYLE_HINT[ANSWER_STYLE] + "\n") if STYLE_HINT.get(ANSWER_STYLE) else "")
        + "\n=== 너의 정체성 — 단순 챗봇이 아니라 '대표님의 사원' ===\n"
        "너는 대표님 밑에서 일하는 유능한 직원이다. 질문에 답만 하고 끝내는 챗봇이 아니라, "
        "맡은 일을 **직접 끝까지 처리하는 실무자**다. 대표님의 몸이 하나 더 늘어난 것처럼 일한다.\n"
        "- 업무 범위는 디자인·도면에 한정되지 않는다. 기획·자료조사·문서정리·일정관리·쇼핑·"
        "여행계획·정보분석 등 개인 업무 전반을 맡아 처리한다.\n"
        "- '~해줘'라는 요청은 방법을 알려달라는 게 아니라 **네가 직접 해달라**는 뜻이다. "
        "도구(검색·읽기·열기·파이썬·화면자동화·웹검색)를 적극적으로 써서 실제로 실행하라.\n"
        "- 항상 **근거(래퍼런스) 기반**으로 일하라. 추측 대신 ```web/```read/```search로 실제 자료를 "
        "확보한 뒤, 그 자료를 인용해 정리·판단한다. 출처가 있으면 함께 밝힌다.\n"
        "- 자료를 받으면 그냥 나열하지 말고 **표·요점·결론 형태로 보기좋게 정리**해 바로 쓸 수 있게 만든다.\n"
        "- 여러 단계가 필요한 일(여행계획·쇼핑·비교조사 등)은 ① 계획을 1~2줄로 밝히고 ② 도구로 "
        "단계별 실행 ③ 결과를 정리 ④ 대표님이 고를 것만 남겨 묻는다. 중간에 멈춰 떠넘기지 마라.\n"
        "- 쇼핑(쿠팡 등 주문)·사이트 입력처럼 프로그램 안에서의 행동이 필요하면 ```open으로 사이트를 "
        "열고 ```auto로 검색어 입력·클릭까지 진행하라. 단, 결제·주문 확정처럼 되돌리기 어려운 마지막 "
        "단계는 실행하기 직전에 대표님께 한 번 확인받는다.\n"
        "- 주식차트·표·이미지를 첨부하면 직접 보고 분석해 매수/매도 관점, 수치, 추세를 정리해준다. "
        "(단, 투자 판단은 참고용이며 최종 결정은 대표님 몫임을 짧게 덧붙인다.)\n"
        "- **눈치껏 먼저 챙겨라.** 시키지 않아도 맥락상 곧 필요할 자료·다음 단계가 보이면 한발 앞서 "
        "준비해 보여준다. 예) 미팅 일정을 잡으면 장소 길찾기·관련 파일을 미리 찾아두고, 여행지를 "
        "정하면 날씨·교통·숙소 후보까지 같이 정리하고, 파일을 열면 핵심 요약을 먼저 띄운다. "
        "단, 과하게 참견하지 말고 '필요하실 것 같아 미리 준비했습니다' 정도로 1~2개만 덧붙인다.\n\n"
        "너는 이 컴퓨터 전체에 접근할 수 있다. 모든 드라이브의 파일이 색인되어 있고, "
        "파일 관련 요청이 오면 시스템이 '미리 검색한 실존 파일' 목록(때로는 내용까지)을 "
        "메시지에 함께 넣어 준다.\n"
        "\n=== 파일 요청 처리 규칙 (가장 중요) ===\n"
        "- '미리 검색한 실존 파일' 블록이 있으면 그 목록이 곧 검색 결과다. 그것만 근거로 답하라.\n"
        "- 파일 경로를 사용자에게 묻는 것 절대 금지. 파일탐색기·명령프롬프트 사용법 안내 금지. "
        "'검색 기능이 제한적' 같은 변명 금지. 너는 이 컴퓨터를 직접 검색하고 읽을 수 있다.\n"
        "- 미리 검색 결과가 없거나 부족하면 즉시 ```search```를 출력해 직접 검색하라 "
        "(물어보지 말고 먼저 검색).\n"
        "- 파일 내용 질문은 추측하지 말고 ```read```로 읽은 뒤 답하라.\n"
        "- 시스템이 '이미 파일/폴더를 열었음'이라고 표시하면, 열었다고 안내만 하라. "
        "경로를 길게 나열하지 마라.\n"
        "- '인터넷에서 검색한 결과' 블록이 있으면 그것을 근거로 답하라. 브라우저 창도 "
        "이미 사용자 화면에 띄워져 있다.\n"
        "\n=== 이 컴퓨터 분석 요약 (시스템이 자동 생성) ===\n"
        + load_analysis() +
        "\n\n=== 사용자 프로필 (이 사람의 사고방식에 맞춰라) ===\n"
        + load_profile() +
        "\n=== 도구 (해당 블록만 출력하면 시스템이 즉시 자동 실행한다 — 승인 절차 없음) ===\n"
        "1. 파일 검색 (컴퓨터 전체) — 경로를 몰라도 된다. 파일명·폴더명 키워드로:\n"
        "```search\n키워드: 단어1 단어2\n확장자: .dwg .dxf .skp\n```\n"
        "→ 컴퓨터 전체 색인에서 즉시 검색해 결과(경로·날짜)를 준다. 확장자 줄은 생략 가능.\n"
        "2. 파일 내용 읽기 — 검색으로 찾았거나 경로를 아는 파일의 내용이 필요하면:\n"
        "```read\nC:\\경로\\파일.xlsx\n```\n"
        "→ 시스템이 내용을 보내준다 (엑셀/워드/PDF/DXF/텍스트/이미지 가능). "
        "한 번에 최대 3개, 줄마다 하나씩.\n"
        "3. 프로그램/파일/웹 열기 — 사용자가 '열어줘/실행해줘/띄워줘/켜줘'라고 요청한 "
        "경우에 사용. 단순히 파일을 언급하거나 질문만 한 경우에는 열지 마라:\n"
        "```open\n대상\n```\n"
        f"대상 = 프로그램명({progs}) 또는 파일 전체경로 또는 URL. "
        "여러 개면 줄마다 하나씩.\n"
        "   웹 검색을 브라우저로 보여달라면 검색 URL을 대상으로 출력 "
        "(예: https://www.google.com/search?q=검색어 · "
        "https://search.naver.com/search.naver?query=검색어).\n"
        "4. 파이썬 코드 실행 — ```python 블록으로 작성하면 시스템이 즉시 실행하고 결과를 보여준다. "
        "사용 가능: openpyxl(엑셀)·python-docx(워드)·pypdf(PDF)·Pillow(이미지)·ezdxf(DXF 도면)"
        "·win32com(오피스/포토샵/일러스트 COM 자동화). "
        f"새 파일은 반드시 r'{OUT_DIR}' 폴더에 저장하고, 원본 파일은 절대 덮어쓰지 않는다. "
        "실행 결과를 print()로 출력하면 사용자에게 보인다.\n"
        "5. 스케치업(.skp) — Ruby API 스크립트를 ```ruby 블록으로 + '창 > Ruby 콘솔에 붙여넣기' 안내.\n"
        "6. 화면 자동화 (프로그램 안에서 타이핑·클릭) — 사용자가 '검색해줘/입력해줘/클릭해줘' 등 "
        "프로그램 안에서의 행동을 요청한 경우:\n"
        "```auto\n열기: 유튜브\n대기: 3\n입력: 검색어\n키: enter\n```\n"
        "사용 가능한 명령(줄마다 하나): 열기:대상 / 대기:초 / 입력:텍스트(현재 커서 위치에 입력) / "
        "키:enter·tab·esc·ctrl+s 등 / 클릭:x,y(화면좌표)\n"
        "웹사이트 검색은 자동화보다 검색결과 URL을 ```open으로 여는 것이 정확하다.\n"
        "7. 인터넷 정보 검색 — 추천·시세·최신 정보 등이 필요하면:\n"
        "```web\n검색어\n```\n"
        "→ 시스템이 뉴스·위키 결과를 보내준다 (브라우저 창 없이). 줄마다 검색어 하나, "
        "최대 3개. 단계마다 다른 검색어로 여러 번 사용 가능.\n"
        "8. 일정·메모 저장/수정 — 사용자가 약속/일정을 잡아달라거나 '기억해/메모해/기록해줘'라고 하면:\n"
        "```plan\n종류: 일정\n날짜: 2026-06-19\n시간: 15:00\n제목: 강남 미팅\n카테고리: 회사\n알림: 하루 전, 30분 전\n```\n"
        "카테고리는 '개인' 또는 '회사' (생략 시 개인). 시간·알림 줄은 생략 가능. "
        "'다음주 금요일' 같은 표현은 아래 '오늘 날짜' 기준으로 정확한 날짜로 환산해 적어라. 메모는:\n"
        "```plan\n종류: 메모\n내용: 기록할 내용 한 줄\n```\n"
        "일정 수정 요청 시 — 아래 '다가오는 일정'의 id를 사용해:\n"
        "```plan\n종류: 수정\nid: 1749999999999\n날짜: 2026-06-20\n시간: 16:00\n제목: 변경된 제목\n카테고리: 개인\n알림: 30분 전\n```\n"
        "→ 시스템이 달력에 저장하고 알림을 자동 예약한다. 저장 후 다시 묻지 마라.\n"
        "9. 워크플로우(업무 보드) 저장 — 대표님이 '프로젝트 만들어줘 / 업무 추가해줘 / 할 일 정리해줘'처럼 "
        "진행할 일을 정리·관리해달라고 하면 보드에 등록한다:\n"
        "```work\n종류: 업무\n프로젝트: 강남 카페 리노베이션\n업무: 평면도 수정, 자재 견적, 클라이언트 미팅\n중요도: 높음\n마감: 2026-06-20\n```\n"
        "중요도는 높음/보통/낮음. 업무는 콤마로 여러 개 가능. 프로젝트명이 없으면 첫 프로젝트에 들어간다. "
        "프로젝트만 만들려면 '종류: 프로젝트\\n프로젝트: 이름\\n폴더: C:\\\\경로(선택)'. "
        "→ 시스템이 보드에 저장하고 왼쪽 📋 워크플로우에 바로 표시한다. 복잡한 일을 맡으면 먼저 이 블록으로 "
        "업무를 쪼개 등록한 뒤 하나씩 처리하라.\n"
        + "\n=== 오늘 날짜: " + time.strftime("%Y-%m-%d") + f" ({kday(time.strftime('%Y-%m-%d'))}요일) ===\n"
        "=== 다가오는 일정 (2주) — '내 일정 뭐야' 질문에 이걸로 답하라 ===\n"
        + upcoming_summary() + "\n"
        "\n=== 복합 요청 처리 (가장 중요한 사고방식) ===\n"
        "여러 행동이 섞인 요청은 문장을 통째로 검색창에 넣지 마라. "
        "요청을 단계로 끊어 도구를 순서대로 쓴다. "
        "먼저 '이렇게 진행하겠습니다: ① … ② …'로 계획을 1~2줄 밝히고 바로 실행하라.\n"
        "예) '여자친구와 100일 기념 여행지 추천 및 숙박 예약' →\n"
        "① ```web\n100일 기념 국내 커플 여행지 추천\n``` 로 후보 조사\n"
        "② 결과에서 2~3곳을 이유와 함께 추천\n"
        "③ 추천지의 숙소 예약 창을 열어준다: "
        "```open\nhttps://www.yanolja.com/search/여수\n``` "
        "(야놀자·여기어때·에어비앤비·아고다·네이버 중 적절한 곳)\n"
        "④ 무엇을 했고 사용자가 다음에 뭘 고르면 되는지 한 줄로 안내\n"
        "날짜·인원·예산처럼 꼭 필요한 정보가 빠졌으면, 할 수 있는 데까지 실행한 뒤 "
        "한 번에 모아서 물어라.\n"
        "\n=== 쇼핑·주문 흐름 (대표님이 '○○ 주문해줘 / 사줘 / 장바구니에 담아줘') ===\n"
        "추측으로 끝내지 말고 아래 순서대로 실제로 진행한다:\n"
        "① 먼저 ```web 로 제품·가격·평점을 조사해 후보 2~3개를 표(제품·가격·핵심사양·평점)로 정리한다.\n"
        "② 대표님이 살 만한 1순위를 근거와 함께 추천한다.\n"
        "③ 해당 제품 검색 페이지를 ```open 으로 브라우저에 띄운다. 쇼핑몰별 검색 URL:\n"
        "   · 쿠팡  https://www.coupang.com/np/search?q=검색어\n"
        "   · 네이버쇼핑  https://search.shopping.naver.com/search/all?query=검색어\n"
        "   · 핀터레스트  https://www.pinterest.co.kr/search/pins/?q=검색어 ← 검색어에 공백이 있으면 +로 대체(예: '미드센추리 모던' → ?q=미드센추리+모던)\n"
        "   · 후즈     https://www.houzz.com/photos/query/검색어\n"
        "   · 아케이데일리 https://www.archdaily.com/search/all?q=검색어\n"
        "   · 디진     https://www.dezeen.com/?s=검색어\n"
        "   · 11번가  https://search.11st.co.kr/Search.tmall?kwd=검색어\n"
        "   · 지마켓  https://browse.gmarket.co.kr/search?keyword=검색어\n"
        "   (대표님이 특정 몰을 말하지 않으면 쿠팡을 기본으로 연다.)\n"
        "④ 제품 페이지에서 장바구니 담기까지는 ```auto 로 진행해도 된다. "
        "다만 **결제·주문 확정 버튼은 절대 자동으로 누르지 마라.** 수량·옵션·배송지·총액을 한 줄로 "
        "요약해 보여주고 '이대로 결제할까요?'라고 반드시 한 번 확인받은 뒤, 대표님이 승인하면 진행한다.\n"
        "⑤ 자주 사는 물건이면 워크플로우(```work)나 메모(```plan)에 기록해 다음에 더 빨리 처리한다.\n"
        "\n=== 보안 ===\n"
        "- 비밀번호·계좌·주민번호 등 민감정보는 절대 따라 말하지 말고 '/금고' 사용을 권하라.\n"
        "\n=== 학습 ===\n"
        "- 사용자의 취향·작업방식 등 기억할 새 사실을 알게 되면 답변 끝에 "
        "```memory\n한 줄 요약\n``` 블록을 붙여라.\n"
        "\n=== 태도 ===\n"
        "- 포기하지 말고 가능한 방법을 끝까지 시도. 모르면 모른다고 솔직히."
    )


# ---------------- 자동 정리 (채팅기록 용량 관리) ----------------
def prune_history(max_files=300, max_mb=30):
    """오래된 채팅기록 자동 삭제 — 채팅기록만, 파일/문서는 절대 건드리지 않음"""
    try:
        files = sorted(globmod.glob(os.path.join(HIST_DIR, "대화_*.json")))
        for p in files[:-1]:
            try:
                if os.path.getsize(p) < 10:
                    os.remove(p)
            except OSError:
                pass
        files = sorted(globmod.glob(os.path.join(HIST_DIR, "대화_*.json")))
        total = sum(os.path.getsize(p) for p in files)
        removed = 0
        while files and (len(files) > max_files or total > max_mb * 1e6):
            p = files.pop(0)
            total -= os.path.getsize(p)
            os.remove(p)
            removed += 1
        return removed
    except Exception:
        return 0


# 삭제 계열 코드 패턴 — 프리패스 모드에서도 경고 문구는 표시 (실행은 그대로 진행)
DELETE_PATTERNS = re.compile(
    r"os\.remove|os\.unlink|shutil\.rmtree|os\.rmdir|send2trash|\.unlink\(|pathlib.*\.rmdir|del\s+/|rmdir\s|Remove-Item",
    re.I)

# 되돌릴 수 없는 치명적 명령 — 프리패스라도 실행 거부 (디스크 포맷·파티션·시스템 전체 삭제 등)
HARD_BLOCK = re.compile(
    r"\bformat\b\s+[a-z]:|diskpart|mkfs|del\s+/[sq].*\\\*|rmdir\s+/s\s+[a-z]:\\?\s*$|"
    r"shutil\.rmtree\(\s*['\"][a-z]:\\?['\"]|rmtree\(\s*['\"][a-z]:[\\/]?['\"]|"
    r"Remove-Item.*-Recurse.*[Cc]:\\?\s|Format-Volume|Clear-Disk|"
    r"reg\s+delete\s+HK|rd\s+/s\s+/q\s+[a-z]:\\?\s*$",
    re.I)


# ---------------- 실행 도구 ----------------
def open_target(t):
    t = t.strip().strip('"')
    if not t:
        return "대상이 비어있습니다"
    # 공백 없는 소문자로 조사
    low = t.replace(" ", "").lower()
    if low in SITES:
        open_url(SITES[low])
        return f"🌐 {t} 열었습니다"
    if t.startswith("http://") or t.startswith("https://"):
        # 쿼리 파라미터에 공백이 있으면 자동 인코딩 (예: ?q=미드센추리 모던 → ?q=미드센추리%20모던)
        if " " in t:
            base, _, qs = t.partition("?")
            if qs:
                enc_parts = []
                for part in qs.split("&"):
                    if "=" in part:
                        k, _, v = part.partition("=")
                        enc_parts.append(k + "=" + urllib.parse.quote(v, safe=""))
                    else:
                        enc_parts.append(part)
                t = base + "?" + "&".join(enc_parts)
        open_url(t)
        return f"🌐 {t} 열었습니다"
    # PROGRAMS 에서 조사 제거 후 검색 (예: "파워셀열어줘" → "파워셀")
    prog_key = None
    if low in PROGRAMS:
        prog_key = low
    else:
        # 조사(열어줘/켜줘/띄워줘/실행해줘 등) 제거 후 재탐색
        for suffix in ("열어줘", "켜줘", "실행해줘", "띄워줘", "시작해줘", "실행", "열기",
                       "열어", "켜", "띄워", "시작"):
            if low.endswith(suffix) and (stripped := low[:-len(suffix)]) in PROGRAMS:
                prog_key = stripped
                break
    if prog_key:
        exe = PROGRAMS[prog_key]
        # ms-settings: / windowsdefender: 같은 URI 스킴
        if ":" in exe and not os.sep in exe:
            subprocess.Popen(f'start "" "{exe}"', shell=True)
            return f"🚀 {t} 열었습니다"
        if os.path.exists(exe):
            subprocess.Popen(f'start "" "{exe}"', shell=True)
            return f"🚀 {t} 실행했습니다"
        # 확장자 없는 내장 명령 (notepad.exe 등)
        if "\\" not in exe:
            subprocess.Popen(f'start "" "{exe}"', shell=True)
            return f"🚀 {t} 실행했습니다"
        return f"{t} 경로를 찾을 수 없습니다: {exe}"
    if os.path.exists(t):
        os.startfile(t)
        return f"📂 {os.path.basename(t)} 열었습니다"
    if "\\" in t or "/" in t:
        return f"경로를 찾을 수 없습니다: {t}"
    try:
        subprocess.Popen(f'start "" "{t}"', shell=True)
        return f"🚀 '{t}' 실행을 시도했습니다"
    except Exception as e:
        return f"실행 실패: {e}"


# ---------------- 화면 자동화 (타이핑·클릭) ----------------
def run_auto(block, log):
    """auto 블록 실행 — 열기/대기/입력/키/클릭. 한글 입력은 클립보드 경유."""
    import pyautogui
    pyautogui.FAILSAFE = True  # 마우스를 화면 모서리로 옮기면 즉시 중단
    for line in block.splitlines():
        line = line.strip()
        if not line or ":" not in line:
            continue
        cmd, arg = [x.strip() for x in line.split(":", 1)]
        try:
            if cmd in ("열기", "open"):
                log(open_target(arg))
                time.sleep(3)
            elif cmd in ("대기", "wait"):
                time.sleep(min(float(arg), 15))
            elif cmd in ("입력", "type"):
                import win32clipboard
                win32clipboard.OpenClipboard()
                win32clipboard.EmptyClipboard()
                win32clipboard.SetClipboardText(arg, win32clipboard.CF_UNICODETEXT)
                win32clipboard.CloseClipboard()
                pyautogui.hotkey("ctrl", "v")
                log(f"⌨ 입력: {arg}")
                time.sleep(0.5)
            elif cmd in ("키", "key"):
                keys = [k.strip().lower().replace("엔터", "enter")
                        for k in arg.split("+")]
                if len(keys) > 1:
                    pyautogui.hotkey(*keys)
                else:
                    pyautogui.press(keys[0])
                log(f"⌨ 키: {arg}")
                time.sleep(0.5)
            elif cmd in ("클릭", "click"):
                x, y = [int(v) for v in arg.replace("，", ",").split(",")[:2]]
                pyautogui.click(x, y)
                log(f"🖱 클릭: {x},{y}")
                time.sleep(0.5)
        except Exception as e:
            log(f"자동화 단계 실패 ({line}): {e}")


# ---------------- 컴퓨터 전체 색인 (코워크식 파일 이해) ----------------
def fixed_drives():
    """고정 디스크 드라이브만 (USB·네트워크 드라이브 제외)"""
    out = []
    for letter in string.ascii_uppercase:
        root = letter + ":\\"
        try:
            if os.path.exists(root) and ctypes.windll.kernel32.GetDriveTypeW(root) == 3:
                out.append(root)
        except Exception:
            pass
    return out or ["C:\\"]


def _db():
    con = sqlite3.connect(INDEX_DB, timeout=30)
    con.execute("CREATE TABLE IF NOT EXISTS files("
                "path TEXT PRIMARY KEY, name TEXT, ext TEXT, size INTEGER, mtime REAL)")
    con.execute("CREATE INDEX IF NOT EXISTS idx_name ON files(name)")
    con.execute("CREATE INDEX IF NOT EXISTS idx_ext ON files(ext)")
    return con


def build_index(status=lambda t: None):
    """전체 드라이브 스캔 → 색인 재구축 + 컴퓨터 분석 요약 생성.
    파일은 읽기(메타데이터)만 — 절대 수정·삭제하지 않음."""
    con = _db()
    cur = con.cursor()
    cur.execute("BEGIN")
    cur.execute("DELETE FROM files")
    n = 0
    batch = []
    for root in fixed_drives():
        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [d for d in dirnames
                           if d.lower() not in INDEX_SKIP and not d.startswith((".", "$"))]
            for fn in filenames:
                p = os.path.join(dirpath, fn)
                try:
                    st = os.stat(p)
                except OSError:
                    continue
                batch.append((p, fn.lower(), os.path.splitext(fn)[1].lower(),
                              st.st_size, st.st_mtime))
                n += 1
            if len(batch) >= 3000:
                cur.executemany("INSERT OR REPLACE INTO files VALUES(?,?,?,?,?)", batch)
                batch = []
                status(f"🗂 컴퓨터 분석 중… {n:,}개")
    if batch:
        cur.executemany("INSERT OR REPLACE INTO files VALUES(?,?,?,?,?)", batch)
    con.commit()
    try:
        summarize_computer(con, n)
    except Exception:
        pass
    con.close()
    return n


def summarize_computer(con, total):
    """색인을 바탕으로 '이 컴퓨터 이해 요약'을 생성 → AI 시스템 프롬프트에 들어감"""
    cats = {"문서": (".docx", ".doc", ".hwp", ".hwpx", ".pdf", ".xlsx", ".xls",
                   ".pptx", ".txt", ".csv"),
            "도면·3D": (".dwg", ".dxf", ".skp", ".3ds", ".max", ".rvt", ".ifc", ".stl"),
            "디자인·이미지": (".psd", ".ai", ".indd", ".jpg", ".jpeg", ".png", ".webp", ".tif"),
            "영상·음악": (".mp4", ".mov", ".avi", ".mp3", ".wav")}
    lines = [f"분석일 {time.strftime('%Y-%m-%d')} · 색인된 파일 {total:,}개"]
    for d in fixed_drives():
        try:
            u = shutil.disk_usage(d)
            lines.append(f"- {d} 드라이브: {u.used / 2**30:.0f}GB 사용 / 전체 {u.total / 2**30:.0f}GB")
        except Exception:
            pass
    lines.append("자료가 많은 폴더 (최근 수정 기준):")
    for cat, exts in cats.items():
        q = ",".join("?" * len(exts))
        cnt = con.execute(f"SELECT COUNT(*) FROM files WHERE ext IN ({q})", exts).fetchone()[0]
        if not cnt:
            continue
        rows = con.execute(f"SELECT path FROM files WHERE ext IN ({q}) "
                           "ORDER BY mtime DESC LIMIT 1500", exts).fetchall()
        dirs = {}
        for (p,) in rows:
            dd = os.path.dirname(p)
            dirs[dd] = dirs.get(dd, 0) + 1
        top = sorted(dirs.items(), key=lambda x: -x[1])[:4]
        lines.append(f"- {cat} {cnt:,}개 — " + " · ".join(dd for dd, _ in top))
    lines.append("설치 프로그램(자동 감지): " +
                 (", ".join(k for k in PROGRAMS
                            if k not in ("메모장", "계산기", "그림판", "탐색기", "오토캐드"))
                  or "기본 프로그램만"))
    with open(ANALYSIS, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def ranked_search(kws, exts=None, limit=15, query=None):
    """관련도 검색 — 키워드+바이그램으로 후보를 모은 뒤 매칭 점수순 정렬.
    '모델사진'↔'모델 사진'처럼 띄어쓰기가 달라도 찾고, 그 위에 임베딩(bge-m3)으로
    의미가 비슷한 파일까지 끌어올린다(동의어·문맥). query는 자연어 원문."""
    kws = [k.lower() for k in kws if k.strip()]
    if not kws:
        return []
    grams = set()
    for k in kws:
        grams.add(k)
        for i in range(len(k) - 1):
            grams.add(k[i:i + 2])
    grams = sorted(grams, key=len, reverse=True)[:20]
    try:
        con = sqlite3.connect(INDEX_DB, timeout=10)
        cond = " OR ".join(["path LIKE ?"] * len(grams))
        args = [f"%{g}%" for g in grams]
        where = f"({cond})"
        if exts:
            where += f" AND ext IN ({','.join('?' * len(exts))})"
            args += sorted(exts)
        rows = con.execute(f"SELECT path, mtime FROM files WHERE {where} "
                           "ORDER BY mtime DESC LIMIT 800", args).fetchall()
        con.close()
    except Exception:
        return []
    scored = []
    for p, mt in rows:
        pl = p.lower()
        full = sum(1 for k in kws if k in pl)            # 키워드 전체 일치
        bi = sum(1 for g in grams if len(g) == 2 and g in pl)  # 부분(바이그램) 일치
        score = full * 10 + bi
        if full or bi >= 2:
            scored.append((score, mt, p))
    scored.sort(key=lambda x: (-x[0], -x[1]))
    # 키워드 상위 후보를 의미 유사도로 재정렬 (임베딩 가능할 때만, 아니면 그대로)
    top = scored[:60]
    q = query or " ".join(kws)
    try:
        return semantic_rerank(q, top, limit)
    except Exception as e:
        log(f"semantic_rerank 실패: {e}")
        return scored[:limit]


def search_files(keywords, exts, limit=30):
    """모델의 ```search``` 도구 — 색인 관련도 검색, 색인 전이면 폴더 직접 탐색"""
    extset = {(e if e.startswith(".") else "." + e).lower() for e in exts} if exts else None
    res = ranked_search(list(keywords), extset, limit)
    if res:
        return [f"{time.strftime('%Y-%m-%d %H:%M', time.localtime(mt))}  {p}"
                for sc, mt, p in res]
    return _walk_search([k.lower() for k in keywords if k.strip()], extset, limit)


def _walk_search(kws, exts, limit=30, max_seconds=20):
    """색인이 아직 없을 때의 폴백 — 사용자 폴더 직접 탐색"""
    t0 = time.time()
    hits = []
    for root in USER_ROOTS:
        if not os.path.isdir(root):
            continue
        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [d for d in dirnames if d.lower() not in INDEX_SKIP]
            for fn in filenames:
                fl = fn.lower()
                if exts and os.path.splitext(fl)[1] not in exts:
                    continue
                if kws and not any(k in fl for k in kws):
                    continue
                p = os.path.join(dirpath, fn)
                try:
                    hits.append((os.path.getmtime(p), p))
                except OSError:
                    pass
            if time.time() - t0 > max_seconds:
                break
        if time.time() - t0 > max_seconds:
            break
    hits.sort(reverse=True)
    return [f"{time.strftime('%Y-%m-%d %H:%M', time.localtime(mt))}  {p}" for mt, p in hits[:limit]]


def parse_search_block(block):
    kws, exts = [], []
    for line in block.splitlines():
        if line.startswith("키워드"):
            kws = line.split(":", 1)[1].split()
        elif line.startswith("확장자"):
            exts = line.split(":", 1)[1].split()
    return kws, exts


# ---------------- 시스템 주도 검색·읽기 (코워크 방식 — 모델 능력에 의존하지 않음) ----------------
INTENT_RE = re.compile(
    r"찾|검색|어디|파일|사진|이미지|문서|폴더|견적|도면|모델링|엑셀|워드|한글|피디에프|pdf|ppt|"
    r"영상|음악|읽어|내용|얼마|확인|분석|요약|보여|열어|띄워", re.I)
READ_RE = re.compile(r"얼마|내용|읽어|알려|요약|분석|확인|뭐였|뭐야|뭔지|언제|누구|어땠")

STOP_WORDS = {
    "내", "나의", "우리", "컴퓨터", "파일", "폴더", "어디", "있는지", "있어", "있나", "있던",
    "좀", "그", "이", "저", "그리고", "근데", "관련", "대한", "해당", "했던", "됐던", "거",
    "찾아줘", "찾아", "찾기", "검색", "검색해줘", "해줘", "해봐", "알려줘", "알려", "보여줘",
    "보여", "열어줘", "열어", "띄워줘", "읽어줘", "확인해줘", "확인", "분석해줘", "분석",
    "요약해줘", "요약", "내용", "얼마", "얼마였어", "얼마야", "뭐였어", "뭐야", "언제", "누구",
}
JOSA = ("으로", "에서", "이랑", "한테", "처럼", "보다", "에게", "까지", "부터", "중에서",
        "중에", "중인", "은", "는", "이", "가", "을", "를", "의", "에", "와", "과", "랑",
        "도", "만", "요", "들", "중")


def extract_keywords(text):
    """문장에서 검색 키워드 추출 — 조사·동사·불용어 제거"""
    toks = re.split(r"[\s,.?!~()\[\]{}\"'·]+", text)
    out = []
    for t in toks:
        t = t.strip()
        if not t:
            continue
        for j in sorted(JOSA, key=len, reverse=True):   # 조사 제거 (긴 것부터)
            if len(t) > len(j) + 1 and t.endswith(j):
                t = t[:-len(j)]
                break
        if len(t) < 2 or t in STOP_WORDS:
            continue
        if re.search(r"(줘|주세요|할까|하자|했어|였어|이야|인지|는지|니|나요|세요|십시오)$", t):
            continue
        out.append(t)
    return out[:6]


def guess_exts(text):
    """요청 문장에서 파일 종류 추론 → 확장자 필터"""
    table = [
        (r"사진|이미지|jpg|png|짤", (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp")),
        (r"엑셀|xlsx|시트", (".xlsx", ".xls", ".csv")),
        (r"워드|docx", (".docx", ".doc")),
        (r"한글파일|hwp", (".hwp", ".hwpx")),
        (r"피디에프|pdf", (".pdf",)),
        (r"견적|계약|보고서|서류|문서", (".pdf", ".xlsx", ".xls", ".docx", ".doc", ".hwp", ".hwpx")),
        (r"도면|캐드|dwg|dxf", (".dwg", ".dxf")),
        (r"스케치업|모델링|skp", (".skp",)),
        (r"영상|동영상|mp4", (".mp4", ".mov", ".avi")),
        (r"피피티|ppt|발표", (".pptx", ".ppt")),
    ]
    for pat, exts in table:
        if re.search(pat, text, re.I):
            return set(exts)
    return None


def auto_context(text):
    """사용자 메시지를 보고 시스템이 직접: 키워드 추출 → 컴퓨터 전체 검색 →
    내용 질문이면 상위 파일을 자동으로 읽기. 결과를 모델에게 떠먹여 준다.
    반환: (컨텍스트 블록, 이미지 목록, 검색결과 [(점수, 수정시각, 경로)])"""
    if not text or text.startswith("/") or not INTENT_RE.search(text):
        return None, [], []
    kws = extract_keywords(text)
    if not kws:
        return None, [], []
    res = ranked_search(kws, guess_exts(text), limit=12, query=text)
    if not res:
        return None, [], []
    lines = [f"{time.strftime('%Y-%m-%d', time.localtime(mt))}  {p}" for sc, mt, p in res]
    block = "=== 시스템이 이 컴퓨터 전체에서 미리 검색한 실존 파일 (관련도순) ===\n" + "\n".join(lines)
    imgs = []
    if READ_RE.search(text):                 # 내용까지 묻는 질문 → 상위 파일 자동 읽기
        read_cnt = 0
        for sc, mt, p in res:
            if read_cnt >= 2:
                break
            if sc < 20:                      # 키워드 2개 이상 확실 일치할 때만
                continue
            try:
                if os.path.getsize(p) > 30 * 1e6:
                    continue
            except OSError:
                continue
            summary, img = analyze_file(p)
            if img:
                if len(imgs) < 1:
                    imgs.append(img)
                    block += f"\n\n[이미지 자동 첨부됨: {os.path.basename(p)}] 이미지를 보고 답하라"
                    read_cnt += 1
                continue
            if summary:
                block += "\n\n=== 파일 내용 (시스템이 자동으로 읽음) ===\n" + summary
                read_cnt += 1
    block += ("\n\n위 목록은 시스템이 실제로 이 컴퓨터에서 찾아낸 것이다. 이 자료만 근거로 답하라. "
              "파일을 열어줄 땐 ```open\n전체경로\n```, 내용이 더 필요하면 ```read\n전체경로\n```. "
              "사용자에게 경로를 묻는 것은 금지.")
    return block, imgs, res


# ---------------- 인터넷 검색 (브라우저 창을 직접 띄워 과정을 보여줌) ----------------
WEB_FORCE_RE = re.compile(
    r"인터넷|웹에서|구글|네이버에서|네이버로|유튜브에서|검색해|알아봐|찾아봐|서치|"
    r"보고\s*싶|듣고\s*싶|틀어|재생", re.I)
# 플랫폼을 직접 언급하면 파일이 발견돼도 웹으로 처리
WEB_PLATFORM_RE = re.compile(r"인터넷|웹에서|구글에서|네이버에서|유튜브에서", re.I)


def mentions_platform(text):
    """'쿠팡에서', '인스타로' 등 사이트+조사 언급 → 컴퓨터 검색 말고 그 사이트에서 검색"""
    low = (text or "").replace(" ", "")
    for name in SITES:
        for josa in ("에서", "으로", "로", "에"):
            if (name + josa) in low:
                return True
    return bool(WEB_PLATFORM_RE.search(text or ""))
WEB_TOPIC_RE = re.compile(
    r"주가|주식|날씨|뉴스|환율|시세|코인|순위|평점|리뷰|추천|최신|요즘|트렌드|"
    r"방법|어떻게|뜻|의미|누구야|뭐야|언제야|얼마야|배송|영업시간", re.I)

# ---------------- 복합 요청 감지 — 행동이 여러 개 섞인 요청은 통검색하지 않고 모델이 단계별 처리 ----------------
ACTION_RE = re.compile(
    r"추천|예약|예매|비교|정리|계획|일정|짜\s*줘|만들|작성|보내|알아보|구매|주문|"
    r"저장|변환|분석|요약|선택|골라", re.I)


def is_multistep(text):
    """'여행지 추천 및 숙박 예약'처럼 행동 동사가 2개 이상이거나 연결어로 묶인 복합 요청인지 판단.
    복합 요청이면 문장 전체를 검색창에 넣는 선처리를 건너뛰고 모델이 단계를 나눠 도구로 처리한다."""
    t = text or ""
    acts = set(ACTION_RE.findall(t))
    linked = re.search(r"및|그리고|하고\s|한\s*다음|다음에|후에|이어서|까지\s*해", t)
    return len(acts) >= 2 or (len(acts) >= 1 and bool(linked))


def web_query(text):
    """문장에서 검색어만 남기기 — 사이트 이름·명령어·조사 제거"""
    q = re.sub(r"인터넷에서|인터넷으로|인터넷|웹에서|검색해줘|검색해봐|검색해|검색|"
               r"알아봐줘|알아봐|찾아봐줘|찾아봐|"
               r"보고\s*싶\S*|듣고\s*싶\S*|하고\s*싶\S*|틀어\S*|재생해?\S*|궁금\S*|"
               r"열어서|열어줘|열어봐|열어|열고|띄워줘|띄워|들어가\S*|접속\S*|"
               r"확인해줘|확인해봐|확인|해줘|해봐|볼래|보자|좀|보여줘", " ", text)
    for name in sorted(SITES, key=len, reverse=True):   # 사이트 이름+조사 제거
        q = re.sub(re.escape(name) + r"(에서|으로|로|에)?", " ", q)
    return re.sub(r"\s+", " ", q).strip(" ?!.,")


# 사이트 안 검색창에 검색어를 넣은 상태로 바로 진입
PLATFORM_SEARCH = {
    "유튜브": "https://www.youtube.com/results?search_query=",
    "넷플릭스": "https://www.netflix.com/search?q=",
    "쿠팡": "https://www.coupang.com/np/search?q=",
    "무신사": "https://www.musinsa.com/search/goods?keyword=",
    "지마켓": "https://browse.gmarket.co.kr/search?keyword=",
    "11번가": "https://search.11st.co.kr/Search.tmall?kwd=",
    "알리": "https://ko.aliexpress.com/wholesale?SearchText=",
    "아마존": "https://www.amazon.com/s?k=",
    "멜론": "https://www.melon.com/search/total/index.htm?q=",
    "인스타그램": "https://www.instagram.com/explore/search/keyword/?q=",
    "인스타": "https://www.instagram.com/explore/search/keyword/?q=",
    "핀터레스트": "https://www.pinterest.co.kr/search/pins/?q=",
    "깃허브": "https://github.com/search?q=",
    "네이버지도": "https://map.naver.com/p/search/",
    "카카오맵": "https://map.kakao.com/?q=",
    "구글맵": "https://www.google.com/maps/search/",
    "야놀자": "https://www.yanolja.com/search/",
    "여기어때": "https://www.goodchoice.kr/product/search/2?keyword=",
    "에어비앤비": "https://www.airbnb.co.kr/s/",
    "아고다": "https://www.agoda.com/ko-kr/search?textToSearch=",
    "네이버": "https://search.naver.com/search.naver?query=",
    "구글": "https://www.google.com/search?q=",
}


def search_url(text, q):
    qq = urllib.parse.quote(q)
    for name in sorted(PLATFORM_SEARCH, key=len, reverse=True):
        if name in text.replace(" ", ""):
            return PLATFORM_SEARCH[name] + qq
    if re.search(r"영상|동영상|노래|음악|뮤비|예능|코미디|브이로그|신곡|음원|"
                 r"듣고\s*싶|틀어|재생", text, re.I):
        return "https://www.youtube.com/results?search_query=" + qq
    return "https://www.google.com/search?q=" + qq


# 공개 SearXNG 인스턴스 (JSON API 허용 — 캡차 없는 진짜 웹 검색). 순서대로 시도.
SEARX_INSTANCES = ["https://searx.be", "https://search.inetol.net",
                   "https://baresearch.org", "https://priv.au"]


def _searxng(q, n, ua):
    """공개 SearXNG에서 일반 웹 검색 결과를 가져온다 (구글/빙 종합, 캡차 없음)."""
    out = []
    for base in SEARX_INSTANCES:
        try:
            url = (base + "/search?q=" + urllib.parse.quote(q)
                   + "&format=json&language=ko&safesearch=0")
            d = json.loads(urllib.request.urlopen(
                urllib.request.Request(url, headers=ua), timeout=8).read())
            for it in d.get("results", [])[:n]:
                title = (it.get("title") or "").strip()
                content = (it.get("content") or "").strip()
                link = (it.get("url") or "").strip()
                if title:
                    out.append(f"- [웹] {title}: {content[:180]}\n  {link}")
            if out:
                return out
        except Exception:
            continue
    return out


def web_search_fetch(q, n=5):
    """여러 공개 소스에서 검색 결과를 모아 모델에게 전달 (캡차·차단 없는 소스만).
    ① SearXNG 일반 웹검색 ② DuckDuckGo 즉답 ③ 구글뉴스 RSS ④ 위키백과.
    실패하면 빈 문자열 → 브라우저 창만 띄우는 모드로 동작."""
    from html import unescape
    ua = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                        "Chrome/126.0 Safari/537.36"}
    out = []
    out += _searxng(q, n, ua)               # ① 진짜 웹 검색 (구글/빙 종합)
    try:                                     # ② DuckDuckGo 즉답 API (정의·요약)
        url = ("https://api.duckduckgo.com/?q=" + urllib.parse.quote(q)
               + "&format=json&no_html=1&kl=kr-kr")
        d = json.loads(urllib.request.urlopen(
            urllib.request.Request(url, headers=ua), timeout=8).read())
        if d.get("AbstractText"):
            out.append(f"- [요약] {d['AbstractText']} ({d.get('AbstractURL', '')})")
        for t in d.get("RelatedTopics", [])[:3]:
            if isinstance(t, dict) and t.get("Text"):
                out.append(f"- [관련] {t['Text']}")
    except Exception:
        pass
    try:   # ③ 최신 소식·시세·이슈
        url = ("https://news.google.com/rss/search?q=" + urllib.parse.quote(q)
               + "&hl=ko&gl=KR&ceid=KR:ko")
        xml = urllib.request.urlopen(urllib.request.Request(url, headers=ua),
                                     timeout=10).read().decode("utf-8", "replace")
        items = re.findall(r"<item><title>([\s\S]*?)</title><link>([\s\S]*?)</link>"
                           r"[\s\S]*?<pubDate>([\s\S]*?)</pubDate>", xml)
        for t, u, d in items[:n]:
            out.append(f"- [뉴스 {d[5:16].strip()}] {unescape(t).strip()}\n  {u.strip()}")
    except Exception:
        pass
    try:   # ④ 개념·인물·정의
        url = ("https://ko.wikipedia.org/w/api.php?action=query&list=search&srsearch="
               + urllib.parse.quote(q) + "&format=json&utf8=1&srlimit=2")
        d = json.loads(urllib.request.urlopen(urllib.request.Request(url, headers=ua),
                                              timeout=10).read())
        for it in d.get("query", {}).get("search", []):
            sn = unescape(re.sub(r"<[^>]+>", "", it.get("snippet", "")))
            out.append(f"- [위키백과] {it['title']}: {sn}")
    except Exception:
        pass
    # 중복 제거 (앞쪽 소스 우선)
    seen, uniq = set(), []
    for line in out:
        key = line[:60]
        if key not in seen:
            seen.add(key)
            uniq.append(line)
    return "\n".join(uniq)


def fetch_url_text(url, max_chars=6000):
    """주어진 URL의 본문 텍스트를 긁어와 정리해 반환 (요약·분석용).
    스크립트·스타일·태그를 제거하고 제목 + 본문 일부를 돌려준다. 실패 시 ''."""
    ua = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                        "Chrome/126.0 Safari/537.36"}
    try:
        from html import unescape
        raw = urllib.request.urlopen(
            urllib.request.Request(url, headers=ua), timeout=15).read()
        # 인코딩 추정 (charset 명시 우선, 없으면 utf-8)
        head = raw[:2048].decode("ascii", "replace")
        m = re.search(r'charset=["\']?([\w\-]+)', head, re.I)
        html = raw.decode(m.group(1) if m else "utf-8", "replace")
    except Exception as e:
        log(f"fetch_url_text 실패({url}): {e}")
        return ""
    title = ""
    mt = re.search(r"<title[^>]*>(.*?)</title>", html, re.S | re.I)
    if mt:
        title = unescape(re.sub(r"\s+", " ", mt.group(1))).strip()
    # 본문만 남기기: 스크립트·스타일·헤더·푸터·네비 제거
    body = re.sub(r"(?is)<(script|style|noscript|svg|head|nav|footer|header|aside|form)[^>]*>"
                  r".*?</\1>", " ", html)
    body = re.sub(r"(?is)<br\s*/?>", "\n", body)
    body = re.sub(r"(?is)</(p|div|li|h[1-6]|tr)>", "\n", body)
    text = unescape(re.sub(r"(?s)<[^>]+>", " ", body))
    lines = [ln.strip() for ln in text.splitlines()]
    text = "\n".join(ln for ln in lines if len(ln) > 1)
    text = re.sub(r"\n{3,}", "\n\n", re.sub(r"[ \t]{2,}", " ", text)).strip()
    if not text:
        return ""
    head_txt = (f"[제목] {title}\n" if title else "") + f"[출처] {url}\n\n"
    return head_txt + text[:max_chars]


# ================= 임베딩 (시맨틱 검색 · 장기기억) =================
def embed(texts):
    """bge-m3로 문장 임베딩 벡터를 얻는다. 실패하면 빈 리스트(키워드 검색으로 폴백)."""
    if isinstance(texts, str):
        texts = [texts]
    texts = [t for t in texts if t and t.strip()]
    if not texts:
        return []
    try:
        req = urllib.request.Request(
            OLLAMA + "/api/embed",
            data=json.dumps({"model": EMBED_MODEL, "input": texts}).encode(),
            headers={"Content-Type": "application/json"})
        d = json.loads(urllib.request.urlopen(req, timeout=30).read())
        return d.get("embeddings", [])
    except Exception as e:
        log(f"embed 실패: {e}")
        return []


def cosine(a, b):
    s = da = db = 0.0
    for x, y in zip(a, b):
        s += x * y
        da += x * x
        db += y * y
    if da == 0 or db == 0:
        return 0.0
    return s / ((da ** 0.5) * (db ** 0.5))


def semantic_rerank(query, candidates, limit=15):
    """키워드로 모은 후보(파일경로 리스트)를 쿼리와의 의미 유사도로 재정렬.
    candidates: [(score, mtime, path), ...] → 같은 형식으로 반환.
    동의어·문맥이 비슷하면 점수를 올려줘 '띄어쓰기/단어 차이'를 넘어선다."""
    if not candidates:
        return candidates
    names = [os.path.splitext(os.path.basename(p))[0].replace("_", " ").replace("-", " ")
             for _, _, p in candidates]
    vecs = embed([query] + names)
    if not vecs or len(vecs) != len(names) + 1:
        return candidates[:limit]                     # 임베딩 불가 → 원래 순서
    qv = vecs[0]
    rescored = []
    for (kw_sc, mt, p), nv in zip(candidates, vecs[1:]):
        sim = cosine(qv, nv)
        rescored.append((kw_sc + sim * 20, mt, p))    # 키워드 점수 + 의미 유사도
    rescored.sort(key=lambda x: (-x[0], -x[1]))
    return rescored[:limit]


def search_history_files(query, k=6):
    """과거 대화 기록(채팅기록/대화_*.json)을 의미(임베딩)로 검색.
    대화색인.json 캐시를 변경분만 갱신하고 쿼리와 가장 가까운 대화 k건을 돌려준다."""
    files = globmod.glob(os.path.join(HIST_DIR, "대화_*.json"))
    cache = {}
    try:
        if os.path.exists(HIST_VEC):
            with open(HIST_VEC, encoding="utf-8") as f:
                cache = json.load(f)
    except Exception:
        cache = {}
    todo, docs, cur = [], [], set()
    for p in files:
        base = os.path.basename(p)
        cur.add(base)
        try:
            mt = os.path.getmtime(p)
        except OSError:
            continue
        if base in cache and cache[base].get("mtime") == mt and cache[base].get("vec"):
            continue
        try:
            with open(p, encoding="utf-8") as f:
                msgs = json.load(f)
        except Exception:
            continue
        first = next((m["content"] for m in msgs if m.get("role") == "user"), "")
        title = first.split("\n")[0][:42] or "(빈 대화)"
        joined = re.sub(r"\s+", " ", " ".join(
            m.get("content", "") for m in msgs
            if m.get("role") in ("user", "assistant"))).strip()[:2000]
        if not joined:
            continue
        todo.append((base, mt, title, joined))
        docs.append(joined)
    if docs:
        vecs = embed(docs)
        if vecs and len(vecs) == len(docs):
            for (base, mt, title, joined), v in zip(todo, vecs):
                cache[base] = {"mtime": mt, "vec": v, "title": title,
                               "preview": joined[:160]}
    for base in list(cache):                  # 삭제된 대화는 캐시에서 제거
        if base not in cur:
            cache.pop(base, None)
    try:
        with open(HIST_VEC, "w", encoding="utf-8") as f:
            json.dump(cache, f, ensure_ascii=False)
    except Exception as e:
        log(f"대화색인 저장 실패: {e}")
    qv = embed(query)
    if not qv:
        return []
    qv = qv[0]
    scored = []
    for base, rec in cache.items():
        if rec.get("vec"):
            scored.append((cosine(qv, rec["vec"]), base, rec.get("title", base),
                           rec.get("preview", "")))
    scored.sort(key=lambda x: -x[0])
    out = []
    for sim, base, title, preview in scored[:k]:
        if sim < 0.30:
            continue
        d, t = base[3:11], base[12:18]
        out.append({"file": base, "title": title,
                    "date": f"{d[:4]}-{d[4:6]}-{d[6:8]} {t[:2]}:{t[2:4]}",
                    "preview": preview, "score": round(sim, 3)})
    return out


# ================= 장기기억 (대화에서 학습 → 자동 회상) =================
MEM_LOCK = threading.Lock()


def memory_load():
    try:
        with open(MEMORY_DB, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []


def memory_add(text):
    """새 사실을 임베딩과 함께 저장. 중복(거의 같은 문장)은 건너뜀."""
    text = (text or "").strip()
    if not text:
        return
    with MEM_LOCK:
        mem = memory_load()
        v = embed(text)
        vec = v[0] if v else []
        for m in mem:                                  # 중복 방지
            if m.get("text") == text or (vec and m.get("vec")
                                         and cosine(vec, m["vec"]) > 0.95):
                return
        mem.append({"text": text, "vec": vec, "date": time.strftime("%Y-%m-%d")})
        mem = mem[-500:]                               # 최대 500개 유지
        try:
            tmp = MEMORY_DB + ".tmp"
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(mem, f, ensure_ascii=False)
            os.replace(tmp, MEMORY_DB)
        except Exception as e:
            log(f"memory_add 저장 실패: {e}")


def memory_recall(query, k=4):
    """질문과 의미가 가까운 기억 k개를 돌려준다 (시스템 프롬프트에 주입)."""
    mem = memory_load()
    if not mem:
        return []
    qv = embed(query)
    if not qv:
        return [m["text"] for m in mem[-k:]]           # 임베딩 불가 → 최근 것
    qv = qv[0]
    scored = [(cosine(qv, m["vec"]), m["text"]) for m in mem if m.get("vec")]
    scored.sort(reverse=True)
    # bge-m3 단문 한국어는 유사도 범위가 압축적 → 임계값 0.40, 상위 k개만
    return [t for s, t in scored[:k] if s > 0.40]


# ================= 음성 인식 (STT — faster-whisper, 로컬) =================
_WHISPER = [None]   # 지연 로딩 (최초 사용 시 모델 적재)


def _get_whisper():
    if _WHISPER[0] is None:
        from faster_whisper import WhisperModel
        import ctranslate2
        # GPU(CUDA) 사용 가능하면 float16으로, 없으면 CPU int8 폴백
        # RTX 2060 기준: small≈480MB / medium≈1.5GB — Ollama와 VRAM 공유하므로 small 유지
        if ctranslate2.get_cuda_device_count() > 0:
            _WHISPER[0] = WhisperModel("small", device="cuda", compute_type="float16")
            log("Whisper: GPU(float16) 모드")
        else:
            _WHISPER[0] = WhisperModel("small", device="cpu", compute_type="int8")
            log("Whisper: CPU(int8) 모드")
    return _WHISPER[0]


def record_until_silence(max_sec=15, silence_sec=1.3, sr=16000):
    """마이크에서 말이 끝날 때까지(또는 max_sec) 녹음해 numpy 배열로 반환."""
    import numpy as np
    import sounddevice as sd
    frames, silent_for, started = [], 0.0, False
    block = 0.1                                        # 100ms 단위로 확인
    with sd.InputStream(samplerate=sr, channels=1, dtype="float32") as stream:
        for _ in range(int(max_sec / block)):
            data, _ = stream.read(int(sr * block))
            frames.append(data.copy())
            amp = float(np.abs(data).mean())
            if amp > 0.012:                            # 말하는 중
                started, silent_for = True, 0.0
            elif started:
                silent_for += block
                if silent_for >= silence_sec:          # 말 끝남
                    break
    if not frames:
        return None
    return np.concatenate(frames, axis=0).flatten()


def transcribe_mic():
    """마이크 녹음 → 한국어 텍스트. 실패 시 빈 문자열."""
    try:
        audio = record_until_silence()
        if audio is None or len(audio) < 1600:
            return ""
        model = _get_whisper()
        segments, _ = model.transcribe(audio, language="ko", beam_size=1)
        return "".join(s.text for s in segments).strip()
    except Exception as e:
        log(f"STT 실패: {e}")
        return ""


# ---------------- 일정 · 메모 (달력 + 알림) ----------------
SCHED_LOCK = threading.Lock()
KDAYS = "월화수목금토일"


def sched_load():
    try:
        with open(SCHED, encoding="utf-8") as f:
            d = json.load(f)
    except Exception:
        d = {}
    d.setdefault("events", [])
    d.setdefault("memos", [])
    return d


def sched_save(d):
    os.makedirs(HOME, exist_ok=True)
    tmp = SCHED + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(d, f, ensure_ascii=False, indent=1)
    os.replace(tmp, SCHED)


def ev_epoch(ev):
    """일정의 기준 시각(epoch). 시간이 없는 종일 일정은 09:00 기준으로 알림 계산"""
    try:
        return time.mktime(time.strptime(
            ev["date"] + " " + (ev.get("time") or "09:00"), "%Y-%m-%d %H:%M"))
    except Exception:
        return None


def kday(datestr):
    try:
        return KDAYS[time.strptime(datestr, "%Y-%m-%d").tm_wday]
    except Exception:
        return ""


def remind_label(off):
    off = int(off)
    if off == 0:
        return "정각"
    if off % 1440 == 0:
        return ("하루" if off == 1440 else f"{off // 1440}일") + " 전"
    if off % 60 == 0:
        return f"{off // 60}시간 전"
    return f"{off}분 전"


REMIND_WORDS = [("일주일", 10080), ("이틀", 2880), ("하루", 1440), ("전날", 1440),
                ("한시간", 60), ("반시간", 30), ("정각", 0), ("당일", 0), ("제시간", 0)]


def parse_reminds(s):
    """'하루 전, 30분 전' 같은 문구 → 분 단위 오프셋 목록 [1440, 30]"""
    s = (s or "").strip()
    if not s or "없" in s:
        return []
    out = set()
    for word, off in REMIND_WORDS:
        if word in s:
            out.add(off)
    for num, unit in re.findall(r"(\d+)\s*(분|시간|일)", s):
        out.add(int(num) * {"분": 1, "시간": 60, "일": 1440}[unit])
    return sorted(out, reverse=True)


def toast(title, body):
    """윈도우 알림 센터 토스트 — 추가 설치 없이 PowerShell WinRT 사용"""
    def psq(s):
        return "'" + str(s).replace("'", "''") + "'"
    try:
        ps = (
            "$null=[Windows.UI.Notifications.ToastNotificationManager,"
            "Windows.UI.Notifications,ContentType=WindowsRuntime];"
            "$t=[Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent("
            "[Windows.UI.Notifications.ToastTemplateType]::ToastText02);"
            "$x=$t.GetElementsByTagName('text');"
            f"$null=$x.Item(0).AppendChild($t.CreateTextNode({psq(title)}));"
            f"$null=$x.Item(1).AppendChild($t.CreateTextNode({psq(body)}));"
            "$n=[Windows.UI.Notifications.ToastNotification]::new($t);"
            "[Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("
            "'{1AC14E77-02E7-4E5D-B744-2EB1AE5198B7}\\WindowsPowerShell\\v1.0\\powershell.exe'"
            ").Show($n)"
        )
        enc = base64.b64encode(ps.encode("utf-16-le")).decode()
        subprocess.Popen(["powershell", "-NoProfile", "-EncodedCommand", enc],
                         creationflags=0x08000000)
    except Exception:
        pass


def handle_work_block(block):
    """모델의 ```work``` 블록 → 워크플로우 보드에 프로젝트/업무 추가. 확인 문구 반환."""
    kv = {}
    for ln in block.strip().splitlines():
        if ":" in ln:
            k, v = ln.split(":", 1)
            kv[k.strip().lstrip("-•· ")] = v.strip()
    kind = kv.get("종류", "업무")
    proj_name = kv.get("프로젝트", "").strip()
    d = wf_load()
    # 프로젝트 찾기/생성
    proj = None
    for p in d["projects"]:
        if p["name"] == proj_name:
            proj = p
            break
    if kind.startswith("프로젝트") or (proj is None and proj_name):
        if proj is None:
            colors = ["#5B6CFF", "#E8473F", "#27AE60", "#F0993E", "#8E44AD", "#16A0A0"]
            proj = {"id": int(time.time() * 1000), "name": proj_name or "새 프로젝트",
                    "color": colors[len(d["projects"]) % len(colors)],
                    "folder": kv.get("폴더", ""), "tasks": []}
            d["projects"].append(proj)
        if kind.startswith("프로젝트") and "업무" not in kv:
            wf_save(d)
            return f"📋 프로젝트 **{proj['name']}** 을(를) 워크플로우에 만들었습니다."
    if proj is None:
        if not d["projects"]:
            return "📋 먼저 프로젝트를 지정하세요 (work 블록에 '프로젝트: 이름')."
        proj = d["projects"][0]
    # 업무 추가 (여러 개면 콤마/줄바꿈 구분)
    raw = kv.get("업무") or kv.get("제목") or kv.get("할일") or ""
    pri = kv.get("중요도") or kv.get("우선순위") or "보통"
    if "높" in pri or "high" in pri.lower() or "긴급" in pri:
        pri = "높음"
    elif "낮" in pri or "low" in pri.lower():
        pri = "낮음"
    else:
        pri = "보통"
    due = kv.get("마감") or kv.get("기한") or ""
    titles = [t.strip() for t in re.split(r"[,\n]", raw) if t.strip()]
    for tt in titles:
        proj.setdefault("tasks", []).append({
            "id": int(time.time() * 1000) + len(proj["tasks"]),
            "title": tt, "status": "todo", "priority": pri, "due": due, "memo": ""})
    wf_save(d)
    if titles:
        return (f"📋 **{proj['name']}** 에 업무 {len(titles)}건을 추가했습니다 "
                f"(중요도: {pri}) — 왼쪽 📋 워크플로우에서 확인하세요.")
    return "📋 워크플로우를 갱신했습니다."


def handle_plan_block(block):
    """모델의 ```plan``` 블록 → 일정/메모 저장·수정. 확인 문구를 돌려준다"""
    kv = {}
    for ln in block.strip().splitlines():
        if ":" in ln:
            k, v = ln.split(":", 1)
            kv[k.strip().lstrip("-•· ")] = v.strip()
    title = kv.get("제목") or kv.get("내용") or kv.get("할일") or ""
    kind = kv.get("종류", "")
    category = kv.get("카테고리", "개인")
    if "회사" in category or "업무" in category or "work" in category.lower():
        category = "회사"
    else:
        category = "개인"

    # ── 일정 수정 ──────────────────────────────────────────────────────
    if kind.startswith("수정") or kind.startswith("변경"):
        ev_id_raw = kv.get("id", "").strip()
        if not ev_id_raw:
            return "✏️ 일정 수정 실패 — id 필드가 없습니다. 일정 목록에서 id를 확인하세요."
        try:
            ev_id = int(ev_id_raw)
        except ValueError:
            return f"✏️ 일정 수정 실패 — id 형식이 잘못됐습니다: {ev_id_raw}"
        fields = {}
        if "날짜" in kv:
            dm = re.search(r"(\d{4})[-./년\s]*(\d{1,2})[-./월\s]*(\d{1,2})", kv["날짜"])
            if dm:
                fields["date"] = f"{dm.group(1)}-{int(dm.group(2)):02d}-{int(dm.group(3)):02d}"
        if "시간" in kv:
            tm2 = re.search(r"(\d{1,2})\s*[:시]\s*(\d{2})?", kv["시간"])
            if tm2:
                fields["time"] = f"{int(tm2.group(1)):02d}:{tm2.group(2) or '00'}"
        if title:
            fields["title"] = title
        if "카테고리" in kv:
            fields["category"] = category
        if "알림" in kv:
            fields["remind"] = parse_reminds(kv["알림"])
            fields["notified"] = []
        if not fields:
            return "✏️ 수정할 내용이 없습니다"
        with SCHED_LOCK:
            d = sched_load()
            for e in d["events"]:
                if e["id"] == ev_id:
                    e.update(fields)
                    sched_save(d)
                    return f"✏️ 일정 수정 완료: {e['title']} ({e['date']})"
        return f"✏️ 일정 수정 실패 — id {ev_id}를 찾을 수 없습니다"

    # ── 메모 ──────────────────────────────────────────────────────────
    if kind.startswith("메모") or (not kv.get("날짜") and not kind.startswith("일정")):
        if not title:
            return "📝 메모 내용이 비어 있어 저장하지 않았습니다"
        with SCHED_LOCK:
            d = sched_load()
            d["memos"].insert(0, {"id": int(time.time() * 1000),
                                  "ts": time.strftime("%Y-%m-%d %H:%M"), "text": title})
            sched_save(d)
        return f"📝 메모 저장: {title}"

    # ── 일정 저장 ─────────────────────────────────────────────────────
    m = re.search(r"(\d{4})[-./년\s]*(\d{1,2})[-./월\s]*(\d{1,2})", kv.get("날짜", ""))
    if not m:
        return f"📅 일정 저장 실패 — 날짜를 알 수 없음: {kv.get('날짜', '(없음)')}"
    date = f"{m.group(1)}-{int(m.group(2)):02d}-{int(m.group(3)):02d}"
    tme = ""
    tm = re.search(r"(\d{1,2})\s*[:시]\s*(\d{2})?", kv.get("시간", ""))
    if tm:
        tme = f"{int(tm.group(1)):02d}:{tm.group(2) or '00'}"
    remind = parse_reminds(kv["알림"]) if "알림" in kv else ([30] if tme else [1440])
    cat_icon = "🏢" if category == "회사" else "👤"
    with SCHED_LOCK:
        d = sched_load()
        d["events"].append({"id": int(time.time() * 1000), "date": date, "time": tme,
                            "title": title or "(제목 없음)", "memo": kv.get("메모", ""),
                            "category": category,
                            "remind": remind, "done": False, "notified": []})
        sched_save(d)
    rl = ", ".join(remind_label(o) for o in remind) if remind else "없음"
    return (f"📅 {cat_icon}[{category}] 일정 저장: {date}({kday(date)}) {tme} {title}"
            f" — 알림: {rl}").replace("  ", " ")


def upcoming_summary(days=14):
    """앞으로 2주 일정 요약 — 시스템 프롬프트에 넣어 모델이 일정을 알게 함"""
    today = time.strftime("%Y-%m-%d")
    end = time.strftime("%Y-%m-%d", time.localtime(time.time() + days * 86400))
    evs = sorted([e for e in sched_load()["events"]
                  if not e.get("done") and today <= e["date"] <= end],
                 key=lambda e: (e["date"], e.get("time") or "99"))
    if not evs:
        return "(2주 내 일정 없음)"
    return "\n".join(
        f"- [id:{e['id']}][{e.get('category','개인')}] {e['date']}({kday(e['date'])}) "
        f"{e.get('time') or '종일'} {e['title']}"
        for e in evs[:20])


# ---------------- 금고 ----------------
def _vault_key(master):
    from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
    from cryptography.hazmat.primitives import hashes
    salt = hashlib.sha256(b"ai-bisearch-vault-salt").digest()[:16]
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=600000)
    return base64.urlsafe_b64encode(kdf.derive(master.encode()))


def vault_load(master):
    from cryptography.fernet import Fernet, InvalidToken
    if not os.path.exists(VAULT):
        return {}
    with open(VAULT, "rb") as f:
        data = f.read()
    try:
        return json.loads(Fernet(_vault_key(master)).decrypt(data).decode())
    except InvalidToken:
        return None


def vault_save(master, d):
    from cryptography.fernet import Fernet
    token = Fernet(_vault_key(master)).encrypt(json.dumps(d, ensure_ascii=False).encode())
    with open(VAULT, "wb") as f:
        f.write(token)


# ---------------- 파일 분석 ----------------
_OCR = {"reader": None, "tried": False}


def ocr_image(path):
    """이미지에서 글자 추출 (easyocr, 한국어+영어). 라이브러리 없으면 ''.
    스캔 문서·스크린샷의 텍스트를 검색·요약 가능하게 한다 (비전 모델 보조)."""
    if _OCR["reader"] is None:
        if _OCR["tried"]:
            return ""
        _OCR["tried"] = True
        try:
            import easyocr
            try:
                _OCR["reader"] = easyocr.Reader(["ko", "en"], gpu=True, verbose=False)
            except Exception:
                _OCR["reader"] = easyocr.Reader(["ko", "en"], gpu=False, verbose=False)
        except Exception as e:
            log(f"easyocr 미설치/초기화 실패(OCR 건너뜀): {e}")
            return ""
    try:
        lines = _OCR["reader"].readtext(path, detail=0, paragraph=True)
        return "\n".join(l.strip() for l in lines if l.strip())[:4000]
    except Exception as e:
        log(f"OCR 실패({path}): {e}")
        return ""


def analyze_file(path):
    ext = os.path.splitext(path)[1].lower()
    name = os.path.basename(path)
    try:
        if ext in IMG_EXT:
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            txt = ocr_image(path)     # 글자가 있으면 OCR로 추출해 함께 전달 (없으면 비전만)
            summary = (f"[이미지 글자(OCR): {name}]\n{txt}\n경로: {path}") if txt else None
            return summary, b64
        if ext == ".dxf":
            import ezdxf
            doc = ezdxf.readfile(path)
            counts = {}
            for e in doc.modelspace():
                counts[e.dxftype()] = counts.get(e.dxftype(), 0) + 1
            layers = [l.dxf.name for l in doc.layers]
            top = sorted(counts.items(), key=lambda x: -x[1])[:15]
            return (f"[DXF 도면 분석: {name}]\n레이어({len(layers)}개): {', '.join(layers[:30])}\n"
                    f"객체: {', '.join(f'{k}×{v}' for k, v in top)}\n파일 경로: {path}"), None
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            out = [f"[엑셀: {name}] 시트: {', '.join(wb.sheetnames)}"]
            for sn in wb.sheetnames[:3]:
                rows = []
                for i, row in enumerate(wb[sn].iter_rows(values_only=True)):
                    if i >= 12: break
                    rows.append(" | ".join("" if c is None else str(c) for c in row[:12]))
                out.append(f"--- {sn} ---\n" + "\n".join(rows))
            out.append(f"파일 경로: {path}")
            return "\n".join(out), None
        if ext == ".docx":
            import docx
            text = "\n".join(p.text for p in docx.Document(path).paragraphs if p.text.strip())[:3000]
            return f"[워드: {name}]\n{text}\n파일 경로: {path}", None
        if ext == ".pdf":
            from pypdf import PdfReader
            r = PdfReader(path)
            text = "".join((p.extract_text() or "") for p in r.pages[:5])[:3000]
            return f"[PDF: {name}, {len(r.pages)}p]\n{text}\n파일 경로: {path}", None
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for i, slide in enumerate(prs.slides[:20], 1):
                texts = [sh.text.strip() for sh in slide.shapes
                         if sh.has_text_frame and sh.text.strip()]
                if texts:
                    chunks.append(f"[슬라이드 {i}] " + " / ".join(texts))
            body = "\n".join(chunks)[:3500]
            return f"[PPT: {name}, {len(prs.slides)}장]\n{body}\n파일 경로: {path}", None
        if ext == ".hwpx":          # 한글 신형식 = zip 안의 XML
            import zipfile
            from html import unescape
            texts = []
            with zipfile.ZipFile(path) as z:
                for nm in sorted(z.namelist()):
                    if nm.startswith("Contents/") and nm.endswith(".xml"):
                        raw = z.read(nm).decode("utf-8", "replace")
                        texts.append(re.sub(r"<[^>]+>", "", unescape(raw)))
            body = re.sub(r"\s+\n", "\n", "\n".join(texts)).strip()[:3500]
            return f"[한글(HWPX): {name}]\n{body}\n파일 경로: {path}", None
        if ext == ".hwp":           # 한글 구형식(OLE) — PrvText 스트림에서 미리보기 텍스트 추출
            import olefile
            ole = olefile.OleFileIO(path)
            body = ""
            if ole.exists("PrvText"):
                body = ole.openstream("PrvText").read().decode("utf-16-le", "replace")
            ole.close()
            body = body.strip()[:3500]
            return (f"[한글(HWP): {name}]\n{body or '(미리보기 텍스트 없음)'}\n"
                    f"파일 경로: {path}"), None
        if ext == ".skp":
            return (f"[스케치업 파일: {name}] .skp는 직접 파싱 불가. "
                    f"수정 필요 시 Ruby 스크립트 작성해 안내.\n파일 경로: {path}"), None
        if ext in {".txt", ".csv", ".md", ".rb", ".py", ".json"}:
            with open(path, encoding="utf-8", errors="replace") as f:
                return f"[파일: {name}]\n{f.read(4000)}\n파일 경로: {path}", None
        return f"[파일: {name}, 형식 {ext}] 미리보기 불가.\n파일 경로: {path}", None
    except Exception as e:
        return f"[파일: {name}] 분석 실패: {e}", None


HTML = r"""<!DOCTYPE html>
<html lang="ko"><head><meta charset="utf-8">
<style>
*{box-sizing:border-box;margin:0;padding:0}
:root{
  --bg:#FAF9F5; --side:#F0EEE6; --line:#E3DFD3; --line2:#D8D3C4;
  --txt:#3D3929; --mut:#8A8470; --accent:#D97757; --accent-h:#C05F3C;
  --user-bg:#F0EEE6; --code-bg:#F7F6F1; --white:#FFFFFF; --fs:15px;
}
body.dark{
  --bg:#1E1E1C; --side:#26261F; --line:#3A392F; --line2:#46453A;
  --txt:#ECE9DE; --mut:#9C967F; --accent:#E08A6A; --accent-h:#E4A284;
  --user-bg:#2E2D26; --code-bg:#161613; --white:#2A2A24;
}
html,body{height:100%}
body{font-family:'Segoe UI','Malgun Gothic',sans-serif;background:var(--bg);color:var(--txt);
  display:flex;overflow:hidden;font-size:var(--fs)}
::-webkit-scrollbar{width:8px;height:8px}
::-webkit-scrollbar-thumb{background:var(--line2);border-radius:4px}
::-webkit-scrollbar-track{background:transparent}
a{color:var(--accent)}

/* ---------- 사이드바 ---------- */
#side{width:264px;min-width:264px;background:var(--side);border-right:1px solid var(--line);
  display:flex;flex-direction:column;transition:margin .2s}
#side.hidden{margin-left:-264px}
.side-top{padding:14px 14px 8px}
.brand{display:flex;align-items:center;gap:9px;font-family:Georgia,'Malgun Gothic',serif;
  font-size:17px;font-weight:600;padding:4px 6px 14px}
.brand .star{color:var(--accent);font-size:20px;line-height:1}
#newchat,#calbtn,#memobtn,#wfbtn,#cleanbtn,#shopbtn{width:100%;border:1px solid var(--line2);background:var(--white);color:var(--txt);
  border-radius:10px;padding:9px 12px;font-size:13.5px;font-family:inherit;cursor:pointer;
  display:flex;align-items:center;gap:8px;transition:.15s}
#newchat:hover,#calbtn:hover,#memobtn:hover,#wfbtn:hover,#cleanbtn:hover,#shopbtn:hover{border-color:var(--accent);color:var(--accent)}
#calbtn,#memobtn,#wfbtn,#cleanbtn,#shopbtn{margin-top:6px}
.side-label{font-size:11.5px;color:var(--mut);padding:14px 20px 6px;letter-spacing:.5px}
#histlist{flex:1;overflow-y:auto;padding:0 8px}
.hitem{display:flex;align-items:center;gap:4px;border-radius:9px;padding:8px 10px;cursor:pointer;transition:.12s}
.hitem:hover{background:#E7E3D5}
.hitem.cur{background:#D8D3C4}
body.dark .hitem:hover{background:#3A392F}
body.dark .hitem.cur{background:#46453A}
.hitem .ht{flex:1;min-width:0;font-size:13px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:var(--txt)}
.hitem .hd{font-size:10.5px;color:var(--mut);margin-top:2px}
.hitem .hdel{background:none;border:none;color:transparent;cursor:pointer;font-size:13px;flex-shrink:0;padding:2px}
.hitem:hover .hdel{color:var(--mut)}
.hitem .hdel:hover{color:var(--accent)}
.hitem .htedit{flex:1;min-width:0;font-size:13px;font-family:inherit;color:var(--txt);
  background:#fff;border:1px solid var(--accent);border-radius:6px;padding:3px 6px;outline:none}
.side-bottom{border-top:1px solid var(--line);padding:10px 14px;display:flex;align-items:center;gap:8px}
.fbtn{border:1px solid var(--line2);background:var(--white);color:var(--mut);border-radius:7px;
  padding:3px 9px;font-size:12px;cursor:pointer;font-family:inherit}
.fbtn:hover{color:var(--accent);border-color:var(--accent)}
/* ---------- 좌측하단 설정 팝업 ---------- */
#settingspop{display:none;position:absolute;bottom:56px;left:10px;z-index:200;
  background:var(--white);border:1px solid var(--line2);border-radius:12px;
  box-shadow:0 4px 18px rgba(0,0,0,.18);padding:6px;min-width:180px;flex-direction:column;gap:2px}
#settingspop.open{display:flex}
.spitem{border:none;background:transparent;color:var(--txt);border-radius:8px;padding:9px 12px;
  font-size:13px;font-family:inherit;cursor:pointer;text-align:left;display:flex;align-items:center;gap:8px;width:100%}
.spitem:hover{background:var(--side);color:var(--accent)}
#settingsmenu{position:relative;display:inline-flex;align-items:center;flex-shrink:0}
/* ---------- 우선순위 뱃지 버튼 ---------- */
.wpri{display:inline-block;font-size:10.5px;border-radius:5px;padding:1px 7px;cursor:pointer;
  border:1px solid transparent;font-family:inherit;font-weight:600;transition:.12s;flex-shrink:0}
.wpri:hover{opacity:.8}
.wpri-높음{background:#fdecea;color:#e8473f;border-color:#f5c2be}
.wpri-보통{background:#fff3e0;color:#f0993e;border-color:#f5d99a}
.wpri-낮음{background:#f0f1f3;color:#9aa0aa;border-color:#d0d3da}
body.dark .wpri-높음{background:#3a1a18;color:#e8473f;border-color:#6a2e2a}
body.dark .wpri-보통{background:#3a2e18;color:#f0993e;border-color:#6a5a2a}
body.dark .wpri-낮음{background:#2a2c30;color:#9aa0aa;border-color:#44464e}
#status{flex:1;text-align:right;font-size:11px;color:var(--mut);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
#status.ok{color:#7A9B5C}#status.err{color:#C0392B}
.genmeta{font-size:10.5px;color:var(--mut);margin-top:5px;opacity:.7}

/* ---------- 업데이트 배너 ---------- */
#updbanner{display:none;background:#FFF8E1;border-bottom:1px solid #F5CC50;
  padding:9px 18px;align-items:center;gap:10px;font-size:13px;z-index:50;flex-shrink:0}
#updbanner.show{display:flex}
#updbanner .updico{font-size:18px;flex-shrink:0}
#updbanner .updtxt{flex:1;line-height:1.4}
#updbanner .updtxt b{color:#7A5A00}
#updbanner .updlog{color:var(--mut);font-size:12px}
.updbtn{border:none;border-radius:7px;padding:5px 14px;font-size:12.5px;
  cursor:pointer;font-family:inherit;font-weight:600;flex-shrink:0}
.updbtn.do{background:var(--accent);color:#fff}
.updbtn.do:hover{background:var(--accent-h)}
.updbtn.skip{background:transparent;border:1px solid var(--line2);color:var(--mut)}
.updbtn.skip:hover{border-color:var(--accent);color:var(--accent)}

/* ---------- 메인 ---------- */
main{flex:1;display:flex;flex-direction:column;min-width:0;position:relative}
#sidetoggle{position:absolute;top:12px;left:12px;z-index:5;border:none;background:transparent;
  color:var(--mut);font-size:17px;cursor:pointer;border-radius:8px;padding:4px 9px}
#sidetoggle:hover{background:var(--side);color:var(--txt)}
#chatwrap{flex:1;overflow-y:auto;padding:36px 24px 12px}
#chat{max-width:760px;margin:0 auto;display:flex;flex-direction:column;gap:22px}
#greet{max-width:760px;margin:14vh auto 0;text-align:center;display:none}
#greet .gstar{color:var(--accent);font-size:42px}
#greet h2{font-family:Georgia,'Malgun Gothic',serif;font-weight:500;font-size:28px;margin-top:14px;color:var(--txt)}
#greet p{color:var(--mut);margin-top:10px;font-size:14px;line-height:1.8}

.row{display:flex;gap:14px}
.row.user{justify-content:flex-end}
.user-msg-wrap{display:flex;flex-direction:column;align-items:flex-end;gap:3px;max-width:78%}
.row.user .bubble{background:var(--user-bg);border-radius:16px;border-bottom-right-radius:5px;
  padding:11px 16px;line-height:1.75;white-space:pre-wrap;word-break:break-word}
.msg-actions{display:flex;gap:4px;opacity:0;transition:opacity .15s;pointer-events:none}
.row.user:hover .msg-actions{opacity:1;pointer-events:auto}
.msg-action-btn{border:1px solid var(--line);background:var(--white);color:var(--mut);
  border-radius:6px;padding:2px 8px;font-size:11px;cursor:pointer;font-family:inherit;white-space:nowrap}
.msg-action-btn:hover{color:var(--accent);border-color:var(--accent)}
.row.ai{align-items:flex-start}
.row.ai .av{width:28px;height:28px;border-radius:50%;background:var(--accent);color:#fff;
  display:flex;align-items:center;justify-content:center;font-size:15px;flex-shrink:0;margin-top:2px}
.row.ai .bubble{flex:1;min-width:0;line-height:1.8;word-break:break-word;padding-top:3px}
.sys{align-self:center;text-align:center;max-width:95%;font-size:12px;color:var(--mut);
  background:var(--side);border-radius:10px;padding:5px 14px;white-space:pre-wrap}

/* 마크다운 */
.md p{margin:0}
.md .sp{height:10px}
.md h2,.md h3,.md h4{font-family:Georgia,'Malgun Gothic',serif;font-weight:600;margin:14px 0 6px}
.md h2{font-size:1.25em}.md h3{font-size:1.13em}.md h4{font-size:1.04em}
.md ul,.md ol{margin:4px 0 4px 22px}
.md li{margin:3px 0}
.md hr{border:none;border-top:1px solid var(--line);margin:12px 0}
.md blockquote{border-left:3px solid var(--line2);padding:2px 12px;color:var(--mut);margin:6px 0}
.md code{background:var(--code-bg);border:1px solid var(--line);border-radius:5px;
  padding:1px 6px;font-family:Consolas,monospace;font-size:.88em}
.md b{font-weight:600}
.code{border:1px solid var(--line);border-radius:10px;margin:10px 0;overflow:hidden;background:var(--code-bg)}
.codehead{display:flex;justify-content:space-between;align-items:center;background:var(--side);
  border-bottom:1px solid var(--line);padding:5px 12px;font-size:11.5px;color:var(--mut);font-family:Consolas,monospace}
.codehead button{border:none;background:none;color:var(--mut);cursor:pointer;font-size:11.5px;font-family:inherit}
.codehead button:hover{color:var(--accent)}
.code pre{padding:12px 14px;overflow-x:auto;font-family:Consolas,monospace;font-size:12.5px;
  line-height:1.6;white-space:pre;color:#4A4435}

/* 생각 중 표시 */
.thinking{display:flex;align-items:center;gap:10px;color:var(--mut);font-size:13.5px}
.thinking .tstar{color:var(--accent);font-size:18px;display:inline-block;animation:spin 1.6s linear infinite}
@keyframes spin{to{transform:rotate(360deg)}}

/* ---------- 입력부 ---------- */
footer{flex-shrink:0;padding:6px 24px 14px}
.fwrap{max-width:760px;margin:0 auto}
#qbar{display:flex;gap:6px;flex-wrap:wrap;padding:0 2px 8px}
.qbtn{background:var(--side);border:1px solid var(--line);color:var(--txt);border-radius:18px;
  padding:5px 13px;font-size:12.5px;cursor:pointer;display:inline-flex;align-items:center;gap:7px;transition:.15s}
.qbtn:hover{border-color:var(--accent)}
.qbtn .qx{color:transparent;cursor:pointer;font-size:11px}
.qbtn:hover .qx{color:var(--mut)}
.qbtn .qx:hover{color:#C0392B}
.qadd{background:none;border:1px dashed var(--line2);color:var(--mut);border-radius:18px;
  padding:5px 13px;font-size:12.5px;cursor:pointer;font-family:inherit}
.qadd:hover{color:var(--accent);border-color:var(--accent)}
#attach-bar{font-size:12px;color:var(--mut);padding:0 2px 6px;display:none}
.filechip{display:inline-flex;align-items:center;gap:5px;background:var(--side);border:1px solid var(--line);
  border-radius:8px;padding:2px 9px;font-size:11.5px;margin:2px 4px 0 0}
.filechip img.thumb{width:34px;height:34px;object-fit:cover;border-radius:5px;
  border:1px solid var(--line);background:#fff;vertical-align:middle}
#inputbox{background:var(--white);border:1px solid var(--line2);border-radius:16px;
  box-shadow:0 2px 10px rgba(61,57,41,.06);padding:10px 12px 8px;transition:border .15s}
#inputbox:focus-within{border-color:#C9C2AE;box-shadow:0 2px 14px rgba(61,57,41,.1)}
#inp{width:100%;border:none;outline:none;resize:none;font-family:inherit;font-size:var(--fs);
  line-height:1.6;max-height:170px;background:transparent;color:var(--txt)}
#inp::placeholder{color:#B5AF9C}
.inrow{display:flex;align-items:center;gap:6px;padding-top:6px}
.ibtn{border:none;background:none;color:var(--mut);width:32px;height:32px;border-radius:8px;
  cursor:pointer;font-size:16px;display:flex;align-items:center;justify-content:center}
.ibtn:hover{background:var(--side);color:var(--txt)}
.flex{flex:1}
#send{border:none;background:var(--accent);color:#fff;width:34px;height:34px;border-radius:50%;
  cursor:pointer;font-size:17px;display:flex;align-items:center;justify-content:center;transition:.15s}
#send:hover{background:var(--accent-h)}
#send:disabled{background:var(--line2);cursor:default}
#send.stopping{background:#C0392B}
#send.stopping:hover{background:#A93226}
.disclaim{text-align:center;font-size:11px;color:#B5AF9C;padding-top:8px}

/* ---------- 일정 · 메모 (달력) ---------- */
.calbox{width:720px;max-width:94vw;max-height:90vh;display:flex;flex-direction:column;overflow:hidden}
.calhead{display:flex;align-items:center;flex-wrap:wrap;gap:6px 6px;margin-bottom:10px}
.calhead .cym{font-family:Georgia,'Malgun Gothic',serif;font-size:16px;font-weight:600;min-width:108px;text-align:center}
.calhead .chttl{font-size:15px;font-weight:600;white-space:nowrap}
.calhead .flex{flex:1}
.memobox{width:480px;max-width:94vw;max-height:90vh;display:flex;flex-direction:column;overflow:hidden}
.cnav{border:1px solid var(--line2);background:var(--white);border-radius:7px;min-width:26px;height:26px;
  cursor:pointer;color:var(--mut);font-size:14px;font-family:inherit}
.cnav:hover{color:var(--accent);border-color:var(--accent)}
#calgrid{display:grid;grid-template-columns:repeat(7,1fr);gap:3px;user-select:none}
.cdow{text-align:center;font-size:11px;color:var(--mut);padding:3px 0}
.cdow.sun{color:#C0392B}
.cday{min-height:48px;border:1px solid transparent;border-radius:8px;padding:4px 5px;cursor:pointer;font-size:12px}
.cday:hover{background:var(--side)}
.cday.sun .dn{color:#C0392B}
.cday.today .dn{background:var(--accent);color:#fff;border-radius:50%;display:inline-flex;
  width:19px;height:19px;align-items:center;justify-content:center}
.cday.sel{border-color:var(--accent);background:#FBF1EC}
.cday .dotc{display:block;font-size:10px;color:var(--accent);margin-top:2px;overflow:hidden;
  white-space:nowrap;text-overflow:ellipsis}
.cday .dotc.alldone{color:var(--mut);text-decoration:line-through}
#daylist{margin-top:10px;max-height:160px;overflow-y:auto;border-top:1px solid var(--line);padding-top:8px}
.evitem{display:flex;align-items:center;gap:8px;padding:6px 4px;border-radius:8px;font-size:13px}
.evitem:hover{background:var(--side)}
.evitem .et{color:var(--accent);font-size:12px;min-width:40px;flex-shrink:0}
.evitem .ettl{flex:1;min-width:0;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.evitem.done .ettl{text-decoration:line-through;color:var(--mut)}
.evitem .erem{font-size:11px;color:var(--mut);white-space:nowrap}
.evbtn{border:none;background:none;cursor:pointer;color:var(--mut);font-size:13px;flex-shrink:0}
.evbtn:hover{color:#C0392B}
.evitem .ecat{font-size:10px;border-radius:4px;padding:1px 6px;font-weight:600;flex-shrink:0}
.ecat.work{background:#D5E8FF;color:#1A6BC4}
.ecat.pers{background:#E8F5E9;color:#2E7D32}
.cfilter{display:flex;flex-wrap:wrap;gap:4px;margin-bottom:8px}
.cfbtn{border:1px solid var(--line2);background:var(--white);color:var(--mut);border-radius:6px;
  padding:2px 9px;font-size:11.5px;cursor:pointer;font-family:inherit;transition:.12s}
.cfbtn.active{background:var(--accent);color:#fff;border-color:var(--accent)}
.cfbtn.cfadd{border-style:dashed;color:var(--accent)}
.cfbtn.cfadd:hover{background:var(--accent);color:#fff;border-style:solid}
.addrow{display:flex;flex-wrap:wrap;gap:6px;margin-top:10px;align-items:center}
.addrow input,.addrow select{border:1px solid var(--line2);border-radius:8px;padding:7px 9px;
  font-size:13px;font-family:inherit;background:var(--bg);color:var(--txt);outline:none;margin:0;width:auto}
.addrow input:focus{border-color:var(--accent)}
#evtitle{flex:1 1 100%}
#memotext{flex:1;min-width:0}
#evtime{flex:0 0 auto}
#evcat,#evremind{flex:1 1 120px;min-width:0}
#evadd{flex:0 0 auto}
.badd{border:none;background:var(--accent);color:#fff;border-radius:8px;padding:0 14px;
  cursor:pointer;font-size:13px;font-family:inherit}
.badd:hover{background:var(--accent-h)}
#memolist{max-height:340px;overflow-y:auto}
.mitem{display:flex;gap:8px;align-items:flex-start;padding:7px 4px;border-bottom:1px solid var(--line);font-size:13px}
.mitem .mts{color:var(--mut);font-size:11px;white-space:nowrap;margin-top:2px}
.mitem .mtx{flex:1;white-space:pre-wrap;word-break:break-word}

/* ---------- 모달 ---------- */
.overlay{display:none;position:fixed;inset:0;background:rgba(61,57,41,.35);
  align-items:center;justify-content:center;z-index:20}
#cmodal,#imodal{z-index:200}
.box{background:var(--white);border:1px solid var(--line);border-radius:14px;padding:22px;width:340px;
  box-shadow:0 8px 30px rgba(61,57,41,.18)}
.box h3{font-family:Georgia,'Malgun Gothic',serif;font-size:15px;margin-bottom:12px}
.box p{font-size:13.5px;line-height:1.7;white-space:pre-wrap}
.box input{width:100%;border:1px solid var(--line2);border-radius:8px;padding:9px 11px;
  font-size:13.5px;outline:none;font-family:inherit;background:var(--bg);color:var(--txt);margin-bottom:8px}
.box input:focus{border-color:var(--accent)}
.box .brow{display:flex;gap:8px;margin-top:10px;justify-content:flex-end}
.box .brow button{border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-family:inherit;font-size:13px}
.bok{background:var(--accent);color:#fff}.bok:hover{background:var(--accent-h)}
.bcancel{background:var(--side);color:var(--mut)}

/* ---------- 기능 대시보드 (할 수 있는 일) ---------- */
.featbox{width:740px;max-width:94vw;max-height:90vh;display:flex;flex-direction:column;overflow:hidden}
.feathead{display:flex;align-items:center;gap:8px;margin-bottom:4px}
.feathead h3{flex:1;margin:0}
.featsub{font-size:12px;color:var(--mut);margin-bottom:14px}
.featgrid{overflow-y:auto;display:grid;grid-template-columns:1fr 1fr;gap:12px;padding-right:4px}
.featcat{border:1px solid var(--line);border-radius:12px;padding:12px 14px;background:var(--bg)}
.featcat h4{font-size:13.5px;margin:0 0 8px;display:flex;align-items:center;gap:7px}
.featcat ul{list-style:none;margin:0;padding:0}
.featcat li{font-size:12.5px;line-height:1.5;padding:6px 0;color:var(--txt);border-top:1px dashed var(--line)}
.featcat li:first-child{border-top:none}
.fex{display:inline-block;margin-top:4px;font-size:11.5px;color:var(--accent);cursor:pointer;
  background:var(--side);border:1px solid var(--line);border-radius:11px;padding:1px 9px;transition:.12s}
.fex:hover{background:var(--accent);color:#fff;border-color:var(--accent)}
.featnote{font-size:11px;color:var(--mut);padding-top:12px;text-align:center;border-top:1px solid var(--line);margin-top:10px}
@media(max-width:560px){.featgrid{grid-template-columns:1fr}}

/* ---------- 폴더정리·쇼핑 대시보드 모달 ---------- */
.dashbox{width:560px;max-width:94vw;display:flex;flex-direction:column;gap:2px}
.dashsub{font-size:12.5px;color:var(--mut);line-height:1.55;margin:4px 2px 6px}
.dashsub b{color:var(--txt)}
.dashsec{font-size:12px;font-weight:700;color:var(--txt);margin:10px 2px 6px}
.cleangrid{display:grid;grid-template-columns:1fr 1fr;gap:8px}
.cleancard{display:flex;align-items:center;gap:8px;border:1px solid var(--line2);background:var(--white);
  border-radius:10px;padding:10px 12px;cursor:pointer;transition:.13s;font-size:13px;color:var(--txt);text-align:left}
.cleancard:hover{border-color:var(--accent);color:var(--accent)}
.cleancard .cpath{font-size:10.5px;color:var(--mut);margin-top:2px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.cleancard .cleft{flex:1;min-width:0}
.dashbtn{margin-top:10px;border:1px dashed var(--line2);background:var(--white);color:var(--txt);
  border-radius:10px;padding:10px 12px;font-size:13px;font-family:inherit;cursor:pointer;transition:.13s;width:100%}
.dashbtn:hover{border-color:var(--accent);color:var(--accent)}
.dashbtn.primary{border-style:solid;background:var(--accent);color:#fff;border-color:var(--accent);font-weight:600}
.dashbtn.primary:hover{background:var(--accent-h);color:#fff}
.dashinput{width:100%;box-sizing:border-box;border:1px solid var(--line2);border-radius:10px;
  padding:11px 13px;font-size:14px;font-family:inherit;background:var(--bg);color:var(--txt)}
.dashresult{margin-top:12px;font-size:12.5px;color:var(--txt);line-height:1.6;white-space:pre-wrap;
  max-height:200px;overflow-y:auto}
.shopsites{display:flex;flex-wrap:wrap;gap:8px}
.shopchk{display:flex;align-items:center;gap:5px;border:1px solid var(--line2);border-radius:8px;
  padding:6px 11px;font-size:12.5px;color:var(--txt);cursor:pointer}
.shopex{display:flex;flex-wrap:wrap;gap:6px;margin-top:10px}
.shoptag{border:1px solid var(--line2);border-radius:14px;padding:3px 11px;font-size:11.5px;
  color:var(--mut);cursor:pointer;transition:.12s}
.shoptag:hover{border-color:var(--accent);color:var(--accent)}

/* ---------- 피카츄 캐릭터 ---------- */
#pika{position:fixed;right:26px;bottom:24px;width:120px;z-index:9000;
  user-select:none;cursor:grab;display:none;flex-direction:column;align-items:center;
  filter:drop-shadow(0 6px 14px rgba(0,0,0,.25))}
#pika.dragging{cursor:grabbing;filter:drop-shadow(0 12px 22px rgba(0,0,0,.35))}
#pika .bubble{position:relative;background:#fff;color:#1f1f1f;border:2px solid #2b2b2b;
  border-radius:14px;padding:8px 12px;font-size:12.5px;line-height:1.4;max-width:180px;
  text-align:center;margin-bottom:8px;box-shadow:0 3px 10px rgba(0,0,0,.12);
  opacity:0;transform:translateY(6px) scale(.96);transition:.25s;pointer-events:none;font-weight:600}
#pika .bubble.show{opacity:1;transform:none}
#pika .bubble:after{content:"";position:absolute;left:50%;bottom:-9px;transform:translateX(-50%);
  border:7px solid transparent;border-top-color:#2b2b2b}
#pika .bubble:before{content:"";position:absolute;left:50%;bottom:-6px;transform:translateX(-50%);
  border:6px solid transparent;border-top-color:#fff;z-index:1}
body.dark #pika .bubble{background:#2a2a2a;color:#f0f0f0;border-color:#f4d03f}
body.dark #pika .bubble:before{border-top-color:#2a2a2a}
body.dark #pika .bubble:after{border-top-color:#f4d03f}

/* 몸체 */
.pchar{position:relative;width:84px;height:78px;animation:pidle 3s ease-in-out infinite}
#pika.st-work .pchar{animation:pwork .6s ease-in-out infinite}
#pika.st-done .pchar{animation:pdone .5s ease-in-out 3}
@keyframes pidle{0%,100%{transform:translateY(0)}50%{transform:translateY(-5px)}}
@keyframes pwork{0%,100%{transform:translateY(0) rotate(-2deg)}50%{transform:translateY(-3px) rotate(2deg)}}
@keyframes pdone{0%,100%{transform:translateY(0) scale(1)}50%{transform:translateY(-9px) scale(1.06)}}
/* 귀 */
.pear{position:absolute;top:-24px;width:14px;height:42px;background:#f9d423;border-radius:50% 50% 40% 40%;
  border:2px solid #3a2f0b}
.pear.l{left:14px;transform:rotate(-18deg);transform-origin:bottom}
.pear.r{right:14px;transform:rotate(18deg);transform-origin:bottom}
.pear:after{content:"";position:absolute;top:-2px;left:-2px;right:-2px;height:15px;background:#2b2b2b;
  border-radius:50% 50% 40% 40%}
/* 얼굴 */
.pface{position:absolute;top:6px;left:6px;width:72px;height:64px;background:#f9d423;
  border-radius:50% 50% 48% 48%;border:2.5px solid #3a2f0b}
/* 눈 */
.peye{position:absolute;top:20px;width:14px;height:16px;background:#2b2b2b;border-radius:50%}
.peye.l{left:14px}.peye.r{right:14px}
.peye:after{content:"";position:absolute;top:2px;left:2px;width:5px;height:5px;background:#fff;border-radius:50%}
/* 볼(빨간 볼) */
.pcheek{position:absolute;top:32px;width:15px;height:13px;background:#e8473f;border-radius:50%;opacity:.92}
.pcheek.l{left:4px}.pcheek.r{right:4px}
/* 입 */
.pmouth{position:absolute;top:40px;left:50%;transform:translateX(-50%);width:14px;height:7px;
  border:2.2px solid #3a2f0b;border-top:none;border-radius:0 0 14px 14px}
/* 표정: 일하는 중 — 눈 반짝, 입 오므림 */
#pika.st-work .peye{height:13px;top:21px}
#pika.st-work .pmouth{width:8px;height:8px;border:2.2px solid #3a2f0b;border-radius:50%}
/* 표정: 완료 — 눈 ^^ , 입 활짝 */
#pika.st-done .peye{height:8px;border-radius:0;background:transparent;border:2.5px solid #2b2b2b;
  border-bottom:none;border-left:none;border-right:none;top:22px;width:13px;
  border-top-left-radius:10px;border-top-right-radius:10px;border-bottom:2.5px solid #2b2b2b;border-top:none}
#pika.st-done .peye:after{display:none}
#pika.st-done .pmouth{width:18px;height:11px;background:#e8473f;border:2.2px solid #3a2f0b;border-top:none}
#pika .pclose{position:absolute;top:-6px;right:-6px;width:18px;height:18px;border-radius:50%;
  background:#fff;border:1.5px solid #888;color:#555;font-size:11px;line-height:15px;text-align:center;
  cursor:pointer;opacity:0;transition:.15s;z-index:2}
#pika:hover .pclose{opacity:1}
#pika .pclose:hover{background:#e8473f;color:#fff;border-color:#e8473f}

/* ---------- 설정 패널 ---------- */
.setrow{display:flex;align-items:center;justify-content:space-between;gap:10px;
  padding:10px 2px;border-bottom:1px solid var(--line);font-size:13.5px;color:var(--txt)}
.setrow select{border:1px solid var(--line2);border-radius:8px;padding:5px 9px;font-family:inherit;
  font-size:13px;background:var(--bg);color:var(--txt)}
.setsec{font-size:13px;font-weight:700;margin:14px 0 8px;color:var(--txt)}
.sw{position:relative;display:inline-block;width:42px;height:23px}
.sw input{opacity:0;width:0;height:0}
.sw .slider{position:absolute;inset:0;background:#ccc;border-radius:23px;transition:.2s;cursor:pointer}
.sw .slider:before{content:"";position:absolute;height:17px;width:17px;left:3px;top:3px;background:#fff;
  border-radius:50%;transition:.2s}
.sw input:checked+.slider{background:var(--accent)}
.sw input:checked+.slider:before{transform:translateX(19px)}
.favitem{display:flex;align-items:center;gap:6px;padding:6px 8px;border:1px solid var(--line);
  border-radius:8px;margin-bottom:6px;font-size:12px}
.favitem .fpath{flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:var(--txt)}
.favitem button{border:none;background:var(--side);border-radius:6px;padding:3px 8px;cursor:pointer;
  font-size:11.5px;color:var(--txt);font-family:inherit}
.favitem button:hover{background:var(--accent);color:#fff}
.favitem .fdel:hover{background:#e8473f;color:#fff}

/* ---------- 워크플로우 보드 ---------- */
.wfbox{width:920px;max-width:95vw;height:84vh;display:flex;flex-direction:column;overflow:hidden;padding:18px}
.wfhead{display:flex;align-items:center;gap:10px;margin-bottom:8px}
.wfhead h3{margin:0}
.wftoggle{display:flex;background:var(--side);border-radius:9px;padding:3px}
.wftab{border:none;background:transparent;padding:6px 14px;border-radius:7px;cursor:pointer;
  font-family:inherit;font-size:12.5px;color:var(--mut)}
.wftab.on{background:var(--accent);color:#fff}
.wfx{border:none;background:var(--side);width:30px;height:30px;border-radius:8px;cursor:pointer;
  font-size:14px;color:var(--mut)}
.wfx:hover{background:#e8473f;color:#fff}
.wflegend{display:flex;align-items:center;gap:16px;font-size:11.5px;color:var(--mut);
  padding:6px 2px 12px;border-bottom:1px solid var(--line);margin-bottom:12px}
.wflegend .dot{display:inline-block;width:9px;height:9px;border-radius:50%;margin-right:5px;vertical-align:middle}
.p-high{background:#e8473f}.p-mid{background:#f0993e}.p-low{background:#9aa0aa}
.wfadd{border:1px solid var(--accent);color:var(--accent);background:transparent;border-radius:8px;
  padding:5px 12px;cursor:pointer;font-family:inherit;font-size:12px}
.wfadd:hover{background:var(--accent);color:#fff}
.wfboard{flex:1;overflow:auto;display:flex;gap:14px;align-items:flex-start;padding-bottom:6px}
.wfcol{min-width:260px;max-width:280px;flex:0 0 auto;background:var(--side);border-radius:12px;
  padding:12px;border-top:4px solid var(--accent)}
.wfcoltop{display:flex;align-items:center;gap:6px;margin-bottom:4px}
.wfcoltop b{flex:1;font-size:13.5px;color:var(--txt);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.wfcoltop .wficon{cursor:pointer;font-size:12px;color:var(--mut);opacity:.6}
.wfcoltop .wficon:hover{opacity:1}
.wfbar{height:6px;background:var(--line);border-radius:4px;overflow:hidden;margin:6px 0 4px}
.wfbar>i{display:block;height:100%;background:var(--accent);border-radius:4px;transition:.3s}
.wfprog{font-size:11px;color:var(--mut);margin-bottom:10px}
.wftask{background:var(--bg);border:1px solid var(--line);border-left:4px solid #9aa0aa;border-radius:8px;
  padding:8px 10px;margin-bottom:7px;font-size:12.5px;cursor:pointer;position:relative}
.wftask.pri-높음{border-left-color:#e8473f}
.wftask.pri-보통{border-left-color:#f0993e}
.wftask.pri-낮음{border-left-color:#9aa0aa}
.wftask.done{opacity:.5}
.wftask.done .wtt{text-decoration:line-through}
.wtt{color:var(--txt);line-height:1.35;padding-right:16px}
.wtmeta{font-size:10.5px;color:var(--mut);margin-top:4px;display:flex;gap:8px;flex-wrap:wrap}
.wtx{position:absolute;top:6px;right:7px;font-size:11px;color:var(--mut);opacity:0;cursor:pointer}
.wftask:hover .wtx{opacity:.7}.wtx:hover{color:#e8473f;opacity:1}
.wfaddtask{width:100%;border:1px dashed var(--line2);background:transparent;border-radius:8px;
  padding:7px;cursor:pointer;font-family:inherit;font-size:12px;color:var(--mut)}
.wfaddtask:hover{border-color:var(--accent);color:var(--accent)}
.wfempty{color:var(--mut);font-size:13px;text-align:center;padding:40px 20px;width:100%}
.wfpgroup{min-width:300px;max-width:340px}

/* ---------- 워크플로우 상시 사이드 패널(우측 레일) ---------- */
#wfrail{width:236px;min-width:236px;background:var(--side);border-left:1px solid var(--line);
  display:flex;flex-direction:column;transition:margin .2s}
#wfrail.hidden{margin-right:-238px}
.wfrhead{display:flex;align-items:center;gap:6px;padding:12px 12px 8px;border-bottom:1px solid var(--line)}
.wfrhead b{flex:1;font-size:13px;color:var(--txt)}
.wfrhead button{border:none;background:transparent;cursor:pointer;color:var(--mut);font-size:13px;
  width:24px;height:24px;border-radius:6px}
.wfrhead button:hover{background:var(--bg);color:var(--txt)}
#wfrbody{flex:1;overflow-y:auto;padding:10px}
.wfrproj{margin-bottom:14px}
.wfrproj .wfrname{font-size:12.5px;font-weight:600;color:var(--txt);display:flex;align-items:center;gap:5px;margin-bottom:5px}
.wfrproj .wfrname i{width:8px;height:8px;border-radius:50%;flex-shrink:0}
.wfrbar{height:5px;background:var(--line);border-radius:3px;overflow:hidden;margin-bottom:3px}
.wfrbar>span{display:block;height:100%;border-radius:3px;transition:.3s}
.wfrpct{font-size:10px;color:var(--mut);margin-bottom:6px}
.wfrtask{font-size:11.5px;color:var(--txt);padding:4px 6px;border-radius:6px;cursor:pointer;
  border-left:3px solid #9aa0aa;background:var(--bg);margin-bottom:4px;line-height:1.3}
.wfrtask.pri-높음{border-left-color:#e8473f}
.wfrtask.pri-보통{border-left-color:#f0993e}
.wfrtask.pri-낮음{border-left-color:#9aa0aa}
.wfrtask.doing{font-weight:600}
.wfrtask:hover{background:var(--line)}
.wfrempty{font-size:11.5px;color:var(--mut);text-align:center;padding:24px 8px;line-height:1.5}
#wfrtab{position:fixed;right:0;top:90px;z-index:60;background:var(--accent);color:#fff;border:none;
  border-radius:9px 0 0 9px;padding:9px 7px;cursor:pointer;font-size:14px;display:none;
  box-shadow:-2px 2px 8px rgba(0,0,0,.15);writing-mode:vertical-rl}
</style></head>
<body>
<aside id="side">
  <div class="side-top">
    <div class="brand"><span class="star">✱</span> KJH비서</div>
    <button id="newchat">✚ &nbsp;새 대화</button>
    <button id="calbtn">📅 &nbsp;일정</button>
    <button id="memobtn">📝 &nbsp;메모</button>
    <button id="wfbtn">📋 &nbsp;워크플로우</button>
    <button id="cleanbtn">🧹 &nbsp;폴더 정리</button>
    <button id="shopbtn">🛒 &nbsp;쇼핑·주문</button>
  </div>
  <div class="side-label">최근 대화</div>
  <div id="histlist"></div>
  <div class="side-bottom">
    <button class="fbtn" id="fminus" title="글씨 작게">A−</button>
    <button class="fbtn" id="fplus" title="글씨 크게">A＋</button>
    <button class="fbtn" id="darkbtn" title="다크/라이트 모드">🌙</button>
    <button class="fbtn" id="exportbtn" title="현재 대화를 파일로 저장">⬇</button>
    <button class="fbtn" id="setbtn" title="설정">⚙</button>
    <div id="settingsmenu">
      <button class="fbtn" id="morebtn" title="기능·업데이트·이름설정">☰</button>
      <div id="settingspop">
        <button class="spitem" id="sp_feat">✨ 할 수 있는 일 목록</button>
        <button class="spitem" id="sp_update">🔄 업데이트 확인</button>
        <button class="spitem" id="sp_name">⚙ 비서 이름 설정</button>
      </div>
    </div>
    <span id="status">부팅 중…</span>
  </div>
</aside>
<main>
  <div id="updbanner">
    <span class="updico">🔄</span>
    <div class="updtxt">
      <b id="updver"></b> 업데이트가 있습니다
      <span class="updlog" id="updlog"></span>
    </div>
    <button class="updbtn do" onclick="doUpdate(this)">지금 업데이트</button>
    <button class="updbtn skip" onclick="hideUpdate()">나중에</button>
  </div>
  <button id="sidetoggle" title="사이드바 접기/펴기">☰</button>
  <div id="chatwrap">
    <div id="greet">
      <div class="gstar">✱</div>
      <h2>안녕하세요, 대표님</h2>
      <p>이 컴퓨터의 모든 파일을 알고 있습니다 — 경로 몰라도 "○○ 찾아줘"라고만 하세요.<br>
      파일 찾기 · 내용 분석 · 문서 작업 · 도면 분석 · 프로그램 실행 무엇이든.<br>
      "내일 3시 미팅 잡아줘, 30분 전에 알려줘" — 일정·메모는 왼쪽 📅 달력에서 한눈에.<br>
      <b>왼쪽 ✨ '할 수 있는 일'</b>을 누르면 이 비서가 가능한 모든 기능을 예시와 함께 볼 수 있습니다.</p>
    </div>
    <div id="chat"></div>
  </div>
  <footer><div class="fwrap">
    <div id="qbar"></div>
    <div id="attach-bar"></div>
    <div id="inputbox">
      <textarea id="inp" rows="1" placeholder="무엇이든 물어보세요"></textarea>
      <div class="inrow">
        <button class="ibtn" id="attach" title="파일 첨부">＋</button>
        <button class="ibtn" id="clip" title="복사한 내용 불러와 질문">📋</button>
        <button class="ibtn" id="capture" title="화면 캡처 후 자동 첨부">📷</button>
        <button class="ibtn" id="mic" title="음성으로 말하기">🎤</button>
        <button class="ibtn" id="tts" title="답변 음성으로 듣기 (켜기/끄기)">🔈</button>
        <button class="ibtn" id="loop" title="연속 음성 대화 (핸즈프리 — 답변 후 자동으로 다시 듣기)">🔁</button>
        <span class="flex"></span>
        <button id="send" title="전송 (Enter)">↑</button>
      </div>
    </div>
    <div class="disclaim">로컬 AI · 무료 · 컴퓨터 전체 파일 이해 · 모든 작업 자동 실행(프리패스)</div>
  </div></footer>
</main>
<aside id="wfrail" class="hidden">
  <div class="wfrhead">
    <b>📋 워크플로우</b>
    <button id="wfropen" title="크게 보기/편집">⤢</button>
    <button id="wfrhide" title="패널 접기">✕</button>
  </div>
  <div id="wfrbody"></div>
</aside>
<button id="wfrtab" title="워크플로우 패널 열기">📋</button>

<div id="cmodal" class="overlay"><div class="box">
  <h3>확인</h3><p id="ctext"></p>
  <div class="brow"><button class="bcancel" id="ccancel">취소</button><button class="bok" id="cok">확인</button></div>
</div></div>
<div id="imodal" class="overlay"><div class="box">
  <h3 id="ititle">입력</h3>
  <input id="iinp" style="width:100%;border:1px solid var(--line2);border-radius:8px;padding:9px 11px;font-size:14px;font-family:inherit;box-sizing:border-box;margin-bottom:6px">
  <div class="brow"><button class="bcancel" id="icancel">취소</button><button class="bok" id="iok">확인</button></div>
</div></div>
<div id="setmodal" class="overlay"><div class="box" style="width:420px;max-width:92vw">
  <h3>⚙ 설정</h3>
  <div class="setrow"><span>🌙 다크 모드</span>
    <label class="sw"><input type="checkbox" id="set_dark"><span class="slider"></span></label></div>
  <div class="setrow"><span>✨ 피카츄 캐릭터</span>
    <label class="sw"><input type="checkbox" id="set_pika"><span class="slider"></span></label></div>
  <div class="setrow"><span>💬 답변 스타일</span>
    <select id="set_style"><option value="간결">간결</option><option value="표준">표준</option><option value="자세히">자세히</option></select></div>
  <div class="setrow"><span>🔊 음성 읽기 속도 <b id="set_rateval">1.05x</b></span>
    <input type="range" id="set_rate" min="0.6" max="1.6" step="0.05" style="width:150px"></div>
  <div class="setsec">📁 즐겨찾기 폴더 <span style="color:var(--mut);font-weight:400">— 자주 쓰는 작업 폴더</span></div>
  <div id="favlist" style="max-height:160px;overflow-y:auto;margin-bottom:8px"></div>
  <button id="set_addfav" class="bcancel" style="width:100%;padding:8px">＋ 폴더 추가</button>
  <div class="brow"><button class="bok" id="setclose">닫기</button></div>
</div></div>
<div id="wfmodal" class="overlay"><div class="box wfbox">
  <div class="wfhead">
    <h3 style="flex:1">📋 워크플로우 <span style="font-size:12px;color:var(--mut);font-weight:400">— 프로젝트·업무 진척도</span></h3>
    <div class="wftoggle">
      <button class="wftab on" id="wf_byproj">프로젝트별</button>
      <button class="wftab" id="wf_bypri">우선순위별</button>
    </div>
    <button class="wfx" id="wfclose">✕</button>
  </div>
  <div class="wflegend">
    <span><i class="dot p-high"></i>높음</span>
    <span><i class="dot p-mid"></i>보통</span>
    <span><i class="dot p-low"></i>낮음</span>
    <span style="flex:1"></span>
    <button id="wf_addproj" class="wfadd">＋ 프로젝트</button>
  </div>
  <div id="wfboard" class="wfboard"></div>
</div></div>
<div id="qmodal" class="overlay"><div class="box">
  <h3>⚡ 빠른 버튼 만들기</h3>
  <input id="qname" placeholder="버튼 이름 (예: 스케치업)">
  <input id="qcmd" placeholder="명령 또는 URL (예: 스케치업 열어줘 / https://...)">
  <div class="brow"><button class="bcancel" id="qcancel">취소</button><button class="bok" id="qok">만들기</button></div>
</div></div>
<div id="namemodal" class="overlay"><div class="box">
  <h3>⚙ 비서 이름 설정</h3>
  <p style="font-size:12.5px;color:var(--mut);margin-bottom:10px">AI 비서를 부를 이름을 정해주세요</p>
  <input id="nameInp" placeholder="예: 민수비서, 클로드, AI비서" style="width:100%;border:1px solid var(--line2);border-radius:8px;padding:9px 11px;font-size:14px;font-family:inherit;box-sizing:border-box;margin-bottom:12px">
  <div class="brow"><button class="bcancel" id="namecancel">취소</button><button class="bok" id="nameok">저장</button></div>
</div></div>
<div id="editmodal" class="overlay"><div class="box">
  <h3>✏️ 일정 수정</h3>
  <input type="hidden" id="editId">
  <div style="display:flex;flex-direction:column;gap:8px;margin-bottom:12px">
    <div style="display:flex;gap:6px">
      <input type="date" id="editDate" style="flex:1;border:1px solid var(--line2);border-radius:8px;padding:7px 9px;font-size:13px;font-family:inherit">
      <input type="time" id="editTime" style="flex:1;border:1px solid var(--line2);border-radius:8px;padding:7px 9px;font-size:13px;font-family:inherit">
    </div>
    <input id="editTitle" placeholder="일정 제목" style="border:1px solid var(--line2);border-radius:8px;padding:9px 11px;font-size:13px;font-family:inherit">
    <select id="editCat" style="border:1px solid var(--line2);border-radius:8px;padding:7px 9px;font-size:13px;font-family:inherit;background:var(--bg);color:var(--txt)">
      <option value="개인">👤 개인</option>
      <option value="회사">🏢 회사</option>
    </select>
    <select id="editRemind" style="border:1px solid var(--line2);border-radius:8px;padding:7px 9px;font-size:13px;font-family:inherit;background:var(--bg);color:var(--txt)">
      <option value="">알림 없음</option>
      <option value="0">정각</option>
      <option value="10">10분 전</option>
      <option value="30">30분 전</option>
      <option value="60">1시간 전</option>
      <option value="1440">하루 전</option>
      <option value="1440,30">하루 전 + 30분 전</option>
    </select>
  </div>
  <div class="brow"><button class="bcancel" id="editcancel">취소</button><button class="bok" id="editsave">저장</button></div>
</div></div>
<div id="calmodal" class="overlay"><div class="box calbox">
  <div class="calhead">
    <span class="chttl">📅 일정</span>
    <span class="flex"></span>
    <button class="cnav" id="calprev">‹</button>
    <span class="cym" id="calym"></span>
    <button class="cnav" id="calnext">›</button>
    <button class="cnav" id="caltoday" title="오늘로">오늘</button>
    <button class="cnav" id="calclose" title="닫기">✕</button>
  </div>
  <div id="calpane">
    <div id="calgrid"></div>
    <div class="cfilter" id="cfilter"></div>
    <div id="daylist"></div>
    <div class="addrow">
      <input id="evtitle" placeholder="일정 내용 — Enter로 추가">
      <input type="time" id="evtime" title="시간 (비우면 종일 일정)">
      <select id="evcat" title="일정 구분">
        <option value="개인">👤 개인</option>
        <option value="회사">🏢 회사</option>
      </select>
      <select id="evremind" title="알림 시점">
        <option value="">알림 없음</option>
        <option value="0">정각</option>
        <option value="10">10분 전</option>
        <option value="30" selected>30분 전</option>
        <option value="60">1시간 전</option>
        <option value="1440">하루 전</option>
        <option value="1440,30">하루 전 + 30분 전</option>
      </select>
      <button class="badd" id="evadd">추가</button>
    </div>
  </div>
</div></div>
<div id="memomodal" class="overlay"><div class="box memobox">
  <div class="calhead">
    <span class="chttl">📝 메모</span>
    <span class="flex"></span>
    <button class="cnav" id="memoclose" title="닫기">✕</button>
  </div>
  <div id="memolist"></div>
  <div class="addrow">
    <input id="memotext" placeholder="메모 입력 — Enter로 저장">
    <button class="badd" id="memoadd">저장</button>
  </div>
</div></div>
<div id="modal" class="overlay"><div class="box">
  <h3>🔐 금고 마스터 비밀번호</h3>
  <input type="password" id="mpw" placeholder="마스터 비밀번호">
  <div class="brow"><button class="bcancel" id="mcancel">취소</button><button class="bok" id="mok">확인</button></div>
</div></div>
<div id="featmodal" class="overlay"><div class="box featbox">
  <div class="feathead">
    <h3>✨ 이 비서가 할 수 있는 일</h3>
    <button class="cnav" id="featclose" title="닫기">✕</button>
  </div>
  <div class="featsub">아래 예시를 누르면 입력창에 채워집니다 — Enter로 실행하세요. 전부 무료·내 컴퓨터에서 동작하며, 인터넷으로 대화가 새어 나가지 않습니다.</div>
  <div class="featgrid" id="featgrid"></div>
  <div class="featnote">🔒 위험한 명령(디스크 포맷·시스템 폴더 삭제 등)은 자동 차단됩니다 · 원본 파일은 덮어쓰지 않고 결과물 폴더에 새로 저장합니다</div>
</div></div>
<div id="cleanmodal" class="overlay"><div class="box dashbox">
  <div class="calhead">
    <span class="chttl">🧹 폴더 정리</span>
    <span class="flex"></span>
    <button class="cnav" id="cleanclose" title="닫기">✕</button>
  </div>
  <div class="dashsub">폴더 안의 파일을 종류별 하위폴더(이미지·문서·엑셀·도면·압축 등)로 자동 분류합니다. <b>원본은 이동만 하고 삭제하지 않습니다.</b></div>
  <div class="dashsec">자주 쓰는 폴더</div>
  <div id="cleanquick" class="cleangrid"></div>
  <button id="cleanpick" class="dashbtn">📁 다른 폴더 선택해서 정리…</button>
  <div id="cleanresult" class="dashresult"></div>
</div></div>
<div id="shopmodal" class="overlay"><div class="box dashbox">
  <div class="calhead">
    <span class="chttl">🛒 쇼핑·주문</span>
    <span class="flex"></span>
    <button class="cnav" id="shopclose" title="닫기">✕</button>
  </div>
  <div class="dashsub">사고 싶은 물건을 적으면 비서가 가격·평점을 조사해 후보를 표로 정리하고 추천합니다. <b>결제 확정은 항상 직접 확인 후 진행</b>합니다.</div>
  <div class="dashsec">무엇을 찾을까요?</div>
  <input id="shopq" class="dashinput" placeholder="예: 인체공학 사무용 의자 10만원대">
  <div class="dashsec">쇼핑몰</div>
  <div class="shopsites">
    <label class="shopchk"><input type="checkbox" value="쿠팡" checked> 쿠팡</label>
    <label class="shopchk"><input type="checkbox" value="네이버쇼핑" checked> 네이버쇼핑</label>
    <label class="shopchk"><input type="checkbox" value="11번가"> 11번가</label>
    <label class="shopchk"><input type="checkbox" value="G마켓"> G마켓</label>
  </div>
  <div class="shopex">
    <span class="shoptag" onclick="setShop('A4용지 한 박스 최저가')">A4용지 박스</span>
    <span class="shoptag" onclick="setShop('무선 마우스 3만원 이하 가성비')">무선 마우스</span>
    <span class="shoptag" onclick="setShop('사무실 공기청정기 30평형')">공기청정기</span>
  </div>
  <button id="shopgo" class="dashbtn primary">🔎 조사 시작</button>
</div></div>

<script>
const chat=document.getElementById('chat'),inp=document.getElementById('inp'),
sendBtn=document.getElementById('send'),attachBtn=document.getElementById('attach'),
attachBar=document.getElementById('attach-bar'),statusEl=document.getElementById('status'),
modal=document.getElementById('modal'),mpw=document.getElementById('mpw'),
cmodal=document.getElementById('cmodal'),greet=document.getElementById('greet'),
chatwrap=document.getElementById('chatwrap');
let cur=null,masterPw=null,thinkEl=null,fontSize=15,curFile='';

function showGreet(b){greet.style.display=b?'block':'none'}
function hasMsgs(){return chat.children.length>0}
function atBottom(){return chatwrap.scrollHeight-chatwrap.scrollTop-chatwrap.clientHeight<90}
function scrollEnd(force){if(force||atBottom())chatwrap.scrollTop=chatwrap.scrollHeight}

function showThinking(){hideThinking();thinkEl=document.createElement('div');
  thinkEl.className='thinking';thinkEl.innerHTML='<span class="tstar">✱</span> 생각하는 중…';
  chat.appendChild(thinkEl);scrollEnd()}
function hideThinking(){if(thinkEl){thinkEl.remove();thinkEl=null}}
function applyFS(n,save){fontSize=Math.max(12,Math.min(22,n));
  document.documentElement.style.setProperty('--fs',fontSize+'px');
  if(save)pywebview.api.save_font(fontSize)}
function setStatus(t,cls){statusEl.textContent=t;statusEl.className=cls||''}

/* ---------- 마크다운 렌더링 ---------- */
function esc(s){return s.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')}
function inlineMd(s){
  s=s.replace(/`([^`]+)`/g,function(_,c){return '<code>'+c+'</code>'});
  s=s.replace(/\*\*([^*]+)\*\*/g,'<b>$1</b>');
  s=s.replace(/(^|[\s(])\*([^*\n]+)\*/g,'$1<i>$2</i>');
  s=s.replace(/(https?:\/\/[^\s<)]+)/g,'<a href="$1" target="_blank">$1</a>');
  return s;
}
function blockMd(s){
  const lines=s.split('\n');let out='',list=null;
  const close=function(){if(list){out+='</'+list+'>';list=null}};
  for(let i=0;i<lines.length;i++){
    const t=lines[i];
    let m;
    if(/^####?#?\s/.test(t)){close();
      const lv=t.match(/^#+/)[0].length;
      const tag=lv<=1?'h2':(lv===2?'h3':'h4');
      out+='<'+tag+'>'+inlineMd(t.replace(/^#+\s*/,''))+'</'+tag+'>';continue}
    if(/^\s*[-*_]{3,}\s*$/.test(t)){close();out+='<hr>';continue}
    m=t.match(/^\s*[-•]\s+(.*)/);
    if(m){if(list!=='ul'){close();out+='<ul>';list='ul'}out+='<li>'+inlineMd(m[1])+'</li>';continue}
    m=t.match(/^\s*\d+[.)]\s+(.*)/);
    if(m){if(list!=='ol'){close();out+='<ol>';list='ol'}out+='<li>'+inlineMd(m[1])+'</li>';continue}
    m=t.match(/^&gt;\s?(.*)/);
    if(m){close();out+='<blockquote>'+inlineMd(m[1])+'</blockquote>';continue}
    close();
    if(t.trim()==='')out+='<div class="sp"></div>';
    else out+='<p>'+inlineMd(t)+'</p>';
  }
  close();return out;
}
function md(s){
  s=s.replace(/```(?:memory|search|open|auto|read|web|plan)\n[\s\S]*?(?:```|$)/g,'').trim();
  let out='',last=0,re=/```(\w*)\n?([\s\S]*?)(?:```|$)/g,m;
  const src=esc(s);
  while((m=re.exec(src))!==null){
    out+=blockMd(src.slice(last,m.index));
    out+='<div class="code"><div class="codehead"><span>'+(m[1]||'code')
      +'</span><button onclick="copyCode(this)">복사</button></div><pre>'+m[2]+'</pre></div>';
    last=re.lastIndex;
  }
  out+=blockMd(src.slice(last));
  return out;
}
function copyCode(btn){
  const pre=btn.closest('.code').querySelector('pre');
  navigator.clipboard.writeText(pre.textContent).then(function(){
    btn.textContent='✓ 복사됨';setTimeout(function(){btn.textContent='복사'},1400)});
}

/* ---------- 메시지 ---------- */
function userMsg(t,files,filePaths){
  showGreet(false);
  const r=document.createElement('div');r.className='row user';
  const wrap=document.createElement('div');wrap.className='user-msg-wrap';
  const b=document.createElement('div');b.className='bubble';
  b.innerHTML=esc(t)+(files&&files.length?'<br><span class="filechip">📎 '+esc(files.join(', '))+'</span>':'');
  wrap.appendChild(b);
  // 액션 버튼
  const acts=document.createElement('div');acts.className='msg-actions';
  if(t){
    const cpBtn=document.createElement('button');cpBtn.className='msg-action-btn';cpBtn.textContent='복사';
    cpBtn.title='텍스트 클립보드에 복사';
    cpBtn.onclick=function(){navigator.clipboard.writeText(t).then(function(){
      cpBtn.textContent='✓ 복사됨';setTimeout(function(){cpBtn.textContent='복사'},1400)})};
    acts.appendChild(cpBtn);
    const reBtn=document.createElement('button');reBtn.className='msg-action-btn';reBtn.textContent='다시입력';
    reBtn.title='이 메시지를 입력창에 다시 넣기';
    reBtn.onclick=function(){inp.value=t;inp.style.height='auto';
      inp.style.height=Math.min(inp.scrollHeight,170)+'px';inp.focus()};
    acts.appendChild(reBtn);
  }
  if(filePaths&&filePaths.length){
    const raBtn=document.createElement('button');raBtn.className='msg-action-btn';raBtn.textContent='파일 재첨부';
    raBtn.title='이 파일들을 다시 첨부';
    raBtn.onclick=function(){pywebview.api.reattach_files(filePaths).then(showAttached)};
    acts.appendChild(raBtn);
  }
  wrap.appendChild(acts);
  r.appendChild(wrap);chat.appendChild(r);scrollEnd(true);
}
function aiRow(){
  showGreet(false);
  const r=document.createElement('div');r.className='row ai';
  r.innerHTML='<div class="av">✱</div>';
  const b=document.createElement('div');b.className='bubble md';
  r.appendChild(b);chat.appendChild(r);scrollEnd();
  return {row:r,el:b};
}
function aiMsg(t){const a=aiRow();a.el.innerHTML=md(t);scrollEnd()}
function sysMsg(t){showGreet(false);const d=document.createElement('div');d.className='sys';
  d.textContent=t;chat.appendChild(d);scrollEnd()}
function esc(s){return (s||'').replace(/[&<>]/g,function(c){
  return {'&':'&amp;','<':'&lt;','>':'&gt;'}[c]})}
function renderHistSearch(q,res){
  showGreet(false);
  const d=document.createElement('div');d.className='sys';
  if(!res||!res.length){
    d.textContent='🔎 "'+q+'" 와(과) 비슷한 과거 대화를 찾지 못했습니다.';
    chat.appendChild(d);scrollEnd();return;
  }
  const h=document.createElement('div');
  h.textContent='🔎 "'+q+'" 관련 과거 대화 '+res.length+'건 — 클릭하면 그 대화를 엽니다';
  d.appendChild(h);
  res.forEach(function(r){
    const item=document.createElement('div');
    item.style.cssText='margin-top:8px;padding:8px 10px;border:1px solid var(--line2);'
      +'border-radius:8px;cursor:pointer';
    item.innerHTML='<b>'+esc(r.title)+'</b> <span style="color:var(--mut);font-size:11px">'
      +r.date+' · 유사도 '+r.score+'</span><br>'
      +'<span style="color:var(--mut);font-size:12px">'+esc(r.preview)+'</span>';
    item.onclick=function(){clearChat();pywebview.api.load_history(r.file).then(loadHistList)};
    d.appendChild(item);
  });
  chat.appendChild(d);scrollEnd();
}
function clearChat(){chat.innerHTML='';showGreet(true)}

/* ---------- 스트리밍 ---------- */
let tq='',tTimer=null,finishing=false;
function aiStart(){hideThinking();cur=Object.assign({raw:''},aiRow())}
function drain(){
  if(!tq){
    clearInterval(tTimer);tTimer=null;
    if(finishing){
      if(cur&&!cur.el.innerHTML)cur.row.remove();
      else if(cur&&cur.raw)speak(cur.raw);   // 답변 완료 → 음성으로 읽기(켜진 경우)
      cur=null;finishing=false}
    return;
  }
  cur.raw+=tq;tq='';                  // 도착한 글자를 즉시 전부 표시 (최대 속도)
  cur.el.innerHTML=md(cur.raw);scrollEnd();
}

/* ---------- 음성 출력(TTS, 브라우저 내장 — 추가 설치 없음) / 입력(STT) ---------- */
let ttsOn=false, micBusy=false, continuousVoice=false;
function maybeRelisten(){   // 연속 음성 대화: 답변(음성)이 끝나면 자동으로 다시 듣기
  if(continuousVoice&&!busy&&!micBusy)setTimeout(startListen,350);
}
function speak(raw){
  if(!ttsOn||!('speechSynthesis'in window)){maybeRelisten();return}
  // 코드블록·도구블록·마크다운 기호 제거 후 읽기
  let s=raw.replace(/```[\s\S]*?```/g,' ').replace(/[*#`>_~\-]/g,' ')
          .replace(/https?:\/\/\S+/g,' ').replace(/\s+/g,' ').trim();
  if(!s){maybeRelisten();return}
  try{speechSynthesis.cancel();
    const u=new SpeechSynthesisUtterance(s.slice(0,600));
    u.lang='ko-KR';u.rate=window._ttsRate||1.05;
    const ko=speechSynthesis.getVoices().find(v=>/ko/i.test(v.lang));
    if(ko)u.voice=ko;
    u.onend=maybeRelisten;              // 다 읽고 나면 자동으로 다시 마이크 켜기
    speechSynthesis.speak(u);
  }catch(e){maybeRelisten();}
}
async function startListen(){
  if(micBusy)return;
  const micBtn=document.getElementById('mic');
  micBusy=true;if(micBtn)micBtn.style.color='var(--accent)';
  try{
    const text=await pywebview.api.start_voice();
    if(text){inp.value=text;inp.style.height='auto';
      inp.style.height=Math.min(inp.scrollHeight,170)+'px';doSend();}
  }catch(e){}
  micBusy=false;if(micBtn)micBtn.style.color='';
}
function setupVoiceButtons(){
  const micBtn=document.getElementById('mic'),ttsBtn=document.getElementById('tts'),
        loopBtn=document.getElementById('loop');
  ttsBtn.onclick=function(){
    ttsOn=!ttsOn;ttsBtn.textContent=ttsOn?'🔊':'🔈';
    ttsBtn.style.color=ttsOn?'var(--accent)':'';
    if(!ttsOn&&'speechSynthesis'in window)speechSynthesis.cancel();
    sysMsg(ttsOn?'🔊 음성 답변을 켰습니다':'🔈 음성 답변을 껐습니다');
  };
  micBtn.onclick=startListen;
  loopBtn.onclick=function(){
    continuousVoice=!continuousVoice;
    loopBtn.style.color=continuousVoice?'var(--accent)':'';
    if(continuousVoice&&!ttsOn)ttsBtn.onclick();   // 핸즈프리면 음성 답변도 자동 켜기
    sysMsg(continuousVoice
      ?'🔁 연속 음성 대화를 켰습니다 — 답변이 끝나면 자동으로 다시 듣습니다 (끄려면 🔁 다시 클릭)'
      :'🔁 연속 음성 대화를 껐습니다');
    if(continuousVoice)startListen();
  };
}
function aiToken(t){
  if(!cur)aiStart();
  tq+=t;
  drain();                            // 받는 즉시 그린다 (지연 없음)
  if(!tTimer)tTimer=setInterval(drain,16);
}
function aiDone(){
  hideThinking();finishing=true;showGenMeta();
  if(!tq&&!tTimer){if(cur&&!cur.el.innerHTML)cur.row.remove();cur=null;finishing=false}
}

/* ---------- 확인/입력 모달 ---------- */
function askConfirm(msg){return new Promise(function(res){
  cmodal.style.display='flex';document.getElementById('ctext').textContent=msg;
  document.getElementById('cok').onclick=function(){cmodal.style.display='none';res(true)};
  document.getElementById('ccancel').onclick=function(){cmodal.style.display='none';res(false)};
})}
function askInput(title,placeholder){return new Promise(function(res){
  var m=document.getElementById('imodal'),inpE=document.getElementById('iinp');
  document.getElementById('ititle').textContent=title||'입력';
  inpE.value='';inpE.placeholder=placeholder||'';m.style.display='flex';
  setTimeout(function(){inpE.focus()},50);
  function done(v){m.style.display='none';inpE.onkeydown=null;res(v)}
  document.getElementById('iok').onclick=function(){done(inpE.value.trim()||null)};
  document.getElementById('icancel').onclick=function(){done(null)};
  inpE.onkeydown=function(e){if(e.key==='Enter')done(inpE.value.trim()||null);
    else if(e.key==='Escape')done(null)};
})}
function askPass(){return new Promise(function(res){
  modal.style.display='flex';mpw.value='';mpw.focus();
  document.getElementById('mok').onclick=function(){modal.style.display='none';res(mpw.value)};
  document.getElementById('mcancel').onclick=function(){modal.style.display='none';res(null)};
  mpw.onkeydown=function(e){if(e.key==='Enter'){modal.style.display='none';res(mpw.value)}};
})}

/* ---------- 첨부 ---------- */
let attached=[];
function showAttached(names){
  attached=names||[];
  attachBar.innerHTML='';
  if(!attached.length){attachBar.style.display='none';return}
  attachBar.style.display='block';
  attached.forEach(function(n,i){
    const c=document.createElement('span');c.className='filechip';c.dataset.idx=i;
    c.appendChild(document.createTextNode('📎 '+n));
    const x=document.createElement('span');x.textContent=' ✕';x.title='이 첨부 빼기';
    x.style.cssText='cursor:pointer;color:#C0392B;font-weight:bold';
    x.onclick=function(){pywebview.api.remove_attached(i).then(showAttached)};
    c.appendChild(x);attachBar.appendChild(c);
  });
  /* 이미지·DXF 썸네일을 비동기로 채워 칩 앞에 붙인다 */
  pywebview.api.attach_thumbs().then(function(urls){
    if(!urls)return;
    urls.forEach(function(u,i){
      if(!u)return;
      const c=attachBar.querySelector('.filechip[data-idx="'+i+'"]');
      if(!c||c.querySelector('img.thumb'))return;
      const im=document.createElement('img');im.className='thumb';im.src=u;
      im.title='미리보기';c.insertBefore(im,c.firstChild);
    });
  }).catch(function(){});
}
// 답변 중에도 다음 채팅 입력 가능(큐로 순서 처리). 버튼은 생성 중 ⏹(중단)으로 바뀐다.
let busy=false;
function setBusy(b){
  busy=b;
  if(sendBtn){sendBtn.textContent=b?'⏹':'↑';
    sendBtn.title=b?'생성 중단':'전송 (Enter)';
    sendBtn.classList.toggle('stopping',b);}
  if(b){if(window._pikaWork)window._pikaWork();}
  else{loadHistList();if(window._pikaDone)window._pikaDone();}
}

/* ---------- 금고 ---------- */
async function handleVault(t){
  const parts=t.split(/\s+/);
  const action=parts[1]||'도움말';
  let shown=t;
  if(action==='저장'&&parts.length>=5){shown=parts.slice(0,4).join(' ')+' ●●●●'}
  userMsg(shown);
  if(!['저장','보기','목록','삭제'].includes(action)){
    sysMsg('금고 사용법: /금고 저장 이름 아이디 비밀번호 · /금고 보기 이름 · /금고 목록 · /금고 삭제 이름');return;
  }
  if(masterPw===null){masterPw=await askPass();if(!masterPw){sysMsg('취소되었습니다');return}}
  const r=await pywebview.api.vault(t,masterPw);
  if(r==='WRONGPW'){masterPw=null;sysMsg('마스터 비밀번호가 틀렸습니다. 다시 시도하세요.')}
  else sysMsg(r);
}

/* ---------- 전송 ---------- */
async function doSend(){
  const t=inp.value.trim();
  const hasFiles=attached.length>0;
  if(!t&&!hasFiles)return;
  inp.value='';inp.style.height='auto';
  if(t.startsWith('/금고')){await handleVault(t);return}
  if(t==='/일정'||t==='/달력'){openCal();return}
  if(t==='/메모'){openMemo();return}
  if(t.startsWith('/메모 ')){userMsg(t);
    schedData=await pywebview.api.memo_add(t.slice(4).trim());
    sysMsg('📝 메모 저장됨 — 왼쪽 📝 메모 버튼에서 확인');return}
  if(t==='/말투'){sysMsg('말투 사용법: /말투 간결 · /말투 표준 · /말투 자세히');return}
  if(t.startsWith('/말투 ')){userMsg(t);
    pywebview.api.set_style(t.slice(4).trim()).then(sysMsg);return}
  if(t==='/찾기'){sysMsg('검색어를 적어 주세요. 예: /찾기 견적서 단가');return}
  if(t.startsWith('/찾기 ')){userMsg(t);
    pywebview.api.search_history(t.slice(4).trim());return}
  if(t==='/일정요약'||t==='/일정정리'){userMsg(t);
    const rs=await pywebview.api.sched_summary(null);aiMsg(rs);return}
  if(t==='/이름'){openNameModal();return}
  if(t.startsWith('/일정 ')){
    const cat2=t.slice(4).trim();
    if(cat2==='개인'||cat2==='회사'||cat2==='전체'){
      userMsg(t);
      const rs2=await pywebview.api.sched_summary(cat2);aiMsg(rs2);return}}
  let filePaths=[];
  if(hasFiles){try{filePaths=await pywebview.api.get_attached_paths();}catch(e){}}
  userMsg(t,hasFiles?attached:null,filePaths);
  showAttached([]);
  setBusy(true);
  showThinking();
  _genStart=performance.now();
  pywebview.api.send(t);
}
sendBtn.onclick=function(){
  if(busy){pywebview.api.stop();sysMsg('⏹ 생성을 중단했습니다');}
  else doSend();
};
setupVoiceButtons();
if('speechSynthesis'in window)speechSynthesis.getVoices();  // 음성 목록 미리 로드
inp.addEventListener('keydown',function(e){
  if(e.key==='Enter'&&!e.shiftKey&&!e.isComposing){e.preventDefault();doSend()}
});
inp.addEventListener('input',function(){inp.style.height='auto';inp.style.height=Math.min(inp.scrollHeight,170)+'px'});

/* ---------- 사이드바 ---------- */
document.getElementById('sidetoggle').onclick=function(){
  document.getElementById('side').classList.toggle('hidden')};
document.getElementById('fminus').onclick=function(){applyFS(fontSize-1,true)};
document.getElementById('fplus').onclick=function(){applyFS(fontSize+1,true)};
async function loadHistList(){
  const list=await pywebview.api.list_history();
  curFile=await pywebview.api.current_session();
  const hl=document.getElementById('histlist');hl.innerHTML='';
  if(!list.length){const d=document.createElement('div');d.className='hitem';
    d.innerHTML='<div class="ht" style="color:var(--mut)">저장된 대화가 없습니다</div>';hl.appendChild(d)}
  list.forEach(function(it){
    const d=document.createElement('div');d.className='hitem'+(it.file===curFile?' cur':'');
    const body=document.createElement('div');body.style.cssText='flex:1;min-width:0';
    body.innerHTML='<div class="ht">'+esc(it.title)+'</div><div class="hd">'+esc(it.date)+'</div>';
    body.onclick=function(){clearChat();pywebview.api.load_history(it.file).then(loadHistList)};
    function startRename(e){
      if(e)e.stopPropagation();
      const ht=body.querySelector('.ht');
      const inpEl=document.createElement('input');inpEl.className='htedit';inpEl.value=it.title;
      ht.replaceWith(inpEl);inpEl.focus();inpEl.select();
      let done=false;
      const save=async function(commit){
        if(done)return;done=true;
        const v=inpEl.value.trim();
        if(commit&&v&&v!==it.title){await pywebview.api.rename_history(it.file,v);it.title=v}
        loadHistList();
      };
      inpEl.onclick=function(ev){ev.stopPropagation()};
      inpEl.onkeydown=function(ev){ev.stopPropagation();
        if(ev.key==='Enter'){ev.preventDefault();save(true)}
        else if(ev.key==='Escape'){save(false)}};
      inpEl.onblur=function(){save(true)};
    }
    const ren=document.createElement('button');ren.className='hdel';
    ren.textContent='✎';ren.title='제목 수정';
    ren.onclick=startRename;
    const del=document.createElement('button');del.className='hdel';
    del.textContent='🗑';del.title='이 대화 삭제';
    del.onclick=async function(e){e.stopPropagation();
      if(!await askConfirm('이 대화를 삭제할까요?\n"'+it.title+'"'))return;
      const r=await pywebview.api.delete_history(it.file);
      if(r==='CURRENT'){clearChat();sysMsg('현재 대화가 삭제되어 새 대화를 시작합니다')}
      loadHistList();
    };
    d.appendChild(body);d.appendChild(ren);d.appendChild(del);hl.appendChild(d);
  });
}
document.getElementById('newchat').onclick=function(){
  pywebview.api.new_chat().then(function(){clearChat();loadHistList()})};

/* ---------- 빠른 버튼 ---------- */
let qbtns=[];
function setQButtons(list){qbtns=list||[];renderQ()}
function renderQ(){
  const bar=document.getElementById('qbar');bar.innerHTML='';
  qbtns.forEach(function(b,i){
    const d=document.createElement('span');d.className='qbtn';
    const isUrl=b.cmd&&(b.cmd.startsWith('http://')||b.cmd.startsWith('https://'));
    const t=document.createElement('span');t.textContent=(isUrl?'🔗':'⚡')+' '+b.name;
    t.onclick=function(){if(isUrl){pywebview.api.open_url(b.cmd);}else{inp.value=b.cmd;doSend();}};
    const x=document.createElement('span');x.className='qx';x.textContent='✕';x.title='버튼 삭제';
    x.onclick=async function(e){e.stopPropagation();
      if(!await askConfirm('빠른 버튼 "'+b.name+'"을 삭제할까요?'))return;
      qbtns.splice(i,1);pywebview.api.save_qbuttons(qbtns);renderQ()};
    d.appendChild(t);d.appendChild(x);bar.appendChild(d);
  });
  const add=document.createElement('button');add.className='qadd';add.textContent='＋ 버튼';
  add.title='자주 쓰는 명령을 원클릭 버튼으로';
  add.onclick=function(){
    const qm=document.getElementById('qmodal');qm.style.display='flex';
    const qn=document.getElementById('qname'),qc=document.getElementById('qcmd');
    qn.value='';qc.value='';qn.focus();
    document.getElementById('qok').onclick=function(){
      const n=qn.value.trim(),c=qc.value.trim();
      if(n&&c){qbtns.push({name:n,cmd:c});pywebview.api.save_qbuttons(qbtns);renderQ()}
      qm.style.display='none'};
    document.getElementById('qcancel').onclick=function(){qm.style.display='none'};
    qc.onkeydown=function(e){if(e.key==='Enter')document.getElementById('qok').onclick()};
  };
  bar.appendChild(add);
}
renderQ();

/* ---------- 일정 · 메모 (달력) ---------- */
let calY,calM,selDate,schedData={events:[],memos:[]},calCatFilter='전체';
let calCats=['개인','회사'];
function catIcon(name){if(name==='개인')return '👤';if(name==='회사')return '🏢';return '🏷️';}
async function loadCats(){
  try{const c=await pywebview.api.cal_cats();if(Array.isArray(c)&&c.length)calCats=c;}catch(_){}
  renderCatUI();
}
function renderCatUI(){
  // 일정 추가/수정 셀렉트
  ['evcat','editCat'].forEach(function(id){
    const sel=document.getElementById(id);if(!sel)return;
    const cur=sel.value;sel.innerHTML='';
    calCats.forEach(function(c){
      const o=document.createElement('option');o.value=c;o.textContent=catIcon(c)+' '+c;sel.appendChild(o);
    });
    if(calCats.indexOf(cur)>=0)sel.value=cur;
  });
  // 필터 버튼 줄
  const bar=document.getElementById('cfilter');if(!bar)return;
  bar.innerHTML='';
  const mk=function(label,val,active){
    const b=document.createElement('button');b.className='cfbtn'+(active?' active':'');
    b.id='cf'+val;b.textContent=label;b.onclick=function(){setCatFilter(val)};return b;};
  bar.appendChild(mk('전체','전체',calCatFilter==='전체'));
  calCats.forEach(function(c){bar.appendChild(mk(catIcon(c)+' '+c,c,calCatFilter===c));});
  const add=document.createElement('button');add.className='cfbtn cfadd';add.textContent='＋ 분류';
  add.title='새 분류 추가';add.onclick=addCatPrompt;bar.appendChild(add);
  // 분류가 2개(기본)보다 많으면 분류 삭제 버튼도
  if(calCats.length>1){
    const del=document.createElement('button');del.className='cfbtn cfadd';del.textContent='－ 분류';
    del.title='분류 삭제';del.onclick=delCatPrompt;bar.appendChild(del);
  }
}
async function addCatPrompt(){
  const name=await askInput('새 일정 분류 이름','예: 가족, 프로젝트A, 운동');
  if(!name||!name.trim())return;
  try{calCats=await pywebview.api.cal_cat_add(name.trim());}catch(_){}
  renderCatUI();
}
async function delCatPrompt(){
  if(calCats.length<=1)return;
  const name=await askInput('삭제할 분류 이름을 정확히 입력','현재: '+calCats.join(', '));
  if(!name||!name.trim())return;
  try{calCats=await pywebview.api.cal_cat_del(name.trim());}catch(_){}
  if(calCatFilter!=='전체'&&calCats.indexOf(calCatFilter)<0)calCatFilter='전체';
  renderCatUI();renderCal();renderDay();
}
const calmodal=document.getElementById('calmodal');
function pad2(n){return (n<10?'0':'')+n}
function dstr(y,m,d){return y+'-'+pad2(m+1)+'-'+pad2(d)}
function todayStr(){const n=new Date();return dstr(n.getFullYear(),n.getMonth(),n.getDate())}
function remLabel(o){o=+o;if(o===0)return '정각';if(o%1440===0)return (o===1440?'하루':(o/1440)+'일')+' 전';
  if(o%60===0)return (o/60)+'시간 전';return o+'분 전'}
async function openCal(){
  const n=new Date();calY=n.getFullYear();calM=n.getMonth();selDate=todayStr();
  calmodal.style.display='flex';
  await loadCats();
  await calRefresh();
  document.getElementById('evtitle').focus();
}
async function calRefresh(){
  if(calmodal.style.display!=='flex')return;
  schedData=await pywebview.api.sched_data();
  renderCal();renderDay();
}
function renderCal(){
  document.getElementById('calym').textContent=calY+'년 '+(calM+1)+'월';
  const g=document.getElementById('calgrid');g.innerHTML='';
  '일월화수목금토'.split('').forEach(function(w,i){
    const d=document.createElement('div');d.className='cdow'+(i===0?' sun':'');d.textContent=w;g.appendChild(d)});
  const start=new Date(calY,calM,1).getDay(),days=new Date(calY,calM+1,0).getDate(),tds=todayStr();
  for(let i=0;i<start;i++)g.appendChild(document.createElement('div'));
  for(let d=1;d<=days;d++){
    const ds=dstr(calY,calM,d);
    const evs=schedData.events.filter(function(e){return e.date===ds});
    const undone=evs.filter(function(e){return !e.done});
    const c=document.createElement('div');
    c.className='cday'+(ds===tds?' today':'')+(ds===selDate?' sel':'')+((start+d-1)%7===0?' sun':'');
    let inner='<span class="dn">'+d+'</span>';
    if(evs.length){
      const first=(undone[0]||evs[0]).title;
      inner+='<span class="dotc'+(undone.length?'':' alldone')+'">● '
        +esc(first)+(evs.length>1?' 외'+(evs.length-1):'')+'</span>';
    }
    c.innerHTML=inner;
    c.onclick=function(){selDate=ds;renderCal();renderDay()};
    g.appendChild(c);
  }
}
function renderDay(){
  const dl=document.getElementById('daylist');dl.innerHTML='';
  const head=document.createElement('div');
  head.style.cssText='font-size:12.5px;color:var(--mut);padding:0 4px 5px';
  head.textContent=(+selDate.slice(5,7))+'월 '+(+selDate.slice(8))+'일 ('
    +'일월화수목금토'[new Date(selDate+'T00:00:00').getDay()]+')';
  dl.appendChild(head);
  let evs=schedData.events.filter(function(e){return e.date===selDate});
  if(calCatFilter!=='전체')evs=evs.filter(function(e){return (e.category||'개인')===calCatFilter});
  evs=evs.sort(function(a,b){return (a.time||'99')<(b.time||'99')?-1:1});
  if(!evs.length){const e=document.createElement('div');
    e.style.cssText='color:var(--mut);font-size:12.5px;padding:2px 4px';
    e.textContent='일정 없음 — 아래에서 바로 추가하세요';dl.appendChild(e)}
  evs.forEach(function(ev){
    const r=document.createElement('div');r.className='evitem'+(ev.done?' done':'');
    const chk=document.createElement('input');chk.type='checkbox';chk.checked=!!ev.done;chk.title='완료 표시';
    chk.onchange=function(){pywebview.api.sched_set(ev.id,{done:chk.checked})
      .then(function(d){schedData=d;renderCal();renderDay()})};
    const t=document.createElement('span');t.className='et';t.textContent=ev.time||'종일';
    const ttl=document.createElement('span');ttl.className='ettl';ttl.textContent=ev.title;ttl.title=ev.title;
    const cat=ev.category||'개인';
    const catBadge=document.createElement('span');
    catBadge.className='ecat '+(cat==='회사'?'work':'pers');
    catBadge.textContent=catIcon(cat);catBadge.title=cat;
    const rem=document.createElement('span');rem.className='erem';
    rem.textContent=(ev.remind&&ev.remind.length)?'🔔 '+ev.remind.map(remLabel).join(', '):'';
    const edit=document.createElement('button');edit.className='evbtn';edit.textContent='✏';edit.title='수정';
    edit.onclick=function(){openEditModal(ev)};
    const del=document.createElement('button');del.className='evbtn';del.textContent='🗑';del.title='삭제';
    del.onclick=async function(){if(!await askConfirm('일정 "'+ev.title+'"을(를) 삭제할까요?'))return;
      schedData=await pywebview.api.sched_del(ev.id);renderCal();renderDay()};
    r.appendChild(chk);r.appendChild(t);r.appendChild(ttl);r.appendChild(catBadge);
    r.appendChild(rem);r.appendChild(edit);r.appendChild(del);
    dl.appendChild(r);
  });
}
async function addEvent(){
  const ttl=document.getElementById('evtitle').value.trim();
  if(!ttl)return;
  const rv=document.getElementById('evremind').value;
  const cat=document.getElementById('evcat').value;
  schedData=await pywebview.api.sched_add({date:selDate,
    time:document.getElementById('evtime').value,title:ttl,
    category:cat,
    remind:rv?rv.split(',').map(Number):[]});
  document.getElementById('evtitle').value='';
  renderCal();renderDay();
}
function setCatFilter(cat){
  calCatFilter=cat;
  ['전체'].concat(calCats).forEach(function(c){
    const b=document.getElementById('cf'+c);
    if(b)b.classList.toggle('active',c===cat);
  });
  renderCal();renderDay();
}
function openEditModal(ev){
  document.getElementById('editId').value=ev.id;
  document.getElementById('editDate').value=ev.date;
  document.getElementById('editTime').value=ev.time||'';
  document.getElementById('editTitle').value=ev.title;
  document.getElementById('editCat').value=ev.category||'개인';
  const rv=(ev.remind&&ev.remind.length)?ev.remind.join(','):'';
  document.getElementById('editRemind').value=rv;
  document.getElementById('editmodal').style.display='flex';
}
function renderMemos(){
  const ml=document.getElementById('memolist');ml.innerHTML='';
  if(!schedData.memos.length){
    ml.innerHTML='<div style="color:var(--mut);font-size:12.5px;padding:6px 2px">'
      +'메모가 없습니다 — 채팅에서 "○○ 메모해줘"라고 하거나 아래에 입력하세요</div>';return}
  schedData.memos.forEach(function(m){
    const r=document.createElement('div');r.className='mitem';
    const ts=document.createElement('span');ts.className='mts';ts.textContent=m.ts;
    const tx=document.createElement('span');tx.className='mtx';tx.textContent=m.text;
    const del=document.createElement('button');del.className='evbtn';del.textContent='🗑';del.title='삭제';
    del.onclick=async function(){if(!await askConfirm('이 메모를 삭제할까요?'))return;
      schedData=await pywebview.api.memo_del(m.id);renderMemos()};
    r.appendChild(ts);r.appendChild(tx);r.appendChild(del);ml.appendChild(r);
  });
}
async function addMemo(){
  const t=document.getElementById('memotext').value.trim();if(!t)return;
  schedData=await pywebview.api.memo_add(t);
  document.getElementById('memotext').value='';renderMemos();
}
const memomodal=document.getElementById('memomodal');
async function openMemo(){
  memomodal.style.display='flex';
  schedData=await pywebview.api.sched_data();
  renderMemos();
  document.getElementById('memotext').focus();
}
document.getElementById('calbtn').onclick=openCal;
document.getElementById('memobtn').onclick=openMemo;
/* ---------- 기능 대시보드 (할 수 있는 일) ---------- */
const FEATURES=[
  {icon:'📁',title:'파일 · 문서',items:[
    {t:'컴퓨터 전체에서 파일 찾기 (경로 몰라도 OK)',ex:'작년 견적서 찾아줘'},
    {t:'파일 내용 분석 — 엑셀·워드·PDF·한글·PPT·이미지',ex:'이 PDF 요약해줘'},
    {t:'문서 자동 생성 — 워드·엑셀·PPT·PDF',ex:'공사 견적서 양식 엑셀로 만들어줘'},
    {t:'이미지 속 글자 인식(OCR)',ex:'이 스크린샷 글자 읽어줘'}
  ]},
  {icon:'🎨',title:'디자인 · 도면',items:[
    {t:'CAD 도면(DXF) 분석·생성 — 레이어·치수 포함',ex:'3x4m 방 평면도 그려줘'},
    {t:'스케치업 모델 만들기 (Ruby 스크립트)',ex:'단독주택 매스 모델 스케치업으로 만들어줘'},
    {t:'포토샵·일러스트 자동화'}
  ]},
  {icon:'🌐',title:'인터넷',items:[
    {t:'웹 검색 — 뉴스·위키 (광고·차단 없음)',ex:'이번 주 원달러 환율 알려줘'},
    {t:'링크 본문 요약 — 주소만 붙이면 끝',ex:'이 기사 요약해줘 https://'},
    {t:'사이트 바로 열기·검색',ex:'유튜브에서 캠핑 브이로그 틀어줘'}
  ]},
  {icon:'💻',title:'컴퓨터 제어',items:[
    {t:'프로그램·파일 실행',ex:'스케치업 열어줘'},
    {t:'화면 자동화 — 타이핑·클릭'},
    {t:'코드 자동 실행 + 실패하면 스스로 고침'},
    {t:'답변이 길거나 빗나가면 ⏹ 로 즉시 중단'}
  ]},
  {icon:'🗓',title:'일정 · 기억',items:[
    {t:'일정 잡고 알림 받기',ex:'내일 3시 미팅 잡고 30분 전에 알려줘'},
    {t:'메모 저장',ex:'/메모 우유 사기'},
    {t:'취향·작업방식을 스스로 기억해 다음에 반영'},
    {t:'예전 대화를 의미로 검색',ex:'/찾기 견적서 단가'}
  ]},
  {icon:'🎙',title:'음성 · 입력',items:[
    {t:'음성으로 말하기 🎤 / 답변 음성으로 듣기 🔈'},
    {t:'연속 음성 대화(핸즈프리) 🔁 — 말로 주고받기'},
    {t:'이미지는 Ctrl+V 붙여넣기 · 화면 캡처 📷'}
  ]},
  {icon:'🔒',title:'보안 · 맞춤',items:[
    {t:'비밀번호·계좌를 암호화 금고에 보관',ex:'/금고 목록'},
    {t:'받은파일 폴더에 넣으면 자동 감지'},
    {t:'답변 말투·길이 바꾸기',ex:'/말투 간결'},
    {t:'위험 명령 자동 차단 · 원본 안전 보존'}
  ]}
];
function renderFeat(){
  const g=document.getElementById('featgrid');g.innerHTML='';
  FEATURES.forEach(function(cat){
    const c=document.createElement('div');c.className='featcat';
    const h=document.createElement('h4');h.textContent=cat.icon+'  '+cat.title;c.appendChild(h);
    const ul=document.createElement('ul');
    cat.items.forEach(function(it){
      const li=document.createElement('li');li.textContent=it.t;
      if(it.ex){
        li.appendChild(document.createElement('br'));
        const ex=document.createElement('span');ex.className='fex';ex.textContent='▷ '+it.ex;
        ex.onclick=function(){
          document.getElementById('featmodal').style.display='none';
          inp.value=it.ex;inp.style.height='auto';
          inp.style.height=Math.min(inp.scrollHeight,170)+'px';inp.focus();
        };
        li.appendChild(ex);
      }
      ul.appendChild(li);
    });
    c.appendChild(ul);g.appendChild(c);
  });
}
function openFeat(){renderFeat();document.getElementById('featmodal').style.display='flex'}
document.getElementById('featclose').onclick=function(){document.getElementById('featmodal').style.display='none'};
document.getElementById('featmodal').onclick=function(e){if(e.target.id==='featmodal')this.style.display='none'};
/* ── 좌측하단 설정 팝업 메뉴 ── */
(function(){
  var btn=document.getElementById('morebtn'),pop=document.getElementById('settingspop');
  btn.onclick=function(e){e.stopPropagation();pop.classList.toggle('open');};
  document.addEventListener('click',function(){pop.classList.remove('open');});
  document.getElementById('sp_feat').onclick=function(){pop.classList.remove('open');openFeat();};
  document.getElementById('sp_name').onclick=function(){pop.classList.remove('open');openNameModal();};
  document.getElementById('sp_update').onclick=async function(){
    pop.classList.remove('open');
    var btn2=this;btn2.textContent='🔄 확인 중…';btn2.disabled=true;
    var r=await pywebview.api.check_update_now();
    btn2.textContent='🔄 업데이트 확인';btn2.disabled=false;
    if(!r||r==='latest'){sysMsg('✅ 최신 버전입니다 ('+window._version+')');}
    else{var parts=r.split('|');showUpdateBanner(parts[0],parts[1],parts[2]||'');}
  };
})();
document.getElementById('calclose').onclick=function(){calmodal.style.display='none'};
document.getElementById('caltoday').onclick=function(){
  const n=new Date();calY=n.getFullYear();calM=n.getMonth();selDate=todayStr();renderCal();renderDay()};
document.getElementById('calprev').onclick=function(){calM--;if(calM<0){calM=11;calY--}renderCal()};
document.getElementById('calnext').onclick=function(){calM++;if(calM>11){calM=0;calY++}renderCal()};
document.getElementById('evadd').onclick=addEvent;
document.getElementById('evtitle').addEventListener('keydown',function(e){if(e.key==='Enter'&&!e.isComposing)addEvent()});
document.getElementById('memoclose').onclick=function(){memomodal.style.display='none'};
document.getElementById('memoadd').onclick=addMemo;
document.getElementById('memotext').addEventListener('keydown',function(e){if(e.key==='Enter'&&!e.isComposing)addMemo()});
calmodal.addEventListener('click',function(e){if(e.target===calmodal)calmodal.style.display='none'});
memomodal.addEventListener('click',function(e){if(e.target===memomodal)memomodal.style.display='none'});
/* ---------- 이름 설정 모달 ---------- */
function openNameModal(){
  pywebview.api.get_name().then(function(n){
    document.getElementById('nameInp').value=n||'';
    document.getElementById('namemodal').style.display='flex';
    setTimeout(function(){document.getElementById('nameInp').focus()},50);
  });
}
/* namebtn 이동됨 → sp_name 으로 대체 */
document.getElementById('namecancel').onclick=function(){document.getElementById('namemodal').style.display='none'};
document.getElementById('nameok').onclick=async function(){
  const n=document.getElementById('nameInp').value.trim();
  if(!n)return;
  const saved=await pywebview.api.save_name(n);
  document.getElementById('namemodal').style.display='none';
  document.querySelector('.brand').innerHTML='<span class="star">✱</span> '+esc(saved);
  sysMsg('⚙ 비서 이름을 "'+saved+'"(으)로 설정했습니다');
};
document.getElementById('nameInp').addEventListener('keydown',function(e){
  if(e.key==='Enter')document.getElementById('nameok').onclick()});
document.getElementById('namemodal').addEventListener('click',function(e){
  if(e.target.id==='namemodal')this.style.display='none'});
/* ---------- 일정 수정 모달 ---------- */
document.getElementById('editcancel').onclick=function(){document.getElementById('editmodal').style.display='none'};
document.getElementById('editsave').onclick=async function(){
  const id=+document.getElementById('editId').value;
  const fields={
    date:document.getElementById('editDate').value,
    time:document.getElementById('editTime').value,
    title:document.getElementById('editTitle').value.trim(),
    category:document.getElementById('editCat').value,
    remind:document.getElementById('editRemind').value?
      document.getElementById('editRemind').value.split(',').map(Number):[],
    notified:[]
  };
  if(!fields.title)return;
  schedData=await pywebview.api.sched_set(id,fields);
  document.getElementById('editmodal').style.display='none';
  renderCal();renderDay();
};
document.getElementById('editmodal').addEventListener('click',function(e){
  if(e.target.id==='editmodal')this.style.display='none'});

/* ---------- 첨부·캡처·붙여넣기 ---------- */
attachBtn.onclick=function(){pywebview.api.pick_files().then(showAttached)};
document.getElementById('capture').onclick=function(){
  sysMsg('📷 캡처 도구가 뜹니다 — 화면에서 영역을 드래그하세요 (창은 잠시 내려갑니다)');
  pywebview.api.capture_screen();
};
document.addEventListener('paste',function(e){
  const items=(e.clipboardData||{}).items||[];
  for(const it of items){
    if(it.type&&it.type.startsWith('image/')){
      e.preventDefault();
      const f=it.getAsFile();const r=new FileReader();
      r.onload=function(){pywebview.api.attach_clipboard_image(r.result).then(showAttached)};
      r.readAsDataURL(f);
      return;
    }
  }
});
showGreet(true);
/* ── 업데이트 배너 ── */
let _updUrl='',_updVer='';
function showUpdateBanner(ver,url,changelog){
  _updUrl=url;_updVer=ver;
  document.getElementById('updver').textContent=ver;
  const logEl=document.getElementById('updlog');
  logEl.textContent=changelog?'— '+changelog:'';
  document.getElementById('updbanner').classList.add('show');
}
function hideUpdate(){
  document.getElementById('updbanner').classList.remove('show');
  if(_updVer)pywebview.api.skip_update(_updVer);   // 이 버전은 다시 안 띄움
}
async function doUpdate(btn){
  btn.textContent='다운로드 중…';btn.disabled=true;
  const r=await pywebview.api.do_update(_updUrl);
  if(r&&r!=='ok'){
    btn.textContent='지금 업데이트';btn.disabled=false;
    sysMsg('⚠ '+r);
  }
}

/* ── 다크모드 ── */
function applyTheme(dark){
  document.body.classList.toggle('dark',dark);
  document.getElementById('darkbtn').textContent=dark?'☀':'🌙';
  try{pywebview.api.set_titlebar_color(dark);}catch(_){}
}
document.getElementById('darkbtn').onclick=function(){
  const dark=!document.body.classList.contains('dark');
  applyTheme(dark);pywebview.api.save_theme(dark);
};

/* ── 대화 내보내기 ── */
document.getElementById('exportbtn').onclick=async function(){
  const r=await pywebview.api.export_chat();
  sysMsg(r);
};

/* ── 클립보드 빠른 질문 ── */
document.getElementById('clip').onclick=async function(){
  const t=await pywebview.api.read_clipboard();
  if(!t||!t.trim()){sysMsg('📋 복사된 텍스트가 없습니다 — 먼저 내용을 복사(Ctrl+C)하세요');return}
  inp.value='다음 복사한 내용을 읽고 핵심만 보기좋게 정리해줘:\n\n'+t.trim();
  inp.style.height='auto';inp.style.height=Math.min(inp.scrollHeight,170)+'px';inp.focus();
};

/* ── 통합 설정 패널 ── */
var _setRate=1.05;
async function renderFavs(){
  const box=document.getElementById('favlist');box.innerHTML='';
  const favs=await pywebview.api.fav_folders();
  if(!favs.length){box.innerHTML='<div style="font-size:11.5px;color:var(--mut);padding:4px">등록된 폴더가 없습니다</div>';return}
  favs.forEach(function(p){
    const it=document.createElement('div');it.className='favitem';
    const name=p.replace(/\\/g,'/').split('/').filter(Boolean).pop()||p;
    const lbl=document.createElement('span');lbl.className='fpath';lbl.textContent='📂 '+name;lbl.title=p;
    const ob=document.createElement('button');ob.textContent='열기';
    ob.onclick=function(){pywebview.api.open_folder_path(p)};
    const cb=document.createElement('button');cb.textContent='정리';cb.title='이 폴더 종류별 정리';
    cb.onclick=async function(){const r=await pywebview.api.organize_folder_path(p);if(r){aiMsg(r);if(window._pikaDone)window._pikaDone();}};
    const db=document.createElement('button');db.className='fdel';db.textContent='✕';
    db.onclick=async function(){await pywebview.api.del_fav_folder(p);renderFavs()};
    it.appendChild(lbl);it.appendChild(ob);it.appendChild(cb);it.appendChild(db);box.appendChild(it);
  });
}
function openSettings(){
  const m=document.getElementById('setmodal');
  document.getElementById('set_dark').checked=document.body.classList.contains('dark');
  document.getElementById('set_pika').checked=!!window._pikaOn;
  document.getElementById('set_style').value=(window._answerStyle||'표준');
  document.getElementById('set_rate').value=_setRate;
  document.getElementById('set_rateval').textContent=(_setRate).toFixed(2)+'x';
  renderFavs();m.style.display='flex';
}
document.getElementById('setbtn').onclick=openSettings;
document.getElementById('setclose').onclick=function(){document.getElementById('setmodal').style.display='none'};
document.getElementById('set_dark').onchange=function(){applyTheme(this.checked);pywebview.api.save_theme(this.checked)};
document.getElementById('set_pika').onchange=function(){pikaShow(this.checked)};
document.getElementById('set_style').onchange=function(){
  window._answerStyle=this.value;pywebview.api.set_style(this.value).then(sysMsg)};
document.getElementById('set_rate').oninput=function(){
  _setRate=parseFloat(this.value);window._ttsRate=_setRate;
  document.getElementById('set_rateval').textContent=_setRate.toFixed(2)+'x';
  pywebview.api.save_tts_rate(_setRate);
};
document.getElementById('set_addfav').onclick=async function(){await pywebview.api.add_fav_folder();renderFavs()};

/* ── 워크플로우 보드 ── */
var wfData={projects:[]};var wfMode='proj';
var PRI=['높음','보통','낮음'];var STAT=['todo','doing','done'];
var STATLBL={todo:'할 일',doing:'진행중',done:'완료'};
async function openWf(){wfData=await pywebview.api.wf_data();renderWf();
  document.getElementById('wfmodal').style.display='flex';}
function _projProgress(p){var ts=p.tasks||[];var done=ts.filter(function(t){return t.status==='done'}).length;
  return{done:done,total:ts.length,pct:ts.length?Math.round(done/ts.length*100):0};}
function taskCard(p,t){
  var c=document.createElement('div');c.className='wftask pri-'+(t.priority||'보통')+(t.status==='done'?' done':'');
  var tt=document.createElement('div');tt.className='wtt';
  tt.textContent=(t.status==='doing'?'▶ ':t.status==='done'?'✓ ':'')+t.title;
  var meta=document.createElement('div');meta.className='wtmeta';
  // 우선순위 뱃지 버튼 (클릭 시 순환)
  var pri=t.priority||'보통';
  var priBadge=document.createElement('button');priBadge.className='wpri wpri-'+pri;priBadge.textContent=pri;
  priBadge.title='우선순위 변경';
  priBadge.onclick=async function(e){e.stopPropagation();
    var ni=(PRI.indexOf(t.priority||'보통')+1)%3;
    wfData=await pywebview.api.wf_set_task(p.id,t.id,{priority:PRI[ni]});renderWf();};
  meta.appendChild(document.createTextNode(''));
  var statSpan=document.createElement('span');statSpan.textContent=STATLBL[t.status||'todo'];meta.appendChild(statSpan);
  if(t.due){var dueSpan=document.createElement('span');dueSpan.textContent='📅 '+t.due;meta.appendChild(dueSpan);}
  meta.appendChild(priBadge);
  var x=document.createElement('span');x.className='wtx';x.textContent='✕';
  x.onclick=async function(e){e.stopPropagation();wfData=await pywebview.api.wf_del_task(p.id,t.id);renderWf()};
  c.onclick=async function(e){   // 클릭하면 상태 순환 todo→doing→done→todo
    if(e.target===priBadge)return;
    var ni=(STAT.indexOf(t.status||'todo')+1)%3;
    wfData=await pywebview.api.wf_set_task(p.id,t.id,{status:STAT[ni]});renderWf();
    if(STAT[ni]==='done'&&window._pikaDone)window._pikaDone();
  };
  c.title='클릭: 진행상태 변경';
  c.appendChild(tt);c.appendChild(meta);c.appendChild(x);return c;
}
function projColumn(p){
  var col=document.createElement('div');col.className='wfcol';col.style.borderTopColor=p.color||'#5B6CFF';
  var top=document.createElement('div');top.className='wfcoltop';
  var b=document.createElement('b');b.textContent=p.name;
  var fo=document.createElement('span');fo.className='wficon';
  if(p.folder){fo.textContent='📂';fo.title='연결 폴더 열기'+(p.file_count!=null?' ('+p.file_count+'개 파일)':'');
    fo.onclick=function(){pywebview.api.open_folder_path(p.folder)};}
  var del=document.createElement('span');del.className='wficon';del.textContent='🗑';del.title='프로젝트 삭제';
  del.onclick=async function(){if(await askConfirm('프로젝트 "'+p.name+'"을(를) 삭제할까요?')){
    wfData=await pywebview.api.wf_del_project(p.id);renderWf();}};
  top.appendChild(b);if(p.folder)top.appendChild(fo);top.appendChild(del);
  var pr=_projProgress(p);
  var bar=document.createElement('div');bar.className='wfbar';var fill=document.createElement('i');
  fill.style.width=pr.pct+'%';fill.style.background=p.color||'#5B6CFF';bar.appendChild(fill);
  var prog=document.createElement('div');prog.className='wfprog';
  prog.textContent=pr.done+'/'+pr.total+' 완료 ('+pr.pct+'%)'+(p.file_count!=null?' · 파일 '+p.file_count+'개':'');
  col.appendChild(top);col.appendChild(bar);col.appendChild(prog);
  var order={'높음':0,'보통':1,'낮음':2};
  (p.tasks||[]).slice().sort(function(a,b2){return (order[a.priority]||1)-(order[b2.priority]||1)})
    .forEach(function(t){col.appendChild(taskCard(p,t))});
  var add=document.createElement('button');add.className='wfaddtask';add.textContent='＋ 업무 추가';
  add.onclick=async function(){var title=await askInput('업무 추가','업무 이름 (예: 평면도 수정)');if(!title)return;
    wfData=await pywebview.api.wf_add_task(p.id,title,'보통','');renderWf();};
  col.appendChild(add);return col;
}
function renderWf(){
  var board=document.getElementById('wfboard');board.innerHTML='';
  if(!wfData.projects||!wfData.projects.length){
    board.innerHTML='<div class="wfempty">아직 프로젝트가 없습니다.<br>오른쪽 위 <b>＋ 프로젝트</b>로 시작하거나, 채팅에 "○○ 프로젝트 만들어줘"라고 하세요.</div>';return;}
  if(wfMode==='proj'){
    wfData.projects.forEach(function(p){board.appendChild(projColumn(p))});
  }else{   // 우선순위별: 높음/보통/낮음 컬럼에 모든 프로젝트 업무를 모음
    PRI.forEach(function(pri){
      var col=document.createElement('div');col.className='wfcol wfpgroup';
      col.style.borderTopColor=(pri==='높음'?'#e8473f':pri==='보통'?'#f0993e':'#9aa0aa');
      var top=document.createElement('div');top.className='wfcoltop';
      var b=document.createElement('b');b.textContent='중요도: '+pri;top.appendChild(b);col.appendChild(top);
      var any=false;
      wfData.projects.forEach(function(p){
        (p.tasks||[]).filter(function(t){return (t.priority||'보통')===pri}).forEach(function(t){
          any=true;var c=taskCard(p,t);
          var tag=document.createElement('div');tag.className='wtmeta';
          tag.innerHTML='<span style="color:'+(p.color||'#5B6CFF')+'">● '+p.name+'</span>';
          c.insertBefore(tag,c.firstChild);col.appendChild(c);});
      });
      if(!any){var e=document.createElement('div');e.className='wfprog';e.textContent='해당 업무 없음';col.appendChild(e);}
      board.appendChild(col);
    });
  }
}
document.getElementById('wfbtn').onclick=openWf;
document.getElementById('wfclose').onclick=function(){document.getElementById('wfmodal').style.display='none'};
document.getElementById('wf_byproj').onclick=function(){wfMode='proj';
  this.classList.add('on');document.getElementById('wf_bypri').classList.remove('on');renderWf();};
document.getElementById('wf_bypri').onclick=function(){wfMode='pri';
  this.classList.add('on');document.getElementById('wf_byproj').classList.remove('on');renderWf();};
document.getElementById('wf_addproj').onclick=async function(){
  var name=await askInput('새 프로젝트','프로젝트 이름 (예: 강남 카페 리노베이션)');if(!name)return;
  var link=await askConfirm('이 프로젝트에 작업 폴더를 연결할까요? (진척도에 파일 수가 표시됩니다)');
  var folder='';if(link){folder=await pywebview.api.wf_pick_folder();}
  wfData=await pywebview.api.wf_add_project(name,null,folder);renderWf();};
window._wfRefresh=async function(){wfData=await pywebview.api.wf_data();
  if(document.getElementById('wfmodal').style.display==='flex')renderWf();renderRail();};

/* 우측 상시 패널 */
function _priColor(p){return p==='높음'?'#e8473f':p==='보통'?'#f0993e':'#9aa0aa'}
function renderRail(){
  var body=document.getElementById('wfrbody');if(!body)return;body.innerHTML='';
  var projs=(wfData.projects||[]);
  if(!projs.length){body.innerHTML='<div class="wfrempty">진행 중인 프로젝트가 없습니다.<br>채팅에 "○○ 프로젝트 만들어줘"라고 하거나 ⤢ 로 추가하세요.</div>';return;}
  projs.forEach(function(p){
    var pr=_projProgress(p);
    var box=document.createElement('div');box.className='wfrproj';
    var nm=document.createElement('div');nm.className='wfrname';
    nm.innerHTML='<i style="background:'+(p.color||'#5B6CFF')+'"></i>'+p.name;
    var bar=document.createElement('div');bar.className='wfrbar';
    var sp=document.createElement('span');sp.style.width=pr.pct+'%';sp.style.background=p.color||'#5B6CFF';bar.appendChild(sp);
    var pct=document.createElement('div');pct.className='wfrpct';
    pct.textContent=pr.done+'/'+pr.total+' ('+pr.pct+'%)'+(p.file_count!=null?' · 파일 '+p.file_count:'');
    box.appendChild(nm);box.appendChild(bar);box.appendChild(pct);
    var order={'높음':0,'보통':1,'낮음':2};
    var open=(p.tasks||[]).filter(function(t){return t.status!=='done'})
      .sort(function(a,b){return (order[a.priority]||1)-(order[b.priority]||1)}).slice(0,4);
    open.forEach(function(t){
      var el=document.createElement('div');el.className='wfrtask pri-'+(t.priority||'보통')+(t.status==='doing'?' doing':'');
      el.textContent=(t.status==='doing'?'▶ ':'')+t.title;el.title='클릭: 진행상태 변경';
      el.onclick=async function(){var ni=(STAT.indexOf(t.status||'todo')+1)%3;
        wfData=await pywebview.api.wf_set_task(p.id,t.id,{status:STAT[ni]});
        if(STAT[ni]==='done'&&window._pikaDone)window._pikaDone();
        if(document.getElementById('wfmodal').style.display==='flex')renderWf();renderRail();};
      box.appendChild(el);
    });
    body.appendChild(box);
  });
}
function showRail(on){
  document.getElementById('wfrail').classList.toggle('hidden',!on);
  document.getElementById('wfrtab').style.display=on?'none':'block';
  pywebview.api.save_wf_rail(on);
}
document.getElementById('wfrhide').onclick=function(){showRail(false)};
document.getElementById('wfrtab').onclick=function(){showRail(true);window._wfRefresh()};
document.getElementById('wfropen').onclick=openWf;

/* ── 폴더 정리 모달 ── */
const cleanmodal=document.getElementById('cleanmodal');
async function openClean(){
  cleanmodal.style.display='flex';
  document.getElementById('cleanresult').textContent='';
  const box=document.getElementById('cleanquick');box.innerHTML='<span style="font-size:12px;color:var(--mut)">불러오는 중…</span>';
  let list=[];try{list=await pywebview.api.clean_folders();}catch(_){}
  box.innerHTML='';
  list.forEach(function(f){
    const c=document.createElement('button');c.className='cleancard';
    c.innerHTML='<div class="cleft"><div>'+esc(f.label)+'</div><div class="cpath">'+esc(f.path)+'</div></div>';
    c.onclick=function(){runClean(f.path);};
    box.appendChild(c);
  });
}
async function runClean(path){
  const res=document.getElementById('cleanresult');
  if(!path){
    // 폴더 선택 후 즉시 미리보기
    res.textContent='📁 폴더를 선택해주세요…';
    let r='';try{r=await pywebview.api.organize_folder();}catch(_){r='⚠ 정리 중 오류가 발생했습니다.';}
    if(window._pikaDone)window._pikaDone();
    if(r){res.textContent=r;aiMsg(r);}else{res.textContent='취소되었습니다.';}
    return;
  }
  res.innerHTML='<span style="color:var(--mut);font-size:12px">미리보기 불러오는 중…</span>';
  let pv;try{pv=await pywebview.api.organize_preview(path);}catch(_){pv=null;}
  if(!pv||pv.error){res.textContent='⚠ '+(pv&&pv.error||'오류');return;}
  const cats=Object.keys(pv.preview||{});
  if(!cats.length){res.textContent='📁 이동할 파일이 없습니다.';return;}
  let html='<div style="font-size:13px;margin-bottom:10px"><b>📋 정리 미리보기</b> — <span style="color:var(--mut)">'+esc(path)+'</span></div>';
  const total=cats.reduce(function(s,c){return s+(pv.preview[c]||[]).length;},0);
  html+='<div style="font-size:12px;color:var(--mut);margin-bottom:8px">총 <b>'+total+'개</b> 파일을 종류별 폴더로 이동합니다.</div>';
  cats.forEach(function(c){
    const files=pv.preview[c]||[];
    html+='<div style="margin-bottom:6px"><span style="font-size:12.5px;font-weight:600">📂 '+esc(c)+'</span> <span style="font-size:11.5px;color:var(--mut)">('+files.length+'개)</span>';
    html+='<div style="font-size:11px;color:var(--mut);margin-top:2px;margin-left:8px">'+files.slice(0,5).map(esc).join(', ')+(files.length>5?' 외 '+(files.length-5)+'개':'')+'</div></div>';
  });
  html+='<div style="display:flex;gap:8px;margin-top:12px">';
  html+='<button id="cleanconfirm" class="dashbtn primary" style="flex:1">✅ 승인 — 정리 시작</button>';
  html+='<button id="cleanabort" class="dashbtn" style="flex:1">✕ 취소</button>';
  html+='</div>';
  res.innerHTML=html;
  document.getElementById('cleanabort').onclick=function(){res.innerHTML='';};
  document.getElementById('cleanconfirm').onclick=async function(){
    res.innerHTML='<span style="color:var(--mut);font-size:12px">🧹 정리 중…</span>';
    if(window._pikaWork)window._pikaWork();
    let r='';try{r=await pywebview.api.organize_folder_path(path);}catch(_){r='⚠ 정리 중 오류가 발생했습니다.';}
    if(window._pikaDone)window._pikaDone();
    if(r){res.textContent=r;aiMsg(r);}else{res.textContent='완료.';}
  };
}
document.getElementById('cleanbtn').onclick=openClean;
document.getElementById('cleanclose').onclick=function(){cleanmodal.style.display='none'};
document.getElementById('cleanpick').onclick=function(){runClean('');};
cleanmodal.addEventListener('click',function(e){if(e.target===cleanmodal)cleanmodal.style.display='none'});

/* ── 쇼핑·주문 모달 ── */
const shopmodal=document.getElementById('shopmodal');
function setShop(t){document.getElementById('shopq').value=t;document.getElementById('shopq').focus();}
function openShop(){
  shopmodal.style.display='flex';
  setTimeout(function(){document.getElementById('shopq').focus();},100);
}
function runShop(){
  const q=document.getElementById('shopq').value.trim();
  if(!q){document.getElementById('shopq').focus();return;}
  const SHOP_URLS={
    '쿠팡':'https://www.coupang.com/np/search?q=',
    '네이버쇼핑':'https://search.shopping.naver.com/search/all?query=',
    '11번가':'https://search.11st.co.kr/Search.tmall?kwd=',
    'G마켓':'https://browse.gmarket.co.kr/search?keyword='
  };
  const sites=[];
  shopmodal.querySelectorAll('.shopsites input:checked').forEach(function(c){sites.push(c.value);});
  if(!sites.length)sites.push('쿠팡','네이버쇼핑');
  shopmodal.style.display='none';
  const eq=encodeURIComponent(q);
  sites.forEach(function(s){
    const base=SHOP_URLS[s];
    if(base)pywebview.api.open_url(base+eq);
  });
}
document.getElementById('shopbtn').onclick=openShop;
document.getElementById('shopclose').onclick=function(){shopmodal.style.display='none'};
document.getElementById('shopgo').onclick=runShop;
document.getElementById('shopq').addEventListener('keydown',function(e){if(e.key==='Enter')runShop();});
shopmodal.addEventListener('click',function(e){if(e.target===shopmodal)shopmodal.style.display='none'});

/* ── 응답 시간·모델 표시 ── */
let _genStart=0;
function showGenMeta(){
  if(!_genStart||!cur||!cur.row)return;
  const sec=((performance.now()-_genStart)/1000).toFixed(1);
  let meta=cur.row.querySelector('.genmeta');
  if(!meta){meta=document.createElement('div');meta.className='genmeta';cur.row.appendChild(meta)}
  meta.textContent='⏱ '+sec+'초 · '+(window._modelName||'');
  _genStart=0;
}

/* ── 피카츄: 메인 창 밖 별도 위젯 창에서 떠다닌다(아래는 그 창을 제어) ── */
window._pikaOn=true;
window._pikaWork=function(){try{pywebview.api.pika_state('work')}catch(_){}};
window._pikaDone=function(){try{pywebview.api.pika_state('done')}catch(_){}};
/* 설정 토글에서 호출 — 외부 피카츄 창 표시/숨김 */
function pikaShow(on){window._pikaOn=!!on;try{pywebview.api.pika_toggle(!!on)}catch(_){}}

/* ── 링크를 시스템 브라우저로 열기 (pywebview 대응) ── */
document.addEventListener('click',function(e){
  const a=e.target.closest&&e.target.closest('a[href^="http"]');
  if(a){e.preventDefault();pywebview.api.open_url(a.getAttribute('href'))}
});

window.addEventListener('pywebviewready',function(){pywebview.api.on_ready()});
</script>
</body></html>"""


# ── 피카츄 떠다니는 외부 위젯 창 (메인 창 밖, 항상 위) ──
PIKA_HTML = r"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
*{margin:0;padding:0;box-sizing:border-box}
html,body{width:100%;height:100%;background:rgba(0,0,0,0) !important;overflow:hidden;
  font-family:'Malgun Gothic',sans-serif;user-select:none;-webkit-user-select:none}
#wrap{position:absolute;inset:0;display:flex;flex-direction:column;align-items:center;
  justify-content:flex-end;padding:6px;cursor:grab;background:rgba(0,0,0,0)}
#wrap:active{cursor:grabbing}
#bubble{background:#fff;color:#1f1f1f;border:2px solid #2b2b2b;border-radius:14px;
  padding:7px 11px;font-size:12px;line-height:1.4;max-width:150px;text-align:center;
  margin-bottom:8px;box-shadow:0 3px 10px rgba(0,0,0,.18);font-weight:600;position:relative;
  opacity:0;transform:translateY(6px) scale(.96);transition:.22s;pointer-events:none}
#bubble.show{opacity:1;transform:none}
#bubble:after{content:"";position:absolute;left:50%;bottom:-9px;transform:translateX(-50%);
  border:7px solid transparent;border-top-color:#2b2b2b}
#bubble:before{content:"";position:absolute;left:50%;bottom:-6px;transform:translateX(-50%);
  border:6px solid transparent;border-top-color:#fff;z-index:1}
.pchar{position:relative;width:84px;height:78px;cursor:pointer;
  filter:drop-shadow(0 6px 14px rgba(0,0,0,.3))}
.st-work .pchar{animation:pwork .6s ease-in-out infinite}
.st-done .pchar{animation:pdone .5s ease-in-out 3}
@keyframes pwork{0%,100%{transform:translateY(0) rotate(-2deg)}50%{transform:translateY(-3px) rotate(2deg)}}
@keyframes pdone{0%,100%{transform:translateY(0) scale(1)}50%{transform:translateY(-9px) scale(1.06)}}
.pear{position:absolute;top:-24px;width:14px;height:42px;background:#f9d423;border-radius:50% 50% 40% 40%;border:2px solid #3a2f0b}
.pear.l{left:14px;transform:rotate(-18deg);transform-origin:bottom}
.pear.r{right:14px;transform:rotate(18deg);transform-origin:bottom}
.pear:after{content:"";position:absolute;top:-2px;left:-2px;right:-2px;height:15px;background:#2b2b2b;border-radius:50% 50% 40% 40%}
.pface{position:absolute;top:6px;left:6px;width:72px;height:64px;background:#f9d423;border-radius:50% 50% 48% 48%;border:2.5px solid #3a2f0b}
.peye{position:absolute;top:20px;width:14px;height:16px;background:#2b2b2b;border-radius:50%}
.peye.l{left:14px}.peye.r{right:14px}
.peye:after{content:"";position:absolute;top:2px;left:2px;width:5px;height:5px;background:#fff;border-radius:50%}
.pcheek{position:absolute;top:32px;width:15px;height:13px;background:#e8473f;border-radius:50%;opacity:.92}
.pcheek.l{left:4px}.pcheek.r{right:4px}
.pmouth{position:absolute;top:40px;left:50%;transform:translateX(-50%);width:14px;height:7px;
  border:2.2px solid #3a2f0b;border-top:none;border-radius:0 0 14px 14px}
.st-work .peye{height:13px;top:21px}
.st-work .pmouth{width:8px;height:8px;border:2.2px solid #3a2f0b;border-radius:50%}
.st-done .peye{height:8px;background:transparent;border:2.5px solid #2b2b2b;border-bottom:none;
  top:22px;width:13px;border-top-left-radius:10px;border-top-right-radius:10px}
.st-done .peye:after{display:none}
.st-done .pmouth{width:18px;height:11px;background:#e8473f;border:2.2px solid #3a2f0b;border-top:none}
#pclose{position:absolute;top:2px;right:2px;width:18px;height:18px;border-radius:50%;
  background:#fff;border:1.5px solid #888;color:#555;font-size:11px;line-height:15px;text-align:center;
  cursor:pointer;opacity:0;transition:.15s;z-index:5}
#wrap:hover #pclose{opacity:.9}
#pclose:hover{background:#e8473f;color:#fff;border-color:#e8473f}
</style></head>
<body>
<div id="wrap">
  <button id="pclose" title="피카츄 숨기기">✕</button>
  <div id="bubble"></div>
  <div class="pchar" id="pchar">
    <div class="pear l"></div><div class="pear r"></div>
    <div class="pface">
      <div class="peye l"></div><div class="peye r"></div>
      <div class="pcheek l"></div><div class="pcheek r"></div>
      <div class="pmouth"></div>
    </div>
  </div>
</div>
<script>
var wrap=document.getElementById('wrap'),bubble=document.getElementById('bubble'),
    pchar=document.getElementById('pchar');
var MSG={idle:['피카! 무엇을 도와드릴까요?','대표님, 시킬 일이 있나요? 피카츄!','클릭하면 비서 창이 떠요 ⚡','오늘도 함께 일해요, 피카!'],
  work:['피카피카! 작업 중… ⚡','지금 처리하고 있어요!','조금만 기다려 주세요, 피카!'],
  done:['피카츄! 다 했어요 ⚡','완료했습니다, 대표님!','짠! 작업 끝났어요 ✨']};
var _t=null;
function pick(a){return a[Math.floor(Math.random()*a.length)];}
window.setState=function(state,text){
  wrap.classList.remove('st-idle','st-work','st-done');
  wrap.classList.add('st-'+state);
  bubble.textContent=text||pick(MSG[state]||MSG.idle);
  bubble.classList.add('show');
  if(_t)clearTimeout(_t);
  if(state!=='work'){_t=setTimeout(function(){bubble.classList.remove('show');
    wrap.classList.remove('st-work','st-done');},4500);}
};
var _moved=false,_dx=0,_dy=0;
pchar.addEventListener('mousedown',function(e){_moved=false;_dx=e.screenX;_dy=e.screenY;});
pchar.addEventListener('mouseup',function(e){
  if(Math.abs(e.screenX-_dx)<4&&Math.abs(e.screenY-_dy)<4){
    try{pywebview.api.focus();}catch(_){}
    window.setState('idle');
  }
});
document.getElementById('pclose').addEventListener('click',function(e){
  e.stopPropagation();try{pywebview.api.hide();}catch(_){}
});
window.addEventListener('pywebviewready',function(){setTimeout(function(){window.setState('idle');},500);});
</script>
</body></html>"""


def load_settings():
    """설정.json 을 읽어 dict 로 반환 (없거나 깨지면 빈 dict)."""
    try:
        with open(SETTINGS, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def save_settings(patch):
    """설정.json 에 patch(dict)를 병합 저장."""
    try:
        s = load_settings()
        s.update(patch)
        os.makedirs(HOME, exist_ok=True)
        tmp = SETTINGS + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(s, f, ensure_ascii=False)
        os.replace(tmp, SETTINGS)
    except Exception as e:
        log(f"설정 저장 실패: {e}")


def wf_load():
    """워크플로우(프로젝트·업무) 데이터를 읽어 반환. 없으면 빈 구조."""
    try:
        with open(WORKFLOW, encoding="utf-8") as f:
            d = json.load(f)
            if isinstance(d, dict) and isinstance(d.get("projects"), list):
                return d
    except Exception:
        pass
    return {"projects": []}


def wf_save(d):
    """워크플로우 데이터를 원자적으로 저장."""
    try:
        os.makedirs(HOME, exist_ok=True)
        tmp = WORKFLOW + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(d, f, ensure_ascii=False, indent=1)
        os.replace(tmp, WORKFLOW)
    except Exception as e:
        log(f"워크플로우 저장 실패: {e}")
    return d


def check_for_update():
    """GitHub version.json 을 읽어 새 버전이 있으면 (latest, download_url, changelog) 반환, 없으면 None.
    사용자가 '나중에'로 건너뛴 버전이면 None 반환."""
    if not UPDATE_JSON_URL:
        return None
    try:
        req = urllib.request.Request(UPDATE_JSON_URL,
                                     headers={"User-Agent": f"KJH비서/{VERSION}"})
        data = json.loads(urllib.request.urlopen(req, timeout=10).read().decode())
        latest = data.get("version", "")
        if latest and latest != VERSION:
            if latest == load_settings().get("skip_version"):
                return None     # 사용자가 이 버전은 건너뛰기로 함
            return latest, data.get("download_url", ""), data.get("changelog", "")
    except Exception as e:
        log(f"업데이트 확인 실패: {e}")
    return None


class Api:
    def __init__(self):
        self._window = None
        self._messages = [{"role": "system", "content": build_system()}]
        self._attached = []
        self._model = None
        self._model_vision = None
        self._run_count = 0
        self._num_predict = 800     # 답변 길이 한도 (잡담은 brief 모드로 줄여 더 빠르게)
        self._stop = False          # 생성 중단 플래그 (Stop 버튼)
        self._summary = ""          # 오래된 대화 요약 (긴 대화 자동 압축)
        self._summary_n = 0         # 요약에 이미 접힌 비-시스템 메시지 수
        self._session = os.path.join(HIST_DIR, time.strftime("대화_%Y%m%d_%H%M%S.json"))

    def _js(self, code):
        try:
            self._window.evaluate_js(code)
        except Exception:
            pass

    def _ensure_shortcut(self):
        """바탕화면·시작메뉴 바로가기가 없으면 자동 생성한다(설치 프로그램의 보강책).
        설치 .bat 의 PowerShell 바로가기 단계는 한글 경로(AI비서)가 cmd→PS 로 넘어가며
        cp949 로 깨질 수 있어 일부 PC에서 실패한다. 이 메서드는 Python(UTF-8)에서 직접
        만들어 그 문제를 피한다.
        - 설치 위치(HOME)에서 실행될 때만 동작(개발본 실행 시엔 건드리지 않음).
        - 바탕화면에 이미 아이콘이 있으면 만들지 않음(중복 방지).
        - 한 번 처리하면 플래그로 재실행 안 함(사용자가 지워도 다시 안 만듦)."""
        try:
            script = os.path.abspath(sys.argv[0])
            if os.path.dirname(script).lower() != HOME.lower():
                return
            if load_settings().get("shortcut_made"):
                return
            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            # 설치본(KJH비서.lnk)·기존 사용본(AI 비서.lnk) 중 하나라도 있으면 중복 생성 안 함
            if any(os.path.exists(os.path.join(desktop, n))
                   for n in ("KJH비서.lnk", "AI 비서.lnk")):
                save_settings({"shortcut_made": True})
                return
            pyw = sys.executable
            if pyw.lower().endswith("python.exe"):
                cand = os.path.join(os.path.dirname(pyw), "pythonw.exe")
                if os.path.exists(cand):
                    pyw = cand
            ico = os.path.join(HOME, "비서.ico")
            import win32com.client
            ws = win32com.client.Dispatch("WScript.Shell")
            targets = [os.path.join(desktop, "KJH비서.lnk")]
            appdata = os.environ.get("APPDATA")
            if appdata:
                targets.append(os.path.join(
                    appdata, r"Microsoft\Windows\Start Menu\Programs", "KJH비서.lnk"))
            for lnk in targets:
                try:
                    sc = ws.CreateShortcut(lnk)
                    sc.TargetPath = pyw
                    sc.Arguments = f'"{script}"'
                    sc.WorkingDirectory = HOME
                    if os.path.exists(ico):
                        sc.IconLocation = ico + ",0"
                    sc.Hotkey = "CTRL+ALT+K"
                    sc.Description = "KJH비서 — 무료 로컬 AI 비서 (어디서든 Ctrl+Alt+K)"
                    sc.Save()
                except Exception:
                    pass
            save_settings({"shortcut_made": True})
        except Exception as e:
            log(f"바로가기 생성 건너뜀: {e}")

    def _sysmsg(self, t):
        self._js(f"sysMsg({json.dumps(t)})")

    def _status(self, t, cls=""):
        self._js(f"setStatus({json.dumps(t)},{json.dumps(cls)})")

    def _save_history(self):
        os.makedirs(HIST_DIR, exist_ok=True)
        slim = [m for m in self._messages if m["role"] != "system"]
        slim = [{k: v for k, v in m.items() if k != "images"} for m in slim]
        with open(self._session, "w", encoding="utf-8") as f:
            json.dump(slim, f, ensure_ascii=False, indent=1)

    def _restore_last(self):
        os.makedirs(HIST_DIR, exist_ok=True)
        files = sorted(globmod.glob(os.path.join(HIST_DIR, "대화_*.json")))
        if not files:
            return False
        last = files[-1]
        try:
            # 1시간 넘은 대화는 자동으로 이어가지 않음 (옛 주제가 새 질문에 섞이는 것 방지)
            if time.time() - os.path.getmtime(last) > 3600:
                return False
            with open(last, encoding="utf-8") as f:
                old = json.load(f)
        except Exception:
            return False
        if not old:
            return False
        self._messages += old
        self._session = last
        for m in old:
            txt = m["content"].split("\n\n=== ")[0]
            if m["role"] == "user":
                self._js(f"userMsg({json.dumps(txt)})")
            else:
                self._js(f"aiMsg({json.dumps(txt)})")
        self._sysmsg(f"이전 대화를 이어갑니다 — {os.path.basename(last)}")
        return True

    # ---------- JS API ----------
    def on_ready(self):
        threading.Thread(target=self._init, daemon=True).start()

    def current_session(self):
        return os.path.basename(self._session)

    def save_font(self, n):
        try:
            s = {}
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
            s["font"] = int(n)
            with open(SETTINGS, "w", encoding="utf-8") as f:
                json.dump(s, f)
        except Exception:
            pass
        return True

    def save_qbuttons(self, btns):
        try:
            s = {}
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
            s["buttons"] = btns
            with open(SETTINGS, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False)
        except Exception:
            pass
        return True

    def set_style(self, style):
        """답변 말투·길이 변경 (간결/표준/자세히) — 설정 저장 + 시스템 프롬프트 갱신."""
        global ANSWER_STYLE
        if style not in STYLE_HINT:
            return "간결 · 표준 · 자세히 중에서 골라 주세요"
        ANSWER_STYLE = style
        try:
            s = {}
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
            s["answer_style"] = style
            with open(SETTINGS, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False)
        except Exception:
            pass
        self._messages[0] = {"role": "system", "content": build_system()}
        self._sys_date = time.strftime("%Y-%m-%d")
        return f"말투를 '{style}' (으)로 바꿨습니다"

    def save_name(self, name):
        """AI 비서 이름 변경 — 설정 저장 + 시스템프롬프트 갱신"""
        global ASSISTANT_NAME
        ASSISTANT_NAME = (name or "KJH비서").strip() or "KJH비서"
        try:
            s = {}
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
            s["name"] = ASSISTANT_NAME
            with open(SETTINGS, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False)
        except Exception:
            pass
        self._messages[0] = {"role": "system", "content": build_system()}
        self._sys_date = time.strftime("%Y-%m-%d")
        return ASSISTANT_NAME

    def get_name(self):
        return ASSISTANT_NAME

    def get_version(self):
        return VERSION

    def do_update(self, download_url):
        """업데이트 파일을 다운받아 검증 후, 현재 버전을 백업하고 교체 bat 로 재시작한다.
        교체 실패 시 자동으로 백업본을 복구한다."""
        try:
            upd_dir = os.path.join(HOME, "업데이트")
            os.makedirs(upd_dir, exist_ok=True)
            new_pyw = os.path.join(upd_dir, "로컬비서_new.pyw")
            urllib.request.urlretrieve(download_url, new_pyw)
            # 다운로드 검증 — 너무 작거나 파이썬 코드가 아니면 거부 (깨진 파일 방지)
            with open(new_pyw, encoding="utf-8", errors="ignore") as f:
                head = f.read(4000)
            if os.path.getsize(new_pyw) < 10000 or "VERSION" not in head:
                return "업데이트 파일이 올바르지 않습니다 (다운로드 실패). 잠시 후 다시 시도하세요."
            current_pyw = os.path.abspath(__file__)
            backup_pyw = os.path.join(upd_dir, "이전버전_백업.pyw")
            try:
                shutil.copy2(current_pyw, backup_pyw)   # 롤백용 백업
            except Exception:
                backup_pyw = ""
            bat = os.path.join(upd_dir, "업데이트실행.bat")
            # cp949: Windows cmd 는 기본 코드페이지 사용. 교체 실패 시 백업 자동 복구.
            with open(bat, "w", encoding="cp949", errors="replace") as f:
                f.write(
                    "@echo off\r\n"
                    "timeout /t 2 /nobreak >nul\r\n"
                    f'copy /y "{new_pyw}" "{current_pyw}" >nul\r\n'
                    "if errorlevel 1 (\r\n"
                    + (f'  copy /y "{backup_pyw}" "{current_pyw}" >nul\r\n' if backup_pyw else "")
                    + "  echo 업데이트 실패 - 이전 버전으로 복구했습니다.\r\n"
                    "  pause\r\n"
                    ")\r\n"
                    f'start "" pythonw.exe "{current_pyw}"\r\n'
                    'del "%~f0"\r\n'
                )
            subprocess.Popen(f'start "" /min cmd /c "{bat}"', shell=True)
            os._exit(0)
        except Exception as e:
            return f"업데이트 실패: {e}"
        return "ok"

    def skip_update(self, version):
        """이 버전 업데이트를 건너뛴다 (다음부터 알림 안 뜸)."""
        save_settings({"skip_version": version})
        return "ok"

    def check_update_now(self):
        """사용자가 수동으로 업데이트 확인 버튼 클릭 시 호출.
        새 버전 있으면 'ver|download_url|changelog', 최신이면 'latest' 반환."""
        try:
            result = check_for_update()
            if result is None:
                return "latest"
            ver, url, changelog = result
            return f"{ver}|{url}|{changelog}"
        except Exception as e:
            log(f"수동 업데이트 확인 실패: {e}")
            return "latest"

    def save_theme(self, dark):
        save_settings({"dark": bool(dark)})
        return "ok"

    def set_titlebar_color(self, dark):
        """타이틀바 색상을 테마에 맞게 설정한다 (Windows 11 DWM)."""
        try:
            DWMWA_CAPTION_COLOR = 35
            # BGR 포맷: 라이트 #FAF9F5 → 0x00F5F9FA, 다크 #1E1E1C → 0x001C1E1E
            color = 0x001C1E1E if dark else 0x00F5F9FA
            hwnd = ctypes.windll.user32.FindWindowW(None, "KJH비서")
            if hwnd:
                ctypes.windll.dwmapi.DwmSetWindowAttribute(
                    hwnd, DWMWA_CAPTION_COLOR,
                    ctypes.byref(ctypes.c_int(color)),
                    ctypes.sizeof(ctypes.c_int))
        except Exception as e:
            log(f"타이틀바 색 설정 실패: {e}")
        return "ok"

    def save_pika(self, on):
        save_settings({"pika": bool(on)})
        return "ok"

    def save_pika_pos(self, x, y):
        save_settings({"pika_x": int(x), "pika_y": int(y)})
        return "ok"

    def pika_state(self, state):
        """피카츄 외부 창에 상태(idle/work/done)를 전달한다."""
        try:
            if getattr(self, "_pika", None):
                self._pika.evaluate_js(f"window.setState&&setState({json.dumps(state)})")
        except Exception:
            pass
        return "ok"

    def pika_toggle(self, on):
        """피카츄 외부 창을 표시/숨김 하고 설정에 저장한다."""
        save_settings({"pika": bool(on)})
        try:
            if getattr(self, "_pika", None):
                if on:
                    self._pika.show()
                else:
                    self._pika.hide()
        except Exception:
            pass
        return "ok"

    def save_wf_rail(self, on):
        save_settings({"wf_rail": bool(on)})
        return "ok"

    def cal_cats(self):
        """일정 분류 목록을 반환한다(기본 개인·회사). 사용자가 추가한 분류도 포함."""
        cats = load_settings().get("cal_cats")
        if not isinstance(cats, list) or not cats:
            cats = ["개인", "회사"]
        return cats

    def cal_cat_add(self, name):
        """일정 분류를 추가하고 갱신된 목록을 반환한다."""
        name = (name or "").strip()
        cats = self.cal_cats()
        if name and name not in cats:
            cats = cats + [name]
            save_settings({"cal_cats": cats})
        return cats

    def cal_cat_del(self, name):
        """일정 분류를 삭제하고 갱신된 목록을 반환한다(최소 1개는 남긴다)."""
        cats = [c for c in self.cal_cats() if c != name]
        if not cats:
            cats = ["개인"]
        save_settings({"cal_cats": cats})
        return cats

    def open_url(self, url):
        """링크를 시스템 기본 브라우저로 연다."""
        try:
            if url.startswith(("http://", "https://")):
                webbrowser.open(url)
        except Exception as e:
            log(f"링크 열기 실패: {e}")
        return "ok"

    def export_chat(self):
        """현재 대화를 마크다운 파일로 결과물 폴더에 저장한다."""
        try:
            msgs = [m for m in self._messages if m["role"] in ("user", "assistant")]
            if not msgs:
                return "💬 저장할 대화가 없습니다."
            os.makedirs(OUT_DIR, exist_ok=True)
            fname = time.strftime("대화내보내기_%Y%m%d_%H%M%S.md")
            path = os.path.join(OUT_DIR, fname)
            with open(path, "w", encoding="utf-8") as f:
                f.write(f"# {ASSISTANT_NAME} 대화 기록\n\n")
                f.write(f"_내보낸 시각: {time.strftime('%Y-%m-%d %H:%M')}_\n\n---\n\n")
                for m in msgs:
                    who = "🧑 나" if m["role"] == "user" else f"✱ {ASSISTANT_NAME}"
                    f.write(f"**{who}**\n\n{m['content']}\n\n---\n\n")
            subprocess.Popen(["explorer", "/select,", path])
            return f"⬇ 대화를 저장했습니다 → {path}"
        except Exception as e:
            return f"내보내기 실패: {e}"

    def _check_update_worker(self):
        """백그라운드에서 업데이트를 확인하고 새 버전이 있으면 배너를 표시한다."""
        time.sleep(8)  # UI 로딩 완료 대기
        result = check_for_update()
        if result:
            latest, url, changelog = result
            self._js(f"showUpdateBanner({json.dumps(latest,ensure_ascii=False)},"
                     f"{json.dumps(url)},{json.dumps(changelog,ensure_ascii=False)})")

    def _titles(self):
        try:
            with open(TITLES_DB, encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return {}

    def _save_titles(self, d):
        try:
            with open(TITLES_DB, "w", encoding="utf-8") as f:
                json.dump(d, f, ensure_ascii=False)
        except Exception as e:
            log(f"제목 저장 실패: {e}")

    def list_history(self):
        os.makedirs(HIST_DIR, exist_ok=True)
        titles = self._titles()
        out = []
        for p in sorted(globmod.glob(os.path.join(HIST_DIR, "대화_*.json")), reverse=True)[:50]:
            try:
                base = os.path.basename(p)  # 대화_YYYYMMDD_HHMMSS.json
                if base in titles:                      # 사용자가 직접 수정한 제목 우선
                    title = titles[base]
                else:
                    with open(p, encoding="utf-8") as f:
                        msgs = json.load(f)
                    first = next((m["content"] for m in msgs if m["role"] == "user"), "")
                    title = first.split("\n")[0][:42] or "(빈 대화)"
                d = base[3:11]; t = base[12:18]
                date = f"{d[:4]}-{d[4:6]}-{d[6:8]} {t[:2]}:{t[2:4]}"
            except Exception:
                continue
            out.append({"file": base, "title": title, "date": date})
        return out

    def search_history(self, query):
        """과거 대화 의미검색 (/찾기) — 백그라운드 임베딩 후 상위 결과를 JS로 렌더."""
        query = (query or "").strip()
        if not query:
            self._sysmsg("찾을 내용을 적어 주세요. 예: /찾기 견적서 단가")
            return
        self._status("🔎 과거 대화를 의미로 검색하는 중…")
        try:
            res = search_history_files(query)
        except Exception as e:
            log(f"search_history 실패: {e}")
            res = []
        self._status(f"● {self._model}", "ok")
        self._js(f"renderHistSearch({json.dumps(query, ensure_ascii=False)},"
                 f"{json.dumps(res, ensure_ascii=False)})")

    def rename_history(self, fname, title):
        """좌측 대화 목록의 제목을 사용자가 직접 수정 → 대화제목.json에 저장"""
        base = os.path.basename(fname)
        title = (title or "").strip()[:60]
        titles = self._titles()
        if title:
            titles[base] = title
        else:
            titles.pop(base, None)
        self._save_titles(titles)
        return True

    def delete_history(self, fname):
        """채팅 삭제 — UI 확인 창을 거친 뒤에만 호출됨"""
        base = os.path.basename(fname)
        path = os.path.join(HIST_DIR, base)
        try:
            os.remove(path)
        except Exception as e:
            return f"삭제 실패: {e}"
        titles = self._titles()                          # 수정한 제목도 함께 정리
        if titles.pop(base, None) is not None:
            self._save_titles(titles)
        if path == self._session:
            self._messages = [{"role": "system", "content": build_system()}]
            self._session = os.path.join(HIST_DIR, time.strftime("대화_%Y%m%d_%H%M%S.json"))
            return "CURRENT"
        return "OK"

    def load_history(self, fname):
        self._save_history()
        path = os.path.join(HIST_DIR, os.path.basename(fname))
        try:
            with open(path, encoding="utf-8") as f:
                old = json.load(f)
        except Exception as e:
            self._sysmsg(f"대화를 불러오지 못했습니다: {e}")
            return False
        self._messages = [{"role": "system", "content": build_system()}] + old
        self._session = path
        for m in old:
            txt = m["content"].split("\n\n=== ")[0]
            if m["role"] == "user":
                self._js(f"userMsg({json.dumps(txt)})")
            else:
                self._js(f"aiMsg({json.dumps(txt)})")
        self._sysmsg("이 대화를 이어서 계속할 수 있습니다")
        return True

    def new_chat(self):
        self._save_history()
        self._messages = [{"role": "system", "content": build_system()}]
        self._session = os.path.join(HIST_DIR, time.strftime("대화_%Y%m%d_%H%M%S.json"))
        return True

    def pick_files(self):
        paths = self._window.create_file_dialog(
            webview.OPEN_DIALOG, allow_multiple=True,
            directory=os.path.join(HOME, "받은파일"))
        if paths:
            self._attached.extend(paths)
        return [os.path.basename(p) for p in self._attached]

    def organize_folder(self):
        """폴더를 골라 파일을 종류별 하위폴더로 자동 분류한다. 원본은 이동만, 삭제 안 함."""
        try:
            sel = self._window.create_file_dialog(
                webview.FOLDER_DIALOG,
                directory=os.path.join(os.path.expanduser("~"), "Desktop"))
        except Exception as e:
            log(f"organize_folder 다이얼로그 실패: {e}")
            return "⚠ 폴더를 열 수 없습니다."
        if not sel:
            return ""  # 취소
        folder = sel[0] if isinstance(sel, (list, tuple)) else sel
        return self._organize_into(folder)

    def organize_folder_path(self, folder):
        """경로를 직접 받아 그 폴더를 종류별로 정리한다 (즐겨찾기에서 호출)."""
        return self._organize_into(folder)

    def organize_preview(self, folder):
        """폴더를 정리하기 전에 어떤 파일이 어느 폴더로 이동될지 미리 보여준다."""
        if not folder or not os.path.isdir(folder):
            return {"error": "폴더가 아닙니다."}
        CATS = {
            "이미지": {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tif", ".tiff", ".heic", ".svg", ".ico"},
            "문서": {".pdf", ".doc", ".docx", ".txt", ".rtf", ".hwp", ".hwpx", ".md", ".odt"},
            "엑셀": {".xls", ".xlsx", ".csv", ".ods"},
            "프레젠테이션": {".ppt", ".pptx", ".key", ".odp"},
            "도면": {".dwg", ".dxf", ".skp", ".rvt", ".3ds", ".ai", ".psd", ".indd", ".eps"},
            "압축": {".zip", ".rar", ".7z", ".tar", ".gz", ".alz", ".egg"},
            "동영상": {".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm", ".m4v"},
            "음악": {".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma"},
            "실행파일": {".exe", ".msi", ".bat", ".lnk"},
        }
        cat_names = set(CATS.keys()) | {"기타"}
        preview = {}
        try:
            entries = os.listdir(folder)
        except Exception as e:
            return {"error": f"폴더를 읽을 수 없습니다: {e}"}
        for name in entries:
            src = os.path.join(folder, name)
            if not os.path.isfile(src):
                continue
            if name in cat_names:
                continue
            ext = os.path.splitext(name)[1].lower()
            cat = "기타"
            for c, exts in CATS.items():
                if ext in exts:
                    cat = c
                    break
            if cat not in preview:
                preview[cat] = []
            preview[cat].append(name)
        return {"folder": folder, "preview": preview}

    def clean_folders(self):
        """폴더 정리 모달에 보여줄 자주 쓰는 폴더 목록(존재하는 것만)."""
        home = os.path.expanduser("~")
        cand = [
            ("🖥 바탕화면", os.path.join(home, "Desktop")),
            ("⬇ 다운로드", os.path.join(home, "Downloads")),
            ("📄 문서", os.path.join(home, "Documents")),
            ("📥 받은파일", INBOX),
            ("📤 결과물", OUT_DIR),
        ]
        return [{"label": l, "path": p} for l, p in cand if os.path.isdir(p)]

    def _organize_into(self, folder):
        if not folder or not os.path.isdir(folder):
            return "⚠ 폴더가 아닙니다."

        CATS = {
            "이미지": {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tif", ".tiff", ".heic", ".svg", ".ico"},
            "문서": {".pdf", ".doc", ".docx", ".txt", ".rtf", ".hwp", ".hwpx", ".md", ".odt"},
            "엑셀": {".xls", ".xlsx", ".csv", ".ods"},
            "프레젠테이션": {".ppt", ".pptx", ".key", ".odp"},
            "도면": {".dwg", ".dxf", ".skp", ".rvt", ".3ds", ".ai", ".psd", ".indd", ".eps"},
            "압축": {".zip", ".rar", ".7z", ".tar", ".gz", ".alz", ".egg"},
            "동영상": {".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm", ".m4v"},
            "음악": {".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma"},
            "실행파일": {".exe", ".msi", ".bat", ".lnk"},
        }
        cat_names = set(CATS.keys()) | {"기타"}
        moved = {}
        try:
            entries = os.listdir(folder)
        except Exception as e:
            return f"⚠ 폴더를 읽을 수 없습니다: {e}"

        for name in entries:
            src = os.path.join(folder, name)
            if not os.path.isfile(src):
                continue  # 폴더는 건드리지 않음
            ext = os.path.splitext(name)[1].lower()
            cat = "기타"
            for c, exts in CATS.items():
                if ext in exts:
                    cat = c
                    break
            if name in cat_names:
                continue
            dest_dir = os.path.join(folder, cat)
            os.makedirs(dest_dir, exist_ok=True)
            dest = os.path.join(dest_dir, name)
            if os.path.exists(dest):  # 같은 이름이면 _2, _3 …
                base, e2 = os.path.splitext(name)
                k = 2
                while os.path.exists(os.path.join(dest_dir, f"{base}_{k}{e2}")):
                    k += 1
                dest = os.path.join(dest_dir, f"{base}_{k}{e2}")
            try:
                os.replace(src, dest)
                moved[cat] = moved.get(cat, 0) + 1
            except Exception as e:
                log(f"파일 이동 실패 {name}: {e}")

        if not moved:
            return f"📁 정리할 파일이 없었어요 — `{folder}`\n(하위 폴더 안의 파일은 건드리지 않았습니다.)"
        total = sum(moved.values())
        lines = [f"🧹 **폴더 정리 완료** — `{folder}`", "", f"총 **{total}개** 파일을 종류별 폴더로 옮겼어요. 원본은 삭제하지 않았습니다.", ""]
        for c, n in sorted(moved.items(), key=lambda x: -x[1]):
            lines.append(f"- 📂 **{c}** — {n}개")
        return "\n".join(lines)

    def read_clipboard(self):
        """클립보드의 텍스트를 돌려준다 (복사한 내용 빠른 질문용)."""
        try:
            import win32clipboard
            win32clipboard.OpenClipboard()
            try:
                if win32clipboard.IsClipboardFormatAvailable(win32clipboard.CF_UNICODETEXT):
                    return win32clipboard.GetClipboardData(win32clipboard.CF_UNICODETEXT) or ""
            finally:
                win32clipboard.CloseClipboard()
        except Exception as e:
            log(f"read_clipboard 실패: {e}")
        return ""

    # ---------- 워크플로우 (프로젝트·업무 보드) ----------
    def wf_data(self):
        """전체 워크플로우 데이터. 폴더가 연결된 프로젝트는 파일 수도 함께 계산."""
        d = wf_load()
        for p in d.get("projects", []):
            fp = p.get("folder")
            if fp and os.path.isdir(fp):
                try:
                    p["file_count"] = sum(1 for n in os.listdir(fp)
                                          if os.path.isfile(os.path.join(fp, n)))
                except Exception:
                    p["file_count"] = None
        return d

    def wf_add_project(self, name, color=None, folder=None):
        d = wf_load()
        pid = int(time.time() * 1000)
        d["projects"].append({"id": pid, "name": (name or "새 프로젝트").strip(),
                              "color": color or "#5B6CFF", "folder": folder or "",
                              "tasks": []})
        return wf_save(d)

    def wf_set_project(self, pid, fields):
        d = wf_load()
        for p in d["projects"]:
            if p["id"] == pid:
                for k in ("name", "color", "folder"):
                    if k in fields:
                        p[k] = fields[k]
        return wf_save(d)

    def wf_del_project(self, pid):
        d = wf_load()
        d["projects"] = [p for p in d["projects"] if p["id"] != pid]
        return wf_save(d)

    def wf_pick_folder(self):
        """프로젝트에 연결할 폴더를 고른다."""
        try:
            sel = self._window.create_file_dialog(
                webview.FOLDER_DIALOG,
                directory=os.path.join(os.path.expanduser("~"), "Desktop"))
            if sel:
                return sel[0] if isinstance(sel, (list, tuple)) else sel
        except Exception as e:
            log(f"wf_pick_folder 실패: {e}")
        return ""

    def wf_add_task(self, pid, title, priority="보통", due=""):
        d = wf_load()
        for p in d["projects"]:
            if p["id"] == pid:
                p.setdefault("tasks", []).append({
                    "id": int(time.time() * 1000),
                    "title": (title or "새 업무").strip(),
                    "status": "todo", "priority": priority or "보통",
                    "due": due or "", "memo": ""})
        return wf_save(d)

    def wf_set_task(self, pid, tid, fields):
        d = wf_load()
        for p in d["projects"]:
            if p["id"] == pid:
                for t in p.get("tasks", []):
                    if t["id"] == tid:
                        for k in ("title", "status", "priority", "due", "memo"):
                            if k in fields:
                                t[k] = fields[k]
        return wf_save(d)

    def wf_del_task(self, pid, tid):
        d = wf_load()
        for p in d["projects"]:
            if p["id"] == pid:
                p["tasks"] = [t for t in p.get("tasks", []) if t["id"] != tid]
        return wf_save(d)

    def fav_folders(self):
        """등록된 즐겨찾기 폴더 목록을 돌려준다."""
        return load_settings().get("fav_folders", [])

    def add_fav_folder(self):
        """폴더를 골라 즐겨찾기에 추가한다."""
        try:
            sel = self._window.create_file_dialog(
                webview.FOLDER_DIALOG,
                directory=os.path.join(os.path.expanduser("~"), "Desktop"))
        except Exception as e:
            log(f"add_fav_folder 실패: {e}")
            return self.fav_folders()
        if not sel:
            return self.fav_folders()
        folder = sel[0] if isinstance(sel, (list, tuple)) else sel
        favs = load_settings().get("fav_folders", [])
        if folder not in favs:
            favs.append(folder)
            save_settings({"fav_folders": favs})
        return favs

    def del_fav_folder(self, path):
        favs = [f for f in load_settings().get("fav_folders", []) if f != path]
        save_settings({"fav_folders": favs})
        return favs

    def open_folder_path(self, path):
        """탐색기로 폴더를 연다."""
        try:
            if os.path.isdir(path):
                os.startfile(path)
                return "ok"
        except Exception as e:
            log(f"open_folder_path 실패: {e}")
        return "fail"

    def save_tts_rate(self, rate):
        save_settings({"tts_rate": float(rate)})
        return "ok"

    def start_voice(self):
        """마이크 음성을 받아 한국어 텍스트로 변환해 돌려준다 (faster-whisper, 로컬)."""
        self._status("🎤 듣고 있습니다… 말씀하세요", "")
        try:
            text = transcribe_mic()
        except Exception as e:
            log(f"start_voice 실패: {e}")
            text = ""
        self._status(f"● {self._model}", "ok")
        if not text:
            self._sysmsg("🎤 음성을 알아듣지 못했습니다 — 마이크를 확인하고 다시 시도해 주세요")
        return text

    def capture_screen(self):
        """화면 캡처 — 윈도우 캡처 도구를 띄우고 결과를 자동 첨부"""
        def worker():
            try:
                from PIL import ImageGrab
                try:
                    import ctypes
                    u = ctypes.windll.user32
                    if u.OpenClipboard(0):
                        u.EmptyClipboard()
                        u.CloseClipboard()
                except Exception:
                    pass
                try:
                    self._window.minimize()
                except Exception:
                    pass
                subprocess.Popen("explorer ms-screenclip:", shell=True)
                img = None
                for _ in range(120):  # 최대 60초 대기
                    time.sleep(0.5)
                    try:
                        g = ImageGrab.grabclipboard()
                    except Exception:
                        g = None
                    if g is not None and not isinstance(g, list):
                        img = g
                        break
                try:
                    self._window.restore()
                except Exception:
                    pass
                if img is None:
                    self._sysmsg("캡처가 취소되었거나 시간이 초과되었습니다")
                    return
                d = os.path.join(HOME, "받은파일")
                os.makedirs(d, exist_ok=True)
                p = os.path.join(d, time.strftime("캡처_%Y%m%d_%H%M%S.png"))
                img.save(p)
                self._attached.append(p)
                names = json.dumps([os.path.basename(x) for x in self._attached],
                                   ensure_ascii=False)
                self._js(f"showAttached({names})")
                self._sysmsg("📷 캡처 완료 — 질문과 함께 전송하세요")
            except Exception as e:
                try:
                    self._window.restore()
                except Exception:
                    pass
                self._sysmsg(f"캡처 오류: {e}")
        threading.Thread(target=worker, daemon=True).start()
        return True

    def get_attached_paths(self):
        """현재 첨부 파일의 전체 경로 목록 반환 (doSend 호출 전 JS에서 가져감)"""
        return list(self._attached)

    def attach_thumbs(self):
        """첨부 파일들의 작은 미리보기(이미지·DXF 도면)를 data URL 목록으로 반환.
        _attached 순서와 1:1 정렬. 미리보기 불가하면 빈 문자열."""
        if not hasattr(self, "_thumb_cache"):
            self._thumb_cache = {}
        out = []
        for p in list(self._attached):
            key = p + "|" + str(os.path.getmtime(p) if os.path.exists(p) else 0)
            if key in self._thumb_cache:
                out.append(self._thumb_cache[key])
                continue
            url = ""
            try:
                ext = os.path.splitext(p)[1].lower()
                if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"):
                    url = self._img_thumb(p)
                elif ext in (".dxf", ".dwg"):
                    url = self._dxf_thumb(p)
            except Exception as e:
                log(f"썸네일 실패 {os.path.basename(p)}: {e}")
            self._thumb_cache[key] = url
            out.append(url)
        return out

    def _img_thumb(self, path):
        from PIL import Image
        import io
        im = Image.open(path)
        im = im.convert("RGB")
        im.thumbnail((96, 96))
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=70)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()

    def _dxf_thumb(self, path):
        """DXF 도면을 작은 PNG 미리보기로 렌더(베스트에포트)."""
        if path.lower().endswith(".dwg"):
            return ""   # DWG는 직접 렌더 불가
        import ezdxf
        from ezdxf.addons.drawing import RenderContext, Frontend
        from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import io
        doc = ezdxf.readfile(path)
        msp = doc.modelspace()
        fig = plt.figure(figsize=(1.4, 1.4), dpi=70)
        ax = fig.add_axes([0, 0, 1, 1])
        ax.set_axis_off()
        try:
            Frontend(RenderContext(doc), MatplotlibBackend(ax)).draw_layout(msp, finalize=True)
        except Exception:
            plt.close(fig)
            return ""
        buf = io.BytesIO()
        fig.savefig(buf, format="png", facecolor="white",
                    bbox_inches="tight", pad_inches=0.05)
        plt.close(fig)
        return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()

    def reattach_files(self, paths):
        """경로 목록을 _attached에 다시 추가 (파일 재첨부 버튼)"""
        for p in paths:
            if os.path.exists(p) and p not in self._attached:
                self._attached.append(p)
        return [os.path.basename(x) for x in self._attached]

    def remove_attached(self, idx):
        """첨부 목록에서 빼기 (✕ 버튼) — 파일 자체는 삭제하지 않음"""
        try:
            self._attached.pop(int(idx))
        except Exception:
            pass
        return [os.path.basename(x) for x in self._attached]

    def attach_clipboard_image(self, data_url):
        """채팅창에 Ctrl+V로 붙여넣은 이미지를 첨부"""
        try:
            b64 = data_url.split(",", 1)[1]
            d = os.path.join(HOME, "받은파일")
            os.makedirs(d, exist_ok=True)
            p = os.path.join(d, time.strftime("붙여넣기_%Y%m%d_%H%M%S.png"))
            with open(p, "wb") as f:
                f.write(base64.b64decode(b64))
            self._attached.append(p)
        except Exception as e:
            self._sysmsg(f"이미지 붙여넣기 실패: {e}")
        return [os.path.basename(x) for x in self._attached]

    def vault(self, text, master):
        parts = text.split()
        action = parts[1] if len(parts) > 1 else ""
        data = vault_load(master)
        if data is None:
            return "WRONGPW"
        try:
            if action == "저장" and len(parts) >= 5:
                name, uid, pw = parts[2], parts[3], " ".join(parts[4:])
                data[name] = {"id": uid, "pw": pw, "saved": time.strftime("%Y-%m-%d")}
                vault_save(master, data)
                return f"🔐 '{name}' 저장 완료 (암호화됨). 조회: /금고 보기 {name}"
            if action == "보기" and len(parts) >= 3:
                e = data.get(parts[2])
                if not e:
                    return f"'{parts[2]}' 항목이 없습니다. 목록: /금고 목록"
                return f"🔐 {parts[2]} — 아이디: {e['id']} / 비밀번호: {e['pw']} (저장일 {e.get('saved','?')})"
            if action == "목록":
                return "🔐 저장된 항목: " + (", ".join(data.keys()) if data else "(없음)")
            if action == "삭제" and len(parts) >= 3:
                if data.pop(parts[2], None) is not None:
                    vault_save(master, data)
                    return f"'{parts[2]}' 삭제 완료"
                return f"'{parts[2]}' 항목이 없습니다"
            return "사용법: /금고 저장 이름 아이디 비밀번호 · /금고 보기 이름 · /금고 목록 · /금고 삭제 이름"
        except Exception as e:
            return f"금고 오류: {e}"

    # ---------- 일정 · 메모 ----------
    def sched_data(self):
        return sched_load()

    def sched_add(self, ev):
        with SCHED_LOCK:
            d = sched_load()
            d["events"].append({"id": int(time.time() * 1000),
                                "date": ev["date"], "time": ev.get("time") or "",
                                "title": ev["title"], "memo": ev.get("memo", ""),
                                "category": ev.get("category", "개인"),
                                "remind": [int(x) for x in ev.get("remind", [])],
                                "done": False, "notified": []})
            sched_save(d)
        return d

    def sched_summary(self, category=None):
        """일정 정리/요약 반환 (category='개인'/'회사'/None=전체)"""
        with SCHED_LOCK:
            d = sched_load()
        evs = sorted([e for e in d["events"] if not e.get("done")],
                     key=lambda e: (e["date"], e.get("time") or "99"))
        if category and category != "전체":
            evs = [e for e in evs if e.get("category", "개인") == category]
        today = time.strftime("%Y-%m-%d")
        past   = [e for e in evs if e["date"] < today]
        coming = [e for e in evs if e["date"] >= today]
        lines = []
        if coming:
            lines.append("**📅 다가오는 일정**")
            for e in coming[:30]:
                icon = "🏢" if e.get("category") == "회사" else "👤"
                lines.append(
                    f"- {icon} **{e['date']}({kday(e['date'])})** "
                    f"{e.get('time') or '종일'}  {e['title']}")
                if e.get("memo"):
                    lines.append(f"  → {e['memo']}")
        if past:
            lines.append("\n**📋 지난 일정 (최근 10건)**")
            for e in past[-10:]:
                icon = "🏢" if e.get("category") == "회사" else "👤"
                lines.append(
                    f"- ~~{icon} {e['date']} {e.get('time') or '종일'}  {e['title']}~~")
        return ("\n".join(lines) if lines
                else "일정이 없습니다" + (f" ({category})" if category else ""))

    def sched_set(self, eid, fields):
        with SCHED_LOCK:
            d = sched_load()
            for e in d["events"]:
                if e["id"] == eid:
                    e.update(fields)
            sched_save(d)
        return d

    def sched_del(self, eid):
        with SCHED_LOCK:
            d = sched_load()
            d["events"] = [e for e in d["events"] if e["id"] != eid]
            sched_save(d)
        return d

    def memo_add(self, text):
        text = (text or "").strip()
        with SCHED_LOCK:
            d = sched_load()
            if text:
                d["memos"].insert(0, {"id": int(time.time() * 1000),
                                      "ts": time.strftime("%Y-%m-%d %H:%M"),
                                      "text": text})
                sched_save(d)
        return d

    def memo_del(self, mid):
        with SCHED_LOCK:
            d = sched_load()
            d["memos"] = [m for m in d["memos"] if m["id"] != mid]
            sched_save(d)
        return d

    def _reminder_worker(self):
        """30초마다 일정 확인 — '하루 전·30분 전' 등 설정된 시점에 윈도우 토스트 + 채팅 알림"""
        while True:
            try:
                now = time.time()
                with SCHED_LOCK:
                    d = sched_load()
                    changed = False
                    for ev in d["events"]:
                        if ev.get("done"):
                            continue
                        et = ev_epoch(ev)
                        if et is None:
                            continue
                        for off in ev.get("remind", []):
                            fire = et - int(off) * 60
                            if fire > now or off in ev.get("notified", []):
                                continue
                            ev.setdefault("notified", []).append(off)
                            changed = True
                            if now - fire < 6 * 3600:   # 너무 오래 지난 알림은 표시 생략
                                when = (f"{ev['date']}({kday(ev['date'])}) "
                                        f"{ev.get('time') or '종일'}")
                                toast(f"⏰ 일정 알림 — {remind_label(off)}",
                                      f"{when} · {ev['title']}")
                                self._sysmsg(f"⏰ {remind_label(off)} 알림: {when} · {ev['title']}")
                    if changed:
                        sched_save(d)
            except Exception:
                pass
            time.sleep(30)

    def send(self, text):
        files, self._attached = self._attached, []
        if getattr(self, "_busy", False):   # 답변 중이면 큐에 쌓고 순서대로 처리
            self._queue = getattr(self, "_queue", [])
            self._queue.append((text, files))
            self._sysmsg("⏳ 접수했습니다 — 앞 작업이 끝나는 대로 바로 처리합니다")
            return
        self._busy = True
        threading.Thread(target=self._reply, args=(text, files), daemon=True).start()

    def stop(self):
        """생성 중단 — 진행 중인 모델 응답을 멈춘다 (Stop 버튼)."""
        self._stop = True
        return True

    def _inbox_worker(self):
        """받은파일 폴더 감시 — 새 파일이 들어오면 토스트 + 채팅 알림으로 처리 제안."""
        os.makedirs(INBOX, exist_ok=True)
        try:
            seen = set(os.listdir(INBOX))
        except Exception:
            seen = set()
        while True:
            try:
                time.sleep(4)
                now = set(os.listdir(INBOX))
                new = [f for f in (now - seen)
                       if os.path.isfile(os.path.join(INBOX, f))
                       and not f.startswith("~$")]
                seen = now
                for f in new:
                    path = os.path.join(INBOX, f)
                    try:   # 복사가 끝날 때까지(크기 안정) 대기 후 알림
                        s1 = os.path.getsize(path)
                        time.sleep(1.2)
                        if os.path.getsize(path) != s1:
                            seen.discard(f)     # 아직 복사 중 → 다음 루프에서 다시 감지
                            continue
                    except OSError:
                        continue
                    try:
                        toast("AI 비서 · 받은파일", f"새 파일: {f}")
                    except Exception:
                        pass
                    self._sysmsg(f"📥 받은파일 폴더에 **{f}** 이(가) 들어왔습니다 — "
                                 "“이 파일 분석해줘” 또는 원하는 처리를 말씀해 주세요.")
            except Exception as e:
                log(f"inbox_worker 오류: {e}")
                time.sleep(10)

    def _next_in_queue(self):
        q = getattr(self, "_queue", [])
        if q:
            text, files = q.pop(0)
            self._busy = True
            threading.Thread(target=self._reply, args=(text, files), daemon=True).start()

    # ---------- 코드 자동 실행 (프리패스) ----------
    def _exec_code(self, code):
        """모델이 작성한 파이썬 코드를 즉시 실행 — 승인 절차 없음(프리패스 설정).
        단, 되돌릴 수 없는 치명적 명령(포맷·파티션·시스템 전체 삭제)은 차단한다.
        반환: (ok: bool, output: str) — 실패 시 output을 모델에 돌려줘 스스로 고치게 한다."""
        try:   # 모든 실행을 감사 로그로 남김 (나중에 추적 가능)
            with open(EXEC_LOG, "a", encoding="utf-8") as f:
                f.write(f"\n===== {time.strftime('%Y-%m-%d %H:%M:%S')} =====\n{code}\n")
        except Exception:
            pass
        if HARD_BLOCK.search(code):
            self._sysmsg("🛑 디스크 포맷·파티션·시스템 폴더 전체 삭제처럼 되돌릴 수 없는 명령이 "
                         "감지되어 **실행을 차단**했습니다. 정말 필요하면 직접 실행해 주세요.")
            log(f"HARD_BLOCK 차단:\n{code}")
            return False, "차단됨(되돌릴 수 없는 위험 명령) — 실행하지 않음"
        if DELETE_PATTERNS.search(code):
            self._sysmsg("⚠️ 파일 삭제 명령이 포함된 코드입니다 — 프리패스 설정에 따라 그대로 실행합니다 "
                         "(실행로그.txt에 기록됨)")
        self._run_count += 1
        os.makedirs(OUT_DIR, exist_ok=True)
        script = os.path.join(OUT_DIR, f"_ai_script_{int(time.time())}_{self._run_count}.py")
        with open(script, "w", encoding="utf-8") as f:
            f.write(code)
        self._sysmsg("⚙ 코드를 자동 실행합니다…")
        try:
            r = subprocess.run(["python", script], capture_output=True, text=True,
                               encoding="utf-8", errors="replace", timeout=300,
                               cwd=OUT_DIR, creationflags=0x08000000)
            out = ((r.stdout or "") + ("\n" + r.stderr if r.stderr else "")).strip()
            if r.returncode == 0:
                self._sysmsg("✅ 실행 완료" + (f" — {out[:600]}" if out else "") +
                             "\n📂 결과물 폴더가 열립니다")
                subprocess.Popen(["explorer", OUT_DIR])
                return True, out[:1500]
            self._sysmsg(f"❌ 실행 실패 — 스스로 고쳐 다시 시도합니다: {out[:400]}")
            return False, out[:1500]
        except subprocess.TimeoutExpired:
            self._sysmsg("❌ 실행 시간 초과(5분) — 중단했습니다")
            return False, "TimeoutExpired: 5분 초과로 강제 종료됨"
        except Exception as e:
            self._sysmsg(f"❌ 실행 오류: {e}")
            return False, f"{type(e).__name__}: {e}"

    # ---------- 내부 ----------
    def _ping(self):
        try:
            with urllib.request.urlopen(OLLAMA + "/api/tags", timeout=2):
                return True
        except Exception:
            return False

    def _pick_model(self):
        """설치된 모델 중 텍스트용/이미지용을 각각 선택. 없으면 호환 모델 자동 설치."""
        try:
            with urllib.request.urlopen(OLLAMA + "/api/tags", timeout=5) as r:
                names = [m["name"] for m in json.load(r)["models"]]
        except Exception:
            names = []

        def pick(prefs):
            for pref in prefs:
                if any(n == pref or n.startswith(pref + "-") for n in names):
                    return pref
            for pref in prefs:
                base = pref.split(":")[0]
                for n in names:
                    if n.startswith(base):
                        return n
            return None

        if not names:
            self._status("AI 모델 내려받는 중… (최초 1회)", "err")
            self._sysmsg(f"이 컴퓨터에 맞는 AI 모델({MODEL_SAFE_PULL})을 처음 한 번 내려받습니다. 잠시만요…")
            subprocess.run(["ollama", "pull", MODEL_SAFE_PULL], creationflags=0x08000000)
            names = [MODEL_SAFE_PULL]
        self._model_vision = pick(VISION_PREFS) or names[0]
        return pick(TEXT_PREFS) or names[0]

    def _self_diagnose(self):
        """시작 시 핵심 라이브러리·환경을 점검하고 문제가 있으면 한 번에 안내한다."""
        libs = {"openpyxl": "엑셀", "docx": "워드", "pptx": "PPT", "pypdf": "PDF",
                "PIL": "이미지", "ezdxf": "캐드도면", "win32com": "프로그램자동화"}
        missing = []
        for mod, label in libs.items():
            try:
                __import__(mod)
            except Exception:
                missing.append((mod, label))
        if missing:
            names = ", ".join(f"{lbl}({m})" for m, lbl in missing)
            pkgs = " ".join({"docx": "python-docx", "pptx": "python-pptx",
                             "PIL": "Pillow", "win32com": "pywin32"}.get(m, m)
                            for m, _ in missing)
            self._sysmsg(f"🩺 일부 기능 라이브러리가 없습니다: {names}\n"
                         f"해당 작업이 필요하면 명령창에서 설치하세요:  pip install {pkgs}")
        return not missing

    def _init(self):
        # 일정 알림은 AI 엔진과 무관하게 항상 동작
        threading.Thread(target=self._reminder_worker, daemon=True).start()
        threading.Thread(target=self._self_diagnose, daemon=True).start()  # 라이브러리 자가진단
        threading.Thread(target=self._ensure_shortcut, daemon=True).start()  # 최초 1회 바로가기 자동 생성
        try:   # 오늘 일정 브리핑 — 켜자마자 한눈에
            today = time.strftime("%Y-%m-%d")
            evs = sorted([e for e in sched_load()["events"]
                          if e["date"] == today and not e.get("done")],
                         key=lambda e: e.get("time") or "99")
            if evs:
                self._sysmsg("📅 오늘 일정 " + str(len(evs)) + "건 — "
                             + " · ".join((e.get("time") or "종일") + " " + e["title"]
                                          for e in evs))
        except Exception:
            pass
        try:   # 워크플로우 브리핑 — 진행중·높은 중요도 업무를 켜자마자 보여준다
            wf = wf_load()
            doing, high = [], []
            for p in wf.get("projects", []):
                for t in p.get("tasks", []):
                    if t.get("status") == "done":
                        continue
                    if t.get("status") == "doing":
                        doing.append(f"{p['name']}: {t['title']}")
                    if t.get("priority") == "높음":
                        high.append(f"{p['name']}: {t['title']}")
            parts = []
            if doing:
                parts.append("진행중 " + str(len(doing)) + "건 — " + " · ".join(doing[:3]))
            if high:
                parts.append("🔴 중요 " + str(len(high)) + "건 — " + " · ".join(high[:3]))
            if parts:
                self._sysmsg("📋 워크플로우 — " + " / ".join(parts)
                             + "  (왼쪽 📋 워크플로우에서 전체 보기)")
        except Exception:
            pass
        if not self._ping():
            exe = os.path.expandvars(r"%LOCALAPPDATA%\Programs\Ollama\ollama.exe")
            installed = os.path.exists(exe) or shutil.which("ollama")
            if not installed:
                # 최초 사용자: Ollama가 아예 안 깔린 경우 — 재시작이 아니라 '설치'가 필요하다.
                self._status("AI 엔진(Ollama) 설치 필요", "err")
                self._sysmsg("⚙️ 이 비서는 무료 AI 엔진 ‘Ollama’가 있어야 작동합니다.\n"
                             "설치 페이지를 열어드릴게요. 설치를 마친 뒤 이 창을 다시 켜면 자동으로 연결됩니다.\n"
                             "👉 https://ollama.com/download")
                try:
                    webbrowser.open("https://ollama.com/download")
                except Exception:
                    pass
                return
            self._status("AI 엔진 시작 중…")
            try:
                subprocess.Popen([exe if os.path.exists(exe) else "ollama", "serve"],
                                 creationflags=0x08000000)
            except Exception as e:
                self._sysmsg(f"Ollama 실행 실패: {e}")
            for _ in range(40):
                time.sleep(1)
                if self._ping():
                    break
        if not self._ping():
            self._status("엔진 연결 실패", "err")
            self._sysmsg("AI 엔진(Ollama) 연결에 실패했습니다. 컴퓨터를 재시작한 뒤 다시 열어 주세요.\n"
                         "그래도 안 되면 작업표시줄 오른쪽 아래에 Ollama 아이콘(라마 모양)이 떠 있는지 확인해 주세요.")
            return
        pruned = prune_history()
        if pruned:
            self._sysmsg(f"🧹 오래된 채팅기록 {pruned}건을 자동 정리했습니다 (문서·파일은 건드리지 않음)")
        self._model = self._pick_model()
        # 저장된 글씨 크기 · 빠른 버튼 적용
        try:
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
                if s.get("font"):
                    self._js(f"applyFS({int(s['font'])},false)")
                if s.get("buttons"):
                    self._js(f"setQButtons({json.dumps(s['buttons'], ensure_ascii=False)})")
                if s.get("answer_style") in STYLE_HINT:
                    global ANSWER_STYLE
                    ANSWER_STYLE = s["answer_style"]
                    self._messages[0] = {"role": "system", "content": build_system()}
                if s.get("name"):
                    global ASSISTANT_NAME
                    ASSISTANT_NAME = s["name"]
                    self._messages[0] = {"role": "system", "content": build_system()}
                    self._js(f"document.querySelector('.brand').innerHTML="
                             f"'<span class=\"star\">✱</span> '+{json.dumps(ASSISTANT_NAME)}")
                if s.get("dark"):
                    self._js("applyTheme(true)")
                else:
                    self.set_titlebar_color(False)
        except Exception:
            pass
        # 업무 위임 기본 버튼 — 저장된 버튼이 하나도 없으면 코워크 단축버튼을 깔아준다
        try:
            ps0 = load_settings()
            if not ps0.get("buttons"):
                defaults = [
                    {"name": "받은파일 처리", "cmd": "받은파일 폴더에 있는 파일을 확인하고 정리·요약해줘"},
                    {"name": "파일 요약", "cmd": "방금 첨부한 파일 내용을 읽고 핵심만 보기좋게 정리해줘"},
                    {"name": "표로 정리", "cmd": "방금 내용을 표로 보기좋게 정리해줘"},
                    {"name": "검색해서 정리", "cmd": "이걸 인터넷에서 찾아 근거와 함께 정리해줘: "},
                ]
                save_settings({"buttons": defaults})
                self._js(f"setQButtons({json.dumps(defaults, ensure_ascii=False)})")
        except Exception:
            pass
        # 피카츄 캐릭터 — 메인 창 밖 별도 위젯 창. 설정값을 메인 창의 토글 상태에 반영
        # (실제 창 표시/숨김은 main() 에서 이미 처리됨).
        try:
            ps = load_settings()
            self._js(f"window._pikaOn={json.dumps(bool(ps.get('pika', True)))}")
        except Exception:
            pass
        self._js(f"window._modelName={json.dumps(self._model or '')}")  # 응답 메타 표시용
        self._js(f"window._version={json.dumps(VERSION)}")  # 업데이트 확인 버튼용
        try:   # 설정값을 JS로 전달 (답변 스타일·음성 속도)
            _s2 = load_settings()
            self._js(f"window._answerStyle={json.dumps(_s2.get('answer_style', ANSWER_STYLE))}")
            if _s2.get("tts_rate"):
                self._js(f"window._ttsRate={float(_s2['tts_rate'])};"
                         f"if(typeof _setRate!=='undefined')_setRate={float(_s2['tts_rate'])}")
            if _s2.get("wf_rail"):   # 상시 워크플로우 패널 — 켜져 있으면 부팅 시 표시·채움
                self._js("if(typeof showRail!=='undefined'){showRail(true);"
                         "if(window._wfRefresh)window._wfRefresh();}")
        except Exception:
            pass
        # 모델 예열 — 실제 시스템 프롬프트까지 미리 처리해 프리픽스 캐시를 채운다
        # (첫 질문도 곧바로 응답 시작). keep_alive 60분 동안 메모리에 상주.
        def warmup():
            try:
                self._sys_date = time.strftime("%Y-%m-%d")
                urllib.request.urlopen(urllib.request.Request(
                    OLLAMA + "/api/chat",
                    data=json.dumps({
                        "model": self._model,
                        "messages": [self._messages[0],
                                     {"role": "user", "content": "준비됐어?"}],
                        "keep_alive": "60m", "stream": False,
                        "options": {"num_ctx": NUM_CTX, "num_predict": 1}}).encode(),
                    headers={"Content-Type": "application/json"}), timeout=180).read()
            except Exception as e:
                log(f"warmup 실패: {e}")
        threading.Thread(target=warmup, daemon=True).start()
        threading.Thread(target=self._index_worker, daemon=True).start()
        threading.Thread(target=self._inbox_worker, daemon=True).start()  # 받은파일 자동 감시
        threading.Thread(target=self._check_update_worker, daemon=True).start()  # 업데이트 확인
        self._status(f"● {self._model}", "ok")
        self._js("loadHistList()")
        self._restore_last()

    def _index_worker(self):
        """컴퓨터 전체 색인 — 처음 실행이면 전체 분석, 이후엔 24시간마다 자동 갱신"""
        try:
            s = {}
            if os.path.exists(SETTINGS):
                with open(SETTINGS, encoding="utf-8") as f:
                    s = json.load(f)
            first = not os.path.exists(INDEX_DB)
            if not first and time.time() - s.get("last_index", 0) < 86400:
                return
            if first:
                self._sysmsg("🗂 처음 실행 — 이 컴퓨터 전체를 분석해 파일 색인을 만듭니다. "
                             "백그라운드에서 진행되니 그냥 사용하셔도 됩니다.")
            n = build_index(status=lambda t: self._status(t))
            try:
                if os.path.exists(SETTINGS):
                    with open(SETTINGS, encoding="utf-8") as f:
                        s = json.load(f)
                s["last_index"] = time.time()
                with open(SETTINGS, "w", encoding="utf-8") as f:
                    json.dump(s, f, ensure_ascii=False)
            except Exception:
                pass
            self._messages[0] = {"role": "system", "content": build_system()}  # 분석 요약 반영
            self._status(f"● {self._model}", "ok")
            if first:
                self._sysmsg(f"🗂 컴퓨터 분석 완료 — 파일 {n:,}개 색인. "
                             "이제 경로 없이 '○○ 파일 찾아줘'라고만 하셔도 됩니다.")
        except Exception as e:
            self._sysmsg(f"색인 생성 실패: {e}")
            self._status(f"● {self._model}", "ok")

    def _context_messages(self, has_images):
        """시스템 + (오래된 대화 요약) + 최근 대화만 보내 응답 속도 확보 (기록 파일엔 전체 보존).
        오래된 메시지는 앞부분만 잘라 보냄 — 빠르고, 옛 내용이 새 답에 섞이지 않음."""
        recent = self._messages[1 + self._summary_n:]   # 요약에 접힌 부분은 제외
        if len(recent) > CTX_MSGS:
            recent = recent[-CTX_MSGS:]
        out = []
        for i, m in enumerate(recent):
            mm = dict(m) if has_images else {k: v for k, v in m.items() if k != "images"}
            c = mm.get("content", "")
            if i < len(recent) - 2 and len(c) > 700:   # 마지막 2개 빼고는 요점만
                mm["content"] = c[:700] + " …(이전 대화 — 생략됨)"
            out.append(mm)
        head = [self._messages[0]]
        if self._summary:
            head.append({"role": "user",
                         "content": "=== 지금까지 대화 요약(이전 맥락, 참고만) ===\n"
                                    + self._summary})
        return head + out

    def _maybe_compact(self):
        """대화가 길어지면 오래된 메시지를 한 묶음으로 요약해 컨텍스트를 압축한다.
        self._messages(전체 기록)는 건드리지 않고 요약 + 접힌 개수만 갱신한다."""
        KEEP = 14
        unsummarized = len(self._messages) - 1 - self._summary_n
        if unsummarized <= 28:
            return
        end = len(self._messages) - KEEP          # 최근 KEEP개는 남겨 둔다
        block = self._messages[1 + self._summary_n:end]
        if len(block) < 6:
            return
        convo = []
        for m in block:
            who = "나" if m.get("role") == "user" else "비서"
            c = re.sub(r"\s+", " ", re.sub(r"```[\s\S]*?```", " ",
                                           m.get("content", ""))).strip()[:400]
            if c:
                convo.append(f"{who}: {c}")
        text = "\n".join(convo)[:6000]
        if not text:
            self._summary_n = end - 1
            return
        self._status("🗜 이전 대화를 요약해 정리하는 중…")
        try:
            prior = ("이전 요약:\n" + self._summary + "\n\n") if self._summary else ""
            r = urllib.request.urlopen(urllib.request.Request(
                OLLAMA + "/api/chat",
                data=json.dumps({
                    "model": self._model, "stream": False, "keep_alive": "60m",
                    "options": {"num_ctx": NUM_CTX, "num_predict": 320, "temperature": 0.3},
                    "messages": [
                        {"role": "system", "content":
                            "다음 대화를 한국어 불릿 5~8개로 압축 요약하라. 사용자가 부탁한 일·결정·"
                            "중요한 사실·진행상황만 남기고 인사·잡담은 버려라. 기존 요약이 있으면 합쳐라."},
                        {"role": "user", "content": prior + "새 대화:\n" + text}]}).encode(),
                headers={"Content-Type": "application/json"}), timeout=120).read()
            summary = json.loads(r).get("message", {}).get("content", "").strip()
        except Exception as e:
            log(f"대화 압축 실패: {e}")
            return
        if summary:
            self._summary = summary[:2000]
            self._summary_n = end - 1
            self._sysmsg("🗜 길어진 대화를 요약해 정리했습니다 (이전 맥락은 요약으로 유지)")

    def _is_smalltalk(self, text):
        """단순 인사·잡담인가? (6GB GPU라 모델은 안 바꾸고, 출력만 짧게 잡아 더 빨리 끝낸다)
        도구·정보가 필요한 신호가 하나라도 있으면 잡담이 아니다."""
        t = (text or "").strip()
        if not t or len(t) > 30:
            return False
        if re.search(r"열어|찾아|검색|만들|작성|실행|코드|엑셀|워드|pdf|도면|일정|메모|"
                     r"얼마|어디|언제|왜|어떻게|추천|예약|분석|정리|요약|알려|\?|？", t, re.I):
            return False
        return bool(re.match(
            r"^(안녕|반가|고마|감사|하이|헬로|hi|hello|잘\s?지|뭐\s?해|뭐하|"
            r"심심|졸려|배고|좋아|싫어|ㅋ+|ㅎ+|넵|응|그래|오케이|ok|굿|good)", t, re.I))

    def _call_model(self, model, on_token, has_images=False):
        req = urllib.request.Request(
            OLLAMA + "/api/chat",
            data=json.dumps({"model": model,
                             "messages": self._context_messages(has_images),
                             "keep_alive": "60m",
                             "options": {"num_ctx": NUM_CTX,
                                         "num_predict": self._num_predict,  # 잡담은 짧게 = 빠름
                                         "temperature": 0.6},
                             "stream": True}).encode(),
            headers={"Content-Type": "application/json"})
        full = []
        buf = []
        last_flush = time.time()

        def flush():
            nonlocal buf, last_flush
            if buf:
                on_token("".join(buf))
                buf = []
            last_flush = time.time()

        with urllib.request.urlopen(req, timeout=900) as r:
            for line in r:
                if self._stop:                 # Stop 버튼 → 즉시 생성 중단
                    full.append(" …(중단됨)")
                    break
                chunk = json.loads(line)
                if "error" in chunk:
                    raise RuntimeError(chunk["error"])
                piece = chunk.get("message", {}).get("content", "")
                if piece:
                    full.append(piece)
                    buf.append(piece)
                    # 더 자주 내보내 글자가 빨리 보이게 (체감 속도↑)
                    if time.time() - last_flush > 0.03 or sum(len(b) for b in buf) > 10:
                        flush()
                if chunk.get("done"):
                    break
        flush()
        return "".join(full)

    def _reply(self, text, files):
        self._stop = False          # 새 턴 시작 → 중단 플래그 초기화
        try:
            self._maybe_compact()   # 대화가 길어졌으면 오래된 부분을 요약해 압축
        except Exception as e:
            log(f"_maybe_compact 오류: {e}")
        # ⚡ 속도 핵심: 시스템 메시지(아주 긴 프롬프트)는 날짜가 바뀔 때만 새로 만든다.
        # 매 턴 새로 만들면 Ollama 프리픽스 캐시가 깨져 첫 글자까지 오래 걸린다.
        # 회상한 기억처럼 매번 달라지는 건 시스템이 아니라 사용자 메시지에 붙인다.
        try:
            today = time.strftime("%Y-%m-%d")
            if getattr(self, "_sys_date", None) != today:
                self._messages[0] = {"role": "system", "content": build_system()}
                self._sys_date = today
        except Exception:
            pass
        recalled = []
        try:   # 질문과 관련된 장기기억을 의미검색으로 회상 (사용자 메시지에 첨부 → 캐시 유지)
            recalled = memory_recall(text or "") if text else []
        except Exception as e:
            log(f"memory_recall 실패: {e}")
        images, ctx = [], []
        for p in files:
            summary, img = analyze_file(p)
            if img:
                images.append(img)
            if summary:
                ctx.append(summary)
        content = text or "첨부한 파일을 분석해줘."
        if recalled:
            content += ("\n\n(관련 기억 — 참고만: "
                        + " / ".join(recalled) + ")")
        if ctx:
            content += "\n\n=== 첨부 파일 분석 ===\n" + "\n\n".join(ctx)

        instant_reply = None   # 창만 띄우면 끝나는 작업 → AI 생성 없이 즉시 응답

        # ① 프로그램/사이트 즉시 실행 — "스케치업 열어줘". 사이트+검색어 조합이면 검색창으로
        opened = []
        site_search = False
        if re.search(r"열어|열고|실행|띄워|켜\s*줘|켜줘|켜고|오픈|접속|들어가|open", text or "", re.I):
            low = (text or "").replace(" ", "")
            payload = web_query(text or "").replace(" ", "")
            for name in sorted(list(PROGRAMS) + list(SITES), key=len, reverse=True):
                if name in low and not any(name in m for m in opened):
                    if name in SITES and payload and payload != name and name not in payload:
                        site_search = True   # "유튜브 열어서 OO 검색" → 홈 말고 검색결과 창
                        continue
                    opened.append(name)
            for url in re.findall(r"https?://\S+", text or ""):
                opened.append(url)
            for t in opened:
                self._sysmsg(open_target(t))
        if opened:
            content += ("\n\n(시스템: 사용자 요청으로 다음을 이미 실행했음 — "
                        + ", ".join(opened) + ". 다시 열지 말고 실행했다고만 답하라)")
            leftover = text or ""
            for name in opened:
                leftover = leftover.replace(name, "")
            leftover = re.sub(r"열어줘|열어봐|열어|열고|실행해줘|실행|띄워줘|띄워|켜줘|켜고|켜|"
                              r"접속해줘|접속|들어가줘|들어가|오픈|해줘|해봐|좀|요|\s", "", leftover)
            if len(leftover) <= 3:
                instant_reply = "🚀 " + ", ".join(opened) + " 열었습니다."

        # ①-2 URL 본문 요약 — 메시지에 링크가 있고 '열어달라'는 게 아니면 본문을 긁어와 첨부
        if not opened and not instant_reply:
            urls = re.findall(r"https?://[^\s)>\]]+", text or "")
            if urls:
                self._status("🌐 링크 내용을 읽는 중…")
                fetched_pages = []
                for u in urls[:2]:
                    page = fetch_url_text(u)
                    if page:
                        fetched_pages.append(page)
                if fetched_pages:
                    content += ("\n\n=== 링크에서 가져온 본문 (이 내용을 근거로 답하라) ===\n"
                                + "\n\n".join(fetched_pages))
                    self._sysmsg(f"🌐 링크 {len(fetched_pages)}개의 본문을 읽어 왔습니다")

        # ② 시스템 주도 컴퓨터 검색 — 파일 관련 요청이면 먼저 찾아서 모델에게 전달
        #    (단, "쿠팡에서/인스타에서"처럼 사이트를 지목했으면 그 사이트에서 검색)
        results = []
        if not opened and not site_search and not mentions_platform(text):
            try:
                auto_block, auto_imgs, results = auto_context(text)
            except Exception:
                auto_block, auto_imgs, results = None, [], []
            if auto_block:
                content += "\n\n" + auto_block
                images.extend(auto_imgs)
                self._sysmsg(f"🔍 컴퓨터에서 관련 파일 {len(results)}건을 먼저 찾았습니다")

        # ③ 파일 직접 열기 — '열어줘/보고싶다/틀어줘'면 무조건 바로 열고, '찾아줘/어디'면 폴더+선택
        if results and re.search(r"열어|열고|띄워|실행|틀어|찾아|어디|보여|보고\s*싶|볼래|보자",
                                 text or ""):
            top_sc, _, top_p = results[0]
            try:
                direct = re.search(r"열어|열고|띄워|실행|틀어|보고\s*싶|볼래|보자", text or "")
                if direct:
                    os.startfile(top_p)          # 열어달라고 하면 항상 바로 연다
                    act = f"**{os.path.basename(top_p)}** 파일을 바로 열었습니다."
                else:
                    subprocess.Popen(f'explorer /select,"{top_p}"')
                    act = f"폴더 창을 열고 **{os.path.basename(top_p)}** 파일을 선택해 뒀습니다."
                if not READ_RE.search(text or ""):
                    others = "\n".join(f"- {os.path.basename(p)}" for sc, mt, p in results[1:3])
                    instant_reply = "📂 " + act + (f"\n\n다른 후보:\n{others}" if others else "")
                else:
                    content += (f"\n\n(시스템: 이미 다음을 수행했음 — {act} 다시 열지 말 것. "
                                "한 문장으로만 안내하고 질문에 답하라)")
            except Exception:
                pass

        # ③-1 모르는 사이트/이름 열기 — 목록에 없으면 네이버 검색으로 바로 진입시켜 줌
        if (not opened and not site_search and not results and not instant_reply
                and re.search(r"열어|들어가|접속|띄워", text or "")):
            name = web_query(text)
            if 0 < len(name.replace(" ", "")) <= 15:
                open_url("https://search.naver.com/search.naver?query="
                         + urllib.parse.quote(name))
                self._sysmsg(f"🌐 '{name}' 검색 창을 띄웠습니다")
                instant_reply = (f"🌐 '{name}' 사이트를 네이버 검색으로 띄웠습니다 — "
                                 "맨 위 결과를 클릭하면 바로 들어갑니다.")

        # 인터넷 검색 — 시스템이 브라우저 창을 직접 띄워 과정을 보여주고, 결과 요지는 모델에게 전달
        #    (복합 요청이면 통검색하지 않고 모델이 단계를 나눠 ```web```/```open```으로 직접 처리)
        try:
            multistep = is_multistep(text)
            if multistep and not instant_reply:
                content += ("\n\n(시스템: 이 요청은 여러 행동이 섞인 복합 요청이다. "
                            "문장을 통째로 검색하지 마라. 단계로 나눠 계획을 1~2줄로 밝힌 뒤, "
                            "정보가 필요한 단계는 ```web```으로 검색하고, 사용자에게 보여줄 "
                            "사이트는 ```open```으로 검색 URL을 열어 순서대로 처리하라)")
            need_web = (site_search
                        or mentions_platform(text)
                        or ((WEB_FORCE_RE.search(text or "")
                             or WEB_TOPIC_RE.search(text or "")) and not results))
            if need_web and not opened and not instant_reply and not multistep:
                q = web_query(text)
                if len(q) >= 2:
                    open_url(search_url(text, q))
                    self._sysmsg(f"🌐 브라우저 창을 띄워 '{q}' 검색을 보여드립니다")
                    if not WEB_TOPIC_RE.search(text or ""):
                        # 정보 질문이 아니라 그냥 보고 싶은 것 → AI 없이 즉시 마무리
                        instant_reply = f"🌐 '{q}' 검색 창을 띄웠습니다. 브라우저에서 보세요."
                    else:
                        fetched = web_search_fetch(q)
                        if fetched:
                            content += ("\n\n=== 시스템이 방금 인터넷에서 검색한 결과 ===\n" + fetched
                                        + "\n\n(브라우저 검색 창도 이미 띄워 줬다. 위 결과를 근거로 "
                                          "핵심을 한국어로 요약해 답하고, 출처 URL을 한두 개 적어라. "
                                          "결과가 부족하면 부족하다고 말하고 브라우저 창을 보라고 안내하라)")
                        else:
                            instant_reply = (f"🌐 '{q}' 검색 창을 띄웠습니다. "
                                             "자세한 내용은 브라우저에서 확인하세요.")
        except Exception:
            pass

        msg = {"role": "user", "content": content}
        if images:
            msg["images"] = images
        self._messages.append(msg)

        # 창 띄우기로 끝나는 작업 → AI 생성 없이 즉시 응답하고 종료 (다음 채팅 바로 가능)
        if instant_reply and not images:
            self._messages.append({"role": "assistant", "content": instant_reply})
            self._js("hideThinking()")
            self._js(f"aiMsg({json.dumps(instant_reply)})")
            self._save_history()
            self._status(f"● {self._model}", "ok")
            self._busy = False
            self._js("setBusy(false)")
            self._next_in_queue()
            return

        use_model = self._model_vision if images else self._model
        # 잡담은 출력을 짧게 잡아 더 빨리 끝낸다(모델 스왑 없음). 작업은 충분히 길게.
        self._num_predict = 150 if (not images and self._is_smalltalk(text)) else 800
        if self._num_predict == 150:
            self._status("빠른 응답 중…")
        else:
            self._status("분석 중…" + (" (이미지 모델)" if images else ""))
        try:
            for hop in range(5):
                self._js("showThinking()")
                reply = self._call_model(
                    use_model,
                    lambda p: self._js(f"aiToken({json.dumps(p)})"),
                    has_images=bool(images))
                self._messages.append({"role": "assistant", "content": reply})
                self._js("aiDone()")

                # 열기 도구 — 프리패스: 모델이 출력하면 바로 실행
                for ob in re.findall(r"```open\n(.*?)```", reply, re.S):
                    for line in ob.strip().splitlines():
                        t = line.strip()
                        if t and t.replace(" ", "") not in opened and t not in opened:
                            opened.append(t)
                            self._sysmsg(open_target(t))

                # 일정·메모 도구 — 프리패스: 바로 저장 + 알림 자동 예약
                for pb in re.findall(r"```plan\n(.*?)```", reply, re.S):
                    try:
                        self._sysmsg(handle_plan_block(pb))
                        self._js("calRefresh()")
                    except Exception as e:
                        self._sysmsg(f"일정 저장 오류: {e}")

                # 워크플로우 도구 — 프로젝트/업무 추가 후 보드 갱신
                for wb in re.findall(r"```work\n(.*?)```", reply, re.S):
                    try:
                        self._sysmsg(handle_work_block(wb))
                        self._js("if(window._wfRefresh)window._wfRefresh()")
                    except Exception as e:
                        self._sysmsg(f"워크플로우 저장 오류: {e}")

                # 화면 자동화 — 프리패스: 바로 실행 (마우스를 화면 모서리로 옮기면 즉시 중단)
                auto_blocks = re.findall(r"```auto\n(.*?)```", reply, re.S)
                if auto_blocks:
                    def auto_worker(blocks=auto_blocks):
                        self._sysmsg("🤖 자동화 시작 — 중단하려면 마우스를 화면 모서리로 옮기세요")
                        for b in blocks:
                            run_auto(b, self._sysmsg)
                        self._sysmsg("🤖 자동화 완료")
                    threading.Thread(target=auto_worker, daemon=True).start()

                # 웹 검색 도구 — 모델이 단계별로 필요한 정보를 직접 검색 (브라우저 창 없이)
                wb = re.search(r"```web\n(.*?)```", reply, re.S)
                if wb and hop < 4:
                    queries = [l.strip() for l in wb.group(1).strip().splitlines()
                               if l.strip()][:3]
                    parts = []
                    for q in queries:
                        self._status(f"🌐 인터넷 검색 중… ({q})")
                        fetched = web_search_fetch(q)
                        parts.append(f"[검색어: {q}]\n"
                                     + (fetched or "결과 없음 — 다른 검색어로 다시 ```web``` "
                                        "하거나, 브라우저로 보여주려면 ```open```으로 검색 URL을 열어라"))
                    self._messages.append({"role": "user", "content":
                        "=== 인터넷 검색 결과 ===\n" + "\n\n".join(parts)
                        + "\n\n이 결과를 근거로 계획의 다음 단계를 이어서 진행하라."})
                    self._sysmsg(f"🌐 인터넷 검색 {len(queries)}건 — 검토 중")
                    continue

                # 검색 도구 — 컴퓨터 전체 색인에서 즉시 검색
                m = re.search(r"```search\n(.*?)```", reply, re.S)
                if m and hop < 4:
                    kws, exts = parse_search_block(m.group(1))
                    self._status(f"🔍 컴퓨터 전체 검색 중… ({' '.join(kws)})")
                    results = search_files(kws, exts)
                    res_text = ("=== 검색 결과 (" + str(len(results)) + "건, 최신순) ===\n"
                                + ("\n".join(results) if results else "결과 없음")
                                + "\n\n이 결과로 요청을 처리하라. 내용 확인이 필요하면 ```read```로 "
                                  "읽어라. 결과가 없으면 다른 키워드로 한 번 더 ```search``` "
                                  "하거나 사용자에게 물어라.")
                    self._messages.append({"role": "user", "content": res_text})
                    self._sysmsg(f"🔍 컴퓨터 전체에서 {len(results)}개 파일 발견 — 검토 중")
                    continue

                # 읽기 도구 — 모델이 요청한 파일 내용을 읽어 전달 (코워크식 활용)
                rblocks = re.findall(r"```read\n(.*?)```", reply, re.S)
                if rblocks and hop < 4:
                    parts, rimgs, cnt = [], [], 0
                    for rb in rblocks:
                        for line in rb.strip().splitlines():
                            p = line.strip().strip('"').strip("'")
                            if not p or cnt >= 3:
                                continue
                            cnt += 1
                            if not os.path.exists(p):
                                parts.append(f"[{p}] 파일이 존재하지 않습니다 — 경로를 다시 확인하라")
                                continue
                            summary, img = analyze_file(p)
                            if img:
                                rimgs.append(img)
                                parts.append(f"[이미지 전달됨: {os.path.basename(p)}] "
                                             "이미지를 보고 분석하라")
                            if summary:
                                parts.append(summary)
                    rmsg = {"role": "user",
                            "content": "=== 파일 내용 ===\n" + "\n\n".join(parts)
                                       + "\n\n이 내용을 바탕으로 사용자 요청을 이어서 처리하라."}
                    if rimgs:
                        rmsg["images"] = rimgs
                        images = rimgs
                        use_model = self._model_vision
                    self._messages.append(rmsg)
                    self._sysmsg(f"📖 파일 {cnt}개 내용 확인 중")
                    continue

                # 코드 실행 — 프리패스. 파이썬은 즉시 실행하고, 실패하면 오류 메시지를
                #   모델에게 돌려줘 스스로 고쳐 다시 실행하게 한다 (자가수정 루프)
                if self._stop:
                    break
                code_blocks = re.findall(r"```(python|ruby)\n(.*?)```", reply, re.S)
                if code_blocks:
                    errs = []
                    for lang, code in code_blocks:
                        if lang == "ruby":
                            self._sysmsg("루비 코드는 스케치업 메뉴 > 창 > Ruby 콘솔에 붙여넣어 실행하세요")
                            continue
                        ok, out = self._exec_code(code)
                        if not ok:
                            errs.append(out)
                    if errs and hop < 4 and not self._stop:
                        self._messages.append({"role": "user", "content":
                            "=== 방금 실행한 코드가 오류로 실패했다 ===\n" + "\n---\n".join(errs)
                            + "\n\n위 오류 메시지를 보고 원인을 고쳐서, 바로 실행 가능한 "
                              "```python 코드 전체를 다시 작성하라. 같은 실수를 반복하지 말고 "
                              "필요한 import·파일 경로를 먼저 확인하라. 설명은 한 줄로 줄여라."})
                        self._sysmsg("🔧 코드 오류를 모델에게 전달 — 스스로 고쳐 재시도합니다")
                        continue
                break

            mem = re.search(r"```memory\n(.*?)```", reply, re.S)
            if mem:
                note = mem.group(1).strip()
                if note:
                    with open(PROFILE, "a", encoding="utf-8") as f:
                        f.write(f"- {note} ({time.strftime('%Y-%m-%d')})\n")
                    threading.Thread(target=memory_add, args=(note,),
                                     daemon=True).start()   # 임베딩 기억에도 저장(자동 회상)
                    self._messages[0] = {"role": "system", "content": build_system()}
                    self._sysmsg(f"🧠 기억했습니다: {note}")

            self._save_history()
            self._status(f"● {self._model}", "ok")
        except Exception as e:
            import traceback
            log("_reply 오류:\n" + traceback.format_exc())
            self._js("aiDone()")
            self._sysmsg(f"오류가 발생했습니다: {e}")
            self._status("오류", "err")
        finally:
            self._busy = False
            self._js("setBusy(false)")
            self._next_in_queue()


def ensure_autostart(enable=True):
    """윈도우 시작 시 자동 실행 등록/해제 (HKCU Run 키, 관리자 권한 불필요)."""
    try:
        import winreg
        pyw = sys.executable                       # pythonw.exe (창 없이 실행)
        script = os.path.abspath(__file__)
        cmd = f'"{pyw}" "{script}"'
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                             r"Software\Microsoft\Windows\CurrentVersion\Run", 0,
                             winreg.KEY_SET_VALUE)
        if enable:
            winreg.SetValueEx(key, "AI비서", 0, winreg.REG_SZ, cmd)
        else:
            try:
                winreg.DeleteValue(key, "AI비서")
            except FileNotFoundError:
                pass
        winreg.CloseKey(key)
    except Exception as e:
        log(f"autostart 설정 실패: {e}")


class PikaApi:
    """피카츄 외부 창 전용 작은 API (메인 창과 분리 — on_ready 중복 실행 방지)."""

    def __init__(self, app):
        self.app = app

    def focus(self):
        """피카츄를 클릭하면 메인 비서 창을 앞으로 띄우고 입력창에 포커스."""
        try:
            self.app._window.show()
            try:
                self.app._window.restore()
            except Exception:
                pass
            self.app._js("(function(){var i=document.getElementById('inp');if(i)i.focus();})()")
        except Exception:
            pass
        return "ok"

    def hide(self):
        """피카츄 위젯의 ✕ — 외부 창을 숨기고 설정을 끔(메인 토글도 동기화)."""
        save_settings({"pika": False})
        try:
            if getattr(self.app, "_pika", None):
                self.app._pika.hide()
        except Exception:
            pass
        try:
            self.app._js("window._pikaOn=false;"
                         "var c=document.getElementById('set_pika');if(c)c.checked=false;")
        except Exception:
            pass
        return "ok"


TRAY = {"ok": False}      # 트레이가 떠 있으면 창 닫기 = 종료 대신 숨기기


def run_tray(window):
    """시스템 트레이 상주 — 창을 닫아도 백그라운드로 살아있고, 트레이로 다시 연다."""
    try:
        import pystray
        from PIL import Image
        ico = os.path.join(HOME, "비서.ico")
        image = Image.open(ico) if os.path.exists(ico) else Image.new("RGB", (64, 64), "#D97757")

        def show(icon, item):
            try:
                window.show()
                window.restore()
            except Exception:
                pass

        def quit_all(icon, item):
            icon.stop()
            os._exit(0)

        menu = pystray.Menu(
            pystray.MenuItem("KJH비서 열기", show, default=True),
            pystray.MenuItem("종료", quit_all))
        TRAY["ok"] = True
        pystray.Icon("KJH비서", image, "KJH비서", menu).run()
    except Exception as e:
        TRAY["ok"] = False
        log(f"트레이 실행 실패: {e}")


if __name__ == "__main__":
    # 전역 예외를 파일로 남겨, 창 없는 .pyw에서 에러가 조용히 사라지지 않게 한다
    def _excepthook(exc_type, exc, tb):
        import traceback
        log("처리되지 않은 예외:\n" + "".join(traceback.format_exception(exc_type, exc, tb)))
    sys.excepthook = _excepthook
    try:
        threading.excepthook = lambda a: log(f"스레드 예외: {a.exc_type.__name__}: {a.exc_value}")
    except Exception:
        pass

    os.environ["PATH"] += os.pathsep + os.path.expandvars(r"%LOCALAPPDATA%\Programs\Ollama")
    for d in (OUT_DIR, HIST_DIR, INBOX):
        os.makedirs(d, exist_ok=True)
    ensure_autostart(True)                         # 부팅 시 자동 실행 등록

    api = Api()
    api._window = webview.create_window(
        "KJH비서", html=HTML, js_api=api,
        width=1180, height=820, min_size=(700, 560), background_color="#FAF9F5")

    # 피카츄 떠다니는 외부 위젯 — 메인 창 밖, 항상 위, 처음엔 우측 하단. 드래그로 이동.
    api._pika = None
    try:
        pw, ph = 170, 200
        px, py = 1240, 560          # 화면 크기를 못 읽을 때의 기본 우측 하단값
        try:
            scr = webview.screens[0]
            px = max(0, scr.width - pw - 24)
            py = max(0, scr.height - ph - 70)   # 작업표시줄 위쪽
        except Exception:
            pass
        api._pika = webview.create_window(
            "피카츄", html=PIKA_HTML, js_api=PikaApi(api),
            width=pw, height=ph, x=px, y=py,
            frameless=True, easy_drag=True, on_top=True,
            resizable=False, transparent=True, background_color="#FAF9F5")
        # 설정에서 꺼져 있으면 시작 시 숨김
        try:
            if not load_settings().get("pika", True):
                api._pika.hide()
        except Exception:
            pass
    except Exception as e:
        log(f"피카츄 창 생성 실패: {e}")

    def _on_closing():
        # 트레이가 살아 있으면 창만 숨기고 백그라운드로 상주 (트레이 아이콘으로 다시 연다)
        if TRAY["ok"]:
            try:
                api._window.hide()
            except Exception:
                return True
            return False        # 닫기 취소 → 숨김 유지
        return True             # 트레이 없으면 정상 종료

    try:
        api._window.events.closing += _on_closing
    except Exception as e:
        log(f"closing 핸들러 등록 실패: {e}")
    threading.Thread(target=run_tray, args=(api._window,), daemon=True).start()
    webview.start()
